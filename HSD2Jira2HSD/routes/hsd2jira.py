"""HSD → Jira modernisation — /hsd2jira"""
import json
import logging
import re
import threading
import uuid
from concurrent.futures import ThreadPoolExecutor, as_completed
from datetime import datetime
from pathlib import Path
import requests
import urllib3
from flask import Blueprint, render_template, request, jsonify, send_file
from requests_kerberos import HTTPKerberosAuth, OPTIONAL

from config import (
    HSD2JIRA_TOKEN_PATH, BASE_DIR as _BASE_DIR,
    JIRA_BASE as _JIRA_BASE, JIRA_PROJECT as _PROJECT,
    JIRA_CF_EPIC_NAME as _EPIC_NAME, JIRA_CF_EPIC_LINK as _EPIC_LINK,
    JIRA_CF_EXT_ISSUE as _EXT_ISSUE_ID, JIRA_CF_TREND_WW as _JIRA_TREND_WW,
    HSD_BASE as _HSD_BASE,
    DEFAULT_COMPONENT as _COMPONENT, DEFAULT_ASSIGNEE as _ASSIGNEE,
)

urllib3.disable_warnings()
logger = logging.getLogger(__name__)

# Templates live in HSD2Jira2HSD/templates/ — works whether this blueprint is
# loaded standalone or imported into data_centre_landing via the loader stub.
_TEMPLATE_DIR = Path(__file__).resolve().parent.parent / 'templates'
bp = Blueprint('hsd2jira', __name__, template_folder=str(_TEMPLATE_DIR))

_JIRA_API   = _JIRA_BASE + '/rest/api/2'
_HSD_REST        = _HSD_BASE + '/rest/article'
_ESSERVICE_URL   = _HSD_BASE + '/ws/ESService'
_ESSERVICE_HDRS  = {
    'APP':          'HSD-ES Article',
    'Accept':       'application/json',
    'Content-Type': 'application/json',
}


# ── Auth helpers ──────────────────────────────────────────────────────────────

def _read_token() -> str:
    """Read Jira bearer token from disk. Raises RuntimeError on failure."""
    try:
        return HSD2JIRA_TOKEN_PATH.read_text().strip()
    except Exception as exc:
        raise RuntimeError(f'Jira token not found — {exc}') from exc


def _token_response() -> 'tuple[str, object]':
    """Return (token, None) or ('', error_flask_response) for use in routes."""
    try:
        return _read_token(), None
    except RuntimeError as exc:
        logger.error('%s', exc)
        return '', (jsonify({'error': str(exc)}), 500)


# ── HSD helper functions ──────────────────────────────────────────────────────

def _new_hsd_session() -> requests.Session:
    """Create a Kerberos-authenticated session for HSD REST calls."""
    s = requests.Session()
    s.auth   = HTTPKerberosAuth(mutual_authentication=OPTIONAL)
    s.verify = False
    return s


def _esservice_request(session: requests.Session, command: str, command_args: dict):
    """Make a request to the HSD ESService API, returning result_table or None."""
    payload = {
        'requests': [{
            'api_client':   'HSD-ES Article',
            'tran_id':      str(uuid.uuid4()).upper(),
            'command':      command,
            'command_args': command_args,
            'var_args':     [],
            'copy_args':    [],
        }]
    }
    try:
        resp = session.post(_ESSERVICE_URL, headers=_ESSERVICE_HDRS,
                            json=payload, timeout=30)
        if resp.status_code == 200:
            result = resp.json()
            if 'responses' in result:
                r0 = result['responses'][0]
                if r0.get('status') == 'success':
                    return r0.get('result_table', [])
                logger.warning('ESService non-success status: %s — %s',
                               r0.get('status'), r0.get('message'))
        else:
            logger.error('ESService HTTP %s', resp.status_code)
    except Exception as exc:
        logger.error('ESService request failed: %s', exc)
    return None


def _fetch_hsd_title(session: requests.Session, hsd_id: str) -> str:
    """Fetch HSD article via REST and return its title."""
    url = f'{_HSD_REST}/{hsd_id}'
    try:
        resp = session.get(url, timeout=20)
        if resp.status_code == 200:
            data  = resp.json()
            items = data.get('data', []) if isinstance(data, dict) else []
            if items:
                title = str(items[0].get('title') or '').strip()
                return title if title else f'HSD {hsd_id}'
    except Exception as exc:
        logger.error('HSD REST fetch failed for %s: %s', hsd_id, exc)
    return f'HSD {hsd_id}'


def _get_ar_children(session: requests.Session, hsd_id: str) -> list:
    """
    Fetch the AR children of hsd_id.
    Strategy:
      1. Call get_related_records without filters — returns all linked records.
      2. Keep only links where subject='ar' and tenant='dg_soc'.
      3. For each, call get_record_by_id to get the full article with dg_soc.ar.* fields.
    """
    links = _esservice_request(session, 'get_related_records', {'id': str(hsd_id)})
    if not links:
        return []

    ar_ids = [
        str(r['id'])
        for r in links
        if isinstance(r, dict)
        and str(r.get('subject', '')).lower() == 'ar'
        and str(r.get('tenant',  '')).lower() in ('dg_soc', 'dg.soc')
        and r.get('id')
    ]
    logger.info('AR link IDs for HSD %s: %s', hsd_id, ar_ids)

    def _fetch_one(ar_id):
        rec = _esservice_request(session, 'get_record_by_id', {'id': ar_id})
        if rec and isinstance(rec, list) and rec[0]:
            return rec[0]
        if isinstance(rec, dict):
            return rec
        return None

    full_records = []
    with ThreadPoolExecutor(max_workers=8) as pool:
        for rec in pool.map(_fetch_one, ar_ids):
            if rec is not None:
                full_records.append(rec)
    return full_records


def _parse_ar_record(rec: dict) -> dict:
    """Extract team/task/jira_key from a full AR article record."""
    def _g(*keys):
        for k in keys:
            v = str(rec.get(k) or '').strip()
            if v and v.lower() not in ('none', 'null', ''):
                return v
        return ''
    def _gjira(*keys):
        """Extract Jira key pattern; field value may be quoted or XML-wrapped."""
        for k in keys:
            raw = str(rec.get(k) or '').strip()
            if not raw or raw.lower() in ('none', 'null'):
                continue
            m = re.search(r'[A-Z]+-\d+', raw)
            if m:
                return m.group(0)
            if raw:
                return raw
        return ''
    return {
        'id':       str(rec.get('id', '')),
        'team':     _g('dg_soc.ar.team',     'dg.soc.ar.team',     'team'),
        'task':     _g('dg_soc.ar.task',     'dg.soc.ar.task',     'task'),
        'jira_key': _gjira('dg_soc.ar.jira_key', 'dg.soc.ar.jira_key', 'jira_key'),
        'status':   _g('dg_soc.ar.status',   'dg.soc.ar.status',   'status'),
    }


def _writeback_jira_key(session: requests.Session, child_id: str,
                        jira_key: str) -> bool:
    """
    Write a Jira key to the jira_key field of the matched AR child record.
    Tries several REST PUT formats first (ESService update_record returns a
    false-positive success without actually writing the field).
    """
    put_bodies = [
        # Format 1 — fieldValues list (same structure used for create)
        {'tenant': 'dg_soc', 'subject': 'ar',
         'fieldValues': [{'dg_soc.ar.jira_key': jira_key}]},
        # Format 2 — flat data list
        {'data': [{'dg_soc.ar.jira_key': jira_key}]},
        # Format 3 — flat top-level field
        {'dg_soc.ar.jira_key': jira_key},
    ]
    for i, body in enumerate(put_bodies, 1):
        try:
            resp = session.put(
                f'{_HSD_REST}/{child_id}',
                json=body,
                headers={'Content-Type': 'application/json'},
                timeout=20,
            )
            if resp.status_code in (200, 201, 204):
                logger.info('HSD write-back OK via REST PUT (format %d): child=%s jira_key=%s',
                            i, child_id, jira_key)
                return True
            logger.warning('HSD REST PUT format %d returned %s: %s',
                           i, resp.status_code, resp.text[:200])
        except Exception as exc:
            logger.error('HSD REST PUT format %d failed for %s: %s', i, child_id, exc)

    # Final fallback — ESService update_record
    result = _esservice_request(session, 'update_record', {
        'id':      str(child_id),
        'tenant':  'dg_soc',
        'subject': 'ar',
        'fields':  {'dg_soc.ar.jira_key': jira_key},
    })
    if result is not None:
        logger.info('HSD write-back via ESService (verify manually): child=%s jira_key=%s',
                    child_id, jira_key)
        return True

    logger.error('HSD write-back failed for child=%s — all methods exhausted', child_id)
    return False


def _update_dg_soc_ar_field(session: requests.Session, child_id: str,
                            field_name: str, value: str) -> bool:
    """Write a single field on a dg_soc.ar child record.
    Tries REST PUT in three body formats then falls back to ESService update_record.
    """
    put_bodies = [
        {'tenant': 'dg_soc', 'subject': 'ar',
         'fieldValues': [{field_name: value}]},
        {'data': [{field_name: value}]},
        {field_name: value},
    ]
    for i, body in enumerate(put_bodies, 1):
        try:
            resp = session.put(
                f'{_HSD_REST}/{child_id}',
                json=body,
                headers={'Content-Type': 'application/json'},
                timeout=20,
            )
            if resp.status_code in (200, 201, 204):
                logger.info('dg_soc.ar field update OK (fmt %d): child=%s %s=%s',
                            i, child_id, field_name, value)
                return True
            logger.warning('dg_soc.ar PUT fmt %d → %s: %s',
                           i, resp.status_code, resp.text[:200])
        except Exception as exc:
            logger.error('dg_soc.ar PUT fmt %d failed for %s: %s', i, child_id, exc)
    result = _esservice_request(session, 'update_record', {
        'id':      str(child_id),
        'tenant':  'dg_soc',
        'subject': 'ar',
        'fields':  {field_name: value},
    })
    if result is not None:
        logger.info('dg_soc.ar field update via ESService: child=%s %s=%s',
                    child_id, field_name, value)
        return True
    logger.error('dg_soc.ar field update failed for child=%s %s — all methods exhausted',
                 child_id, field_name)
    return False


# ── Query-processing helpers ──────────────────────────────────────────────────

def _fetch_hsd_info(session: requests.Session, hsd_id: str) -> dict:
    """Fetch HSD article and return title, tenant, subject, release, family_affected, status, state."""
    url = f'{_HSD_REST}/{hsd_id}'
    def _rdy(rec, name):
        """Try plain then tenant-prefixed key names for date/readiness fields."""
        for k in (name, f'ip_hw_graphics.dcn.{name}', f'dg_soc.dcn.{name}',
                  f'ip_hw_graphics.feature.{name}', f'dg_soc.feature.{name}'):
            v = str(rec.get(k) or '').strip()
            if v and v.lower() not in ('none', 'null', ''):
                return v
        return ''
    try:
        resp = session.get(url, params={'include_text_fields': 'Y'}, timeout=20)
        if resp.status_code == 200:
            data  = resp.json()
            items = data.get('data', []) if isinstance(data, dict) else []
            if items:
                r = items[0]
                return {
                    'title':            str(r.get('title')   or '').strip() or f'HSD {hsd_id}',
                    'tenant':           str(r.get('tenant')  or '').strip(),
                    'subject':          str(r.get('subject') or '').strip(),
                    'release':          str(r.get('release') or '').strip(),
                    'family_affected':  str(r.get('family_affected')  or '').strip(),
                    'release_affected': str(r.get('release_affected') or '').strip(),
                    'status':           str(r.get('status')  or '').strip(),
                    'state':            str(r.get('ip_hw_graphics.feature.state') or
                                            r.get('dg_soc.feature.state') or
                                            r.get('state') or '').strip(),
                    'fulsim_rdy':       _rdy(r, 'fulsim_rdy_for_drv_ww'),
                    'emu_rdy':          _rdy(r, 'rtl_actual_rdy_for_drv_ww'),
                    'rtl_trend_rdy':    _rdy(r, 'rtl_trend_rdy_for_drv_ww'),
                    'turnin_trend_ww':  _rdy(r, 'turnin_trend_ww'),
                    'study_priority':   str(r.get('ip_hw_graphics.feature.study_priority') or r.get('study_priority') or '').strip(),
                }
    except Exception as exc:
        logger.error('HSD info fetch failed for %s: %s', hsd_id, exc)
    return {'title': f'HSD {hsd_id}', 'tenant': '', 'subject': '', 'release': '',
            'family_affected': '', 'release_affected': '', 'status': '', 'state': '',
            'fulsim_rdy': '', 'emu_rdy': '', 'rtl_trend_rdy': '',
            'turnin_trend_ww': '', 'study_priority': ''}


def _is_feature_hsd(info: dict) -> bool:
    """Return True if the HSD article is a dg_soc.feature (tenant=dg_soc, subject=feature)."""
    tenant  = info.get('tenant',  '').lower().replace('.', '_')
    subject = info.get('subject', '').lower()
    return tenant in ('dg_soc', 'dg.soc') and subject == 'feature'


def _is_server_platf_feature(info: dict) -> bool:
    """Return True if the HSD article is a server_platf.feature."""
    tenant  = info.get('tenant',  '').lower().replace('.', '_').replace('-', '_')
    subject = info.get('subject', '').lower()
    return 'server_platf' in tenant and subject == 'feature'


def _title_has_kmd(title: str) -> bool:
    """Return True if the title contains any KMD variant (kmd, xekmd, Xe KMD, …)."""
    return bool(re.search(r'kmd', title, re.IGNORECASE))


def _get_server_platf_ar_children(session: requests.Session, hsd_id: str) -> list:
    """
    Fetch server_platf.ar children of hsd_id.
    Like _get_ar_children but filters for server_platf tenant.
    """
    links = _esservice_request(session, 'get_related_records', {'id': str(hsd_id)})
    if not links:
        return []

    ar_ids = [
        str(r['id'])
        for r in links
        if isinstance(r, dict)
        and str(r.get('subject', '')).lower() == 'ar'
        and 'server_platf' in str(r.get('tenant', '')).lower().replace('-', '_').replace('.', '_')
        and r.get('id')
    ]
    logger.info('server_platf AR link IDs for HSD %s: %s', hsd_id, ar_ids)

    def _fetch_one(ar_id):
        rec = _esservice_request(session, 'get_record_by_id', {'id': ar_id})
        if rec and isinstance(rec, list) and rec[0]:
            return rec[0]
        if isinstance(rec, dict):
            return rec
        return None

    full_records = []
    with ThreadPoolExecutor(max_workers=8) as pool:
        for rec in pool.map(_fetch_one, ar_ids):
            if rec is not None:
                full_records.append(rec)
    return full_records


def _parse_server_platf_ar(rec: dict) -> dict:
    """Extract fields from a server_platf.ar record."""
    def _g(*keys):
        for k in keys:
            v = str(rec.get(k) or '').strip()
            if v and v.lower() not in ('none', 'null', ''):
                return v
        return ''
    return {
        'id':     str(rec.get('id', '')),
        'title':  _g('title', 'server_platf.ar.title'),
        'tag':    _g('tag'),           # bare 'tag' field — no tenant prefix
        'team':   _g('server_platf.ar.team', 'team'),
        'task':   _g('server_platf.ar.task', 'task'),
        'status': _g('status'),
        'trend':  _g('server_platf.ar.trend', 'trend'),
    }


def _update_server_platf_ar_field(session: requests.Session, child_id: str,
                                  field_name: str, value: str) -> bool:
    """Write a single field on a server_platf.ar child record via ESService."""
    put_bodies = [
        {'tenant': 'server_platf', 'subject': 'ar',
         'fieldValues': [{field_name: value}]},
        {'data': [{field_name: value}]},
        {field_name: value},
    ]
    for i, body in enumerate(put_bodies, 1):
        try:
            resp = session.put(
                f'{_HSD_REST}/{child_id}',
                json=body,
                headers={'Content-Type': 'application/json'},
                timeout=20,
            )
            if resp.status_code in (200, 201, 204):
                logger.info('server_platf.ar field update OK (fmt %d): child=%s %s=%s',
                            i, child_id, field_name, value)
                return True
            logger.warning('server_platf.ar PUT fmt %d → %s: %s',
                           i, resp.status_code, resp.text[:200])
        except Exception as exc:
            logger.error('server_platf.ar PUT fmt %d failed for %s: %s', i, child_id, exc)
    result = _esservice_request(session, 'update_record', {
        'id':      str(child_id),
        'tenant':  'server_platf',
        'subject': 'ar',
        'fields':  {field_name: value},
    })
    if result is not None:
        logger.info('server_platf.ar field update via ESService: child=%s %s=%s',
                    child_id, field_name, value)
        return True
    logger.error('server_platf.ar field update failed for child=%s %s — all methods exhausted',
                 child_id, field_name)
    return False


def _writeback_tag_field(session: requests.Session, child_id: str,
                         jira_key: str) -> bool:
    """
    Write a Jira key to the tag field of a server_platf.ar record.
    The HSD field is the common 'tag' field (no tenant prefix).
    """
    put_bodies = [
        # Format 1 — fieldValues list with bare 'tag' key (confirmed field name from debug)
        {'tenant': 'server_platf', 'subject': 'ar',
         'fieldValues': [{'tag': jira_key}]},
        # Format 2 — flat top-level field
        {'tag': jira_key},
        # Format 3 — data list
        {'data': [{'tag': jira_key}]},
    ]
    for i, body in enumerate(put_bodies, 1):
        try:
            resp = session.put(
                f'{_HSD_REST}/{child_id}',
                json=body,
                headers={'Content-Type': 'application/json'},
                timeout=20,
            )
            if resp.status_code in (200, 201, 204):
                logger.info('HSD tag writeback OK via REST PUT (format %d): child=%s tag=%s',
                            i, child_id, jira_key)
                return True
            logger.warning('HSD tag PUT format %d returned %s: %s',
                           i, resp.status_code, resp.text[:200])
        except Exception as exc:
            logger.error('HSD tag PUT format %d failed for %s: %s', i, child_id, exc)

    # Final fallback — ESService update_record
    result = _esservice_request(session, 'update_record', {
        'id':      str(child_id),
        'tenant':  'server_platf',
        'subject': 'ar',
        'fields':  {'tag': jira_key},
    })
    if result is not None:
        logger.info('HSD tag writeback via ESService (verify manually): child=%s tag=%s',
                    child_id, jira_key)
        return True

    logger.error('HSD tag writeback failed for child=%s — all methods exhausted', child_id)
    return False


def _is_sw_wa_hsd(info: dict) -> bool:
    """Return True if the HSD article is an ip_hw_graphics.bugeco (SW WA bug)."""
    tenant  = info.get('tenant', '').lower().replace('.', '_').replace('-', '_')
    subject = info.get('subject', '').lower()
    return 'ip_hw_graphics' in tenant and subject == 'bugeco'


def _get_sw_impact_children(session: requests.Session, hsd_id: str) -> list:
    """
    Fetch ip_hw_graphics.sw_impact children of an ip_hw_graphics.bugeco parent.
    Uses ESService get_related_records then fetches each full record.
    """
    links = _esservice_request(session, 'get_related_records', {'id': str(hsd_id)})
    if not links:
        return []

    sw_ids = [
        str(r['id'])
        for r in links
        if isinstance(r, dict)
        and str(r.get('subject', '')).lower() == 'sw_impact'
        and 'ip_hw_graphics' in str(r.get('tenant', '')).lower().replace('.', '_').replace('-', '_')
        and r.get('id')
    ]
    logger.info('sw_impact link IDs for HSD %s: %s', hsd_id, sw_ids)

    full_records = []
    for sw_id in sw_ids:
        rec = _esservice_request(session, 'get_record_by_id', {'id': sw_id})
        if rec and isinstance(rec, list) and rec[0]:
            full_records.append(rec[0])
        elif isinstance(rec, dict):
            full_records.append(rec)
    return full_records


def _parse_sw_impact_record(rec: dict) -> dict:
    """Extract fields from an ip_hw_graphics.sw_impact record."""
    def _g(*keys):
        for k in keys:
            v = str(rec.get(k) or '').strip()
            if v and v.lower() not in ('none', 'null', ''):
                return v
        return ''
    return {
        'id':           str(rec.get('id', '')),
        'title':        _g('title', 'ip_hw_graphics.sw_impact.title'),
        'sw_component': _g('ip_hw_graphics.sw_impact.sw_component', 'sw_component'),
        'func_impact':  _g('ip_hw_graphics.sw_impact.func_impact',  'func_impact'),
        'sw_task':      _g('ip_hw_graphics.sw_impact.sw_task',      'sw_task'),
        'sw_exposure':  _g('ip_hw_graphics.sw_impact.sw_exposure',  'sw_exposure'),
        'sw_record':    _g('ip_hw_graphics.sw_impact.sw_record',     'sw_record'),
        'explanation':  _g('ip_hw_graphics.sw_impact.explanation',   'explanation'),
        'os':           _g('ip_hw_graphics.sw_impact.os',            'os'),
        'done':         _g('ip_hw_graphics.sw_impact.done',          'done'),
        'trend':        _g('ip_hw_graphics.sw_impact.trend',         'trend'),
    }


def _fetch_jira_status_and_trend(token: str, key: str) -> tuple[str, str]:
    """Fetch status name and Actual Trend WW from a Jira issue.
    Returns (status_name, trend_ww) — empty strings on any failure.
    """
    hdrs = {'Authorization': f'Bearer {token}', 'Accept': 'application/json'}
    try:
        r = requests.get(
            f'{_JIRA_API}/issue/{key}',
            params={'fields': f'status,{_JIRA_TREND_WW}'},
            headers=hdrs, verify=False, timeout=20,
        )
        if not r.ok:
            logger.warning('Jira GET %s returned %s', key, r.status_code)
            return '', ''
        fields = r.json().get('fields', {})
        status = (fields.get('status') or {}).get('name', '') or ''
        trend  = str(fields.get(_JIRA_TREND_WW) or '').strip()
        # Strip sub-WW decimal suffix: "26WW47.2" → "26WW47"
        trend  = trend.split('.')[0] if trend else trend
        return status.strip(), trend
    except Exception as exc:
        logger.error('Failed to fetch Jira %s: %s', key, exc)
        return '', ''


def _ww_to_hsd_ww(ww_str: str) -> str:
    """Convert Intel WW notation ('26WW25') to HSD YYYYWWx format ('2026WW25').
    Returns the original string unchanged if parsing fails.
    """
    m = re.match(r'^(\d{2})WW(\d{1,2})$', ww_str.strip(), re.IGNORECASE)
    if not m:
        return ww_str
    year = 2000 + int(m.group(1))
    week = int(m.group(2))
    return f'{year}WW{week}'


def _ww_to_date(ww_str: str) -> str:
    """Convert Intel WW notation ('26WW47') to a calendar date (Friday of that week).
    Returns the original string unchanged if parsing fails.
    """
    m = re.match(r'^(\d{2})WW(\d{1,2})$', ww_str.strip(), re.IGNORECASE)
    if not m:
        return ww_str
    year = 2000 + int(m.group(1))
    week = int(m.group(2))
    try:
        d = datetime.strptime(f'{year}-W{week:02d}-5', '%G-W%V-%u')  # Friday
        return d.strftime('%Y-%m-%d')
    except ValueError:
        return ww_str


def _update_sw_impact_field(session: requests.Session, child_id: str,
                            field_name: str, value: str) -> bool:
    """Write a single field on an ip_hw_graphics.sw_impact child record.
    Tries REST PUT in three formats then falls back to ESService update_record.
    """
    put_bodies = [
        {'tenant': 'ip_hw_graphics', 'subject': 'sw_impact',
         'fieldValues': [{field_name: value}]},
        {'data': [{field_name: value}]},
        {field_name: value},
    ]
    for i, body in enumerate(put_bodies, 1):
        try:
            resp = session.put(
                f'{_HSD_REST}/{child_id}',
                json=body,
                headers={'Content-Type': 'application/json'},
                timeout=20,
            )
            if resp.status_code in (200, 201, 204):
                logger.info('sw_impact field update OK (format %d): child=%s %s=%s',
                            i, child_id, field_name, value)
                return True
            logger.warning('sw_impact field PUT format %d returned %s: %s',
                           i, resp.status_code, resp.text[:200])
        except Exception as exc:
            logger.error('sw_impact field PUT format %d failed for %s: %s', i, child_id, exc)
    result = _esservice_request(session, 'update_record', {
        'id':      str(child_id),
        'tenant':  'ip_hw_graphics',
        'subject': 'sw_impact',
        'fields':  {field_name: value},
    })
    if result is not None:
        logger.info('sw_impact field update via ESService (verify manually): child=%s %s=%s',
                    child_id, field_name, value)
        return True
    logger.error('sw_impact field update failed for child=%s field=%s — all methods exhausted',
                 child_id, field_name)
    return False


def _is_ip_hw_feature_hsd(info: dict) -> bool:
    """Return True if the HSD article is an ip_hw_graphics.feature."""
    tenant  = info.get('tenant', '').lower().replace('.', '_').replace('-', '_')
    subject = info.get('subject', '').lower()
    return 'ip_hw_graphics' in tenant and subject == 'feature'


def _create_sw_impact_child(session: requests.Session, parent_id: str,
                             sw_exposure: str, owner: str = _ASSIGNEE) -> str | None:
    """
    Create a new ip_hw_graphics.sw_impact child under parent_id.
    Sets sw_component=i915_kmd, sw_task=development, os=linux.
    Returns the new record ID as a string, or None on failure.
    """
    field_values = [
        {'parent_id':                              str(parent_id)},
        {'ip_hw_graphics.sw_impact.sw_component': 'i915_kmd'},
        {'ip_hw_graphics.sw_impact.sw_task':      'development'},
        {'ip_hw_graphics.sw_impact.os':           'linux'},
        {'ip_hw_graphics.sw_impact.sw_exposure':   sw_exposure},
        {'owner':                                  owner},
    ]
    try:
        resp = session.post(
            f'{_HSD_BASE}/rest/article',
            json={'tenant': 'ip_hw_graphics', 'subject': 'sw_impact',
                  'fieldValues': field_values},
            headers={'Content-Type': 'application/json'},
            timeout=20,
        )
        if resp.status_code in (200, 201):
            data   = resp.json()
            new_id = str(data.get('new_id', ''))
            if new_id and new_id != 'None':
                logger.info('ip_hw_graphics.sw_impact child created: %s under %s', new_id, parent_id)
                return new_id
            logger.warning('sw_impact POST 200 but no new_id: %s', str(data)[:300])
        else:
            logger.warning('sw_impact POST returned %s: %s', resp.status_code, resp.text[:500])
    except Exception as exc:
        logger.error('sw_impact child creation failed: %s', exc)
    return None


def _writeback_sw_record(session: requests.Session, child_id: str,
                         jira_key: str) -> bool:
    """
    Write a Jira key to ip_hw_graphics.sw_impact.sw_record of a WA child record.
    Mirrors the pattern used for dg_soc and server_platf write-backs.
    """
    put_bodies = [
        {'tenant': 'ip_hw_graphics', 'subject': 'sw_impact',
         'fieldValues': [{'ip_hw_graphics.sw_impact.sw_record': jira_key}]},
        {'data': [{'ip_hw_graphics.sw_impact.sw_record': jira_key}]},
        {'ip_hw_graphics.sw_impact.sw_record': jira_key},
    ]
    for i, body in enumerate(put_bodies, 1):
        try:
            resp = session.put(
                f'{_HSD_REST}/{child_id}',
                json=body,
                headers={'Content-Type': 'application/json'},
                timeout=20,
            )
            if resp.status_code in (200, 201, 204):
                logger.info('SW WA sw_record writeback OK (format %d): child=%s key=%s',
                            i, child_id, jira_key)
                return True
            logger.warning('SW WA REST PUT format %d returned %s: %s',
                           i, resp.status_code, resp.text[:200])
        except Exception as exc:
            logger.error('SW WA REST PUT format %d failed for %s: %s', i, child_id, exc)

    result = _esservice_request(session, 'update_record', {
        'id':      str(child_id),
        'tenant':  'ip_hw_graphics',
        'subject': 'sw_impact',
        'fields':  {'ip_hw_graphics.sw_impact.sw_record': jira_key},
    })
    if result is not None:
        logger.info('SW WA writeback via ESService (verify manually): child=%s key=%s',
                    child_id, jira_key)
        return True

    logger.error('SW WA writeback failed for child=%s — all methods exhausted', child_id)
    return False


def _platform_from_release(release: str) -> str:
    """Return Jira platform value derived from HSD release field, or '' if no match."""
    if 'alpine' in release.lower():
        return 'tiger_shores'
    return ''


def _platform_from_family_affected(family_affected: str) -> tuple[str, str]:
    """Return (jira_platform, auto_project_tag) derived from HSD family_affected field."""
    fa = family_affected.lower()
    if 'jaguar shores' in fa:
        return 'Jaguar Shores', 'JGS'
    return '', ''


def _detect_platform_and_tag(info: dict) -> tuple[str, str]:
    """
    Auto-detect Jira platform and project tag from HSD info fields
    (family_affected, release_affected, tenant, subject).
    Returns (platform, tag) — both empty strings if no match.

    Rules per tenant:
      dg_soc.feature       : Alpine → tiger_shores/TGS | Crescent Island → crescent_island/CRI
      ip_hw_graphics.feature: xe5 → tiger_shores/TGS
                              xe3p + Xe3p_v1c_XPC → crescent_island/CRI
                              xe4  + Xe4_XPC       → Jaguar Shores/JGS
      server_platf.feature : Crescent Island → crescent_island/CRI
      server.feature       : Jaguar Shores   → Jaguar Shores/JGS
    """
    tenant  = info.get('tenant',  '').lower().replace('.', '_').replace('-', '_')
    subject = info.get('subject', '').lower()
    family  = info.get('family_affected',  '').lower()
    release = info.get('release_affected', '').lower()

    if 'dg_soc' in tenant and subject == 'feature':
        if 'alpine' in family:
            return 'tiger_shores', 'TGS'
        if 'crescent island' in family:
            return 'crescent_island', 'CRI'

    elif 'ip_hw_graphics' in tenant and subject == 'feature':
        if 'xe5' in family:
            return 'tiger_shores', 'TGS'
        if 'xe3p' in family and 'xe3p_v1c_xpc' in release:
            return 'crescent_island', 'CRI'
        if 'xe4' in family and 'xe4_xpc' in release:
            return 'Jaguar Shores', 'JGS'

    elif 'server_platf' in tenant and subject == 'feature':
        if 'crescent island' in family:
            return 'crescent_island', 'CRI'

    elif tenant == 'server' and subject == 'feature':
        if 'jaguar shores' in family:
            return 'Jaguar Shores', 'JGS'

    return '', ''


# ── server.feature helpers ─────────────────────────────────────────────────────

_SERVER_AR_VALID_COMPONENTS = ('sw.xe_kmd', 'sw.ual_kmd')


def _is_server_feature(info: dict) -> bool:
    """Return True if the HSD article is a server.feature (tenant=server, subject=feature)."""
    tenant  = info.get('tenant',  '').lower().strip()
    subject = info.get('subject', '').lower().strip()
    return tenant == 'server' and subject == 'feature'


def _get_server_ar_children(session: requests.Session, hsd_id: str) -> list:
    """Fetch server.ar children of hsd_id."""
    links = _esservice_request(session, 'get_related_records', {'id': str(hsd_id)})
    if not links:
        return []
    ar_ids = [
        str(r['id'])
        for r in links
        if isinstance(r, dict)
        and str(r.get('subject', '')).lower() == 'ar'
        and str(r.get('tenant',  '')).lower().strip() == 'server'
        and r.get('id')
    ]
    logger.info('server.ar link IDs for HSD %s: %s', hsd_id, ar_ids)

    def _fetch_one(ar_id):
        rec = _esservice_request(session, 'get_record_by_id', {'id': ar_id})
        if rec and isinstance(rec, list) and rec[0]:
            return rec[0]
        if isinstance(rec, dict):
            return rec
        return None

    full_records = []
    with ThreadPoolExecutor(max_workers=8) as pool:
        for rec in pool.map(_fetch_one, ar_ids):
            if rec is not None:
                full_records.append(rec)
    return full_records


def _parse_server_ar_record(rec: dict) -> dict:
    """Extract fields from a server.ar record."""
    def _g(*keys):
        for k in keys:
            v = str(rec.get(k) or '').strip()
            if v and v.lower() not in ('none', 'null', ''):
                return v
        return ''
    def _graw(*keys):
        """Like _g but preserves 'none' — meaningful for exposure."""
        for k in keys:
            v = str(rec.get(k) or '').strip()
            if v and v.lower() != 'null':
                return v
        return ''
    def _gjira(*keys):
        """Extract Jira key from field value; field may be XML blob."""
        for k in keys:
            raw = str(rec.get(k) or '').strip()
            if not raw or raw.lower() in ('none', 'null'):
                continue
            # Try to extract a Jira key pattern (e.g. VLK-76349)
            m = re.search(r'[A-Z]+-\d+', raw)
            if m:
                return m.group(0)
            if raw:
                return raw  # fallback: return raw if no pattern found
        return ''
    return {
        'id':                  str(rec.get('id', '')),
        'title':               _g('title', 'server.ar.title'),
        'component':           _g('server.ar.component', 'component'),
        'exposure':            _graw('server.ar.exposure', 'exposure'),
        'jira_key':            _gjira('server.ar.JIRA_KEY', 'server.ar.jira_key', 'jira_key'),
        'new_projected_dates': _g('server.ar.new_projected_dates', 'new_projected_dates'),
        'reason':              _g('server.ar.reason', 'reason'),
    }


def _server_ar_jira_config(component: str) -> dict:
    """Map server.ar component value to Jira component name and assignee."""
    cl = component.lower()
    if 'ual_kmd' in cl:
        return {'jira_component': 'Kernel - UAL', 'assignee': 'myossefi'}
    # default: sw.xe_kmd → XeKMD
    return {'jira_component': 'XeKMD', 'assignee': _ASSIGNEE}


def _create_server_ar_dev_child(session: requests.Session, parent_id: str,
                                 component: str, exposure: str) -> str | None:
    """Create a server.ar 'SW Development Review' child under parent_id."""
    field_values = [
        {'title':               'SW Development Review'},
        {'parent_id':           str(parent_id)},
        {'server.ar.component': component},
        {'server.ar.exposure':  exposure},
    ]
    try:
        resp = session.post(
            f'{_HSD_BASE}/rest/article',
            json={'tenant': 'server', 'subject': 'ar', 'fieldValues': field_values},
            headers={'Content-Type': 'application/json'},
            timeout=20,
        )
        if resp.status_code in (200, 201):
            data   = resp.json()
            new_id = str(data.get('new_id', ''))
            if new_id and new_id != 'None':
                logger.info('server.ar dev child created: %s under %s', new_id, parent_id)
                return new_id
            logger.warning('server.ar POST 200 but no new_id: %s', str(data)[:300])
        else:
            logger.warning('server.ar POST returned %s: %s', resp.status_code, resp.text[:500])
    except Exception as exc:
        logger.error('server.ar child creation failed: %s', exc)
    return None


def _writeback_server_ar_jira_key(session: requests.Session, child_id: str,
                                   jira_key: str) -> bool:
    """Write jira_key to server.ar child record."""
    put_bodies = [
        {'tenant': 'server', 'subject': 'ar',
         'fieldValues': [{'server.ar.jira_key': jira_key}]},
        {'data': [{'server.ar.jira_key': jira_key}]},
        {'server.ar.jira_key': jira_key},
    ]
    for i, body in enumerate(put_bodies, 1):
        try:
            resp = session.put(
                f'{_HSD_REST}/{child_id}',
                json=body,
                headers={'Content-Type': 'application/json'},
                timeout=20,
            )
            if resp.status_code in (200, 201, 204):
                logger.info('server.ar jira_key writeback OK (format %d): child=%s key=%s',
                            i, child_id, jira_key)
                return True
            logger.warning('server.ar PUT format %d returned %s: %s',
                           i, resp.status_code, resp.text[:200])
        except Exception as exc:
            logger.error('server.ar PUT format %d failed for %s: %s', i, child_id, exc)
    result = _esservice_request(session, 'update_record', {
        'id':      str(child_id),
        'tenant':  'server',
        'subject': 'ar',
        'fields':  {'server.ar.jira_key': jira_key},
    })
    if result is not None:
        logger.info('server.ar writeback via ESService (verify manually): child=%s key=%s',
                    child_id, jira_key)
        return True
    logger.error('server.ar jira_key writeback failed for child=%s — all methods exhausted', child_id)
    return False


def _update_server_ar_field(session: requests.Session, child_id: str,
                            field_name: str, value: str) -> bool:
    """Write a single field on a server.ar child record.
    Tries REST PUT in three body formats then falls back to ESService update_record.
    """
    put_bodies = [
        {'tenant': 'server', 'subject': 'ar',
         'fieldValues': [{field_name: value}]},
        {'data': [{field_name: value}]},
        {field_name: value},
    ]
    for i, body in enumerate(put_bodies, 1):
        try:
            resp = session.put(
                f'{_HSD_REST}/{child_id}',
                json=body,
                headers={'Content-Type': 'application/json'},
                timeout=20,
            )
            if resp.status_code in (200, 201, 204):
                logger.info('server.ar field update OK (fmt %d): child=%s %s=%s',
                            i, child_id, field_name, value)
                return True
            logger.warning('server.ar PUT fmt %d → %s: %s',
                           i, resp.status_code, resp.text[:200])
        except Exception as exc:
            logger.error('server.ar PUT fmt %d failed for %s: %s', i, child_id, exc)
    result = _esservice_request(session, 'update_record', {
        'id':      str(child_id),
        'tenant':  'server',
        'subject': 'ar',
        'fields':  {field_name: value},
    })
    if result is not None:
        logger.info('server.ar field update via ESService: child=%s %s=%s',
                    child_id, field_name, value)
        return True
    logger.error('server.ar field update failed for child=%s %s — all methods exhausted',
                 child_id, field_name)
    return False


def _parse_ar_record_extended(rec: dict) -> dict:
    """Like _parse_ar_record but also extracts exposure, type_, and title."""
    def _g(*keys):
        for k in keys:
            v = str(rec.get(k) or '').strip()
            if v and v.lower() not in ('none', 'null', ''):
                return v
        return ''
    def _graw(*keys):
        """Like _g but preserves 'none' — it is a meaningful exposure value."""
        for k in keys:
            v = str(rec.get(k) or '').strip()
            if v and v.lower() != 'null':
                return v
        return ''
    def _gjira(*keys):
        """Extract Jira key pattern; field value may be quoted or XML-wrapped."""
        for k in keys:
            raw = str(rec.get(k) or '').strip()
            if not raw or raw.lower() in ('none', 'null'):
                continue
            m = re.search(r'[A-Z]+-\d+', raw)
            if m:
                return m.group(0)
            if raw:
                return raw
        return ''
    return {
        'id':       str(rec.get('id', '')),
        'title':    _g('title', 'dg_soc.ar.title',    'dg.soc.ar.title'),
        'team':     _g('dg_soc.ar.team',     'dg.soc.ar.team',     'team'),
        'task':     _g('dg_soc.ar.task',     'dg.soc.ar.task',     'task'),
        'jira_key': _gjira('dg_soc.ar.jira_key', 'dg.soc.ar.jira_key', 'jira_key'),
        'status':   _g('dg_soc.ar.status',   'dg.soc.ar.status',   'status'),
        'exposure': _graw('dg_soc.ar.exposure', 'dg.soc.ar.exposure', 'exposure'),
        'type_':    _g('dg_soc.ar.type',     'dg.soc.ar.type',     'type'),
        'eta_ww':   _g('dg_soc.ar.eta_ww',   'dg.soc.ar.eta_ww',   'ar.eta_wk'),
        'phase':    _g('dg_soc.ar.phase',    'dg.soc.ar.phase',    'phase'),
    }


def _execute_hsd_query(session: requests.Session, query_id: str):
    """
    Fetch HSD article IDs from a saved HSD query.
    Returns list of ID strings, or None on total failure.
    """
    # Try ESService execute_saved_query
    result = _esservice_request(session, 'execute_saved_query', {
        'query_id': str(query_id),
        'fields': ['id'],
    })
    if result is not None:
        ids = [str(r['id']) for r in result if isinstance(r, dict) and r.get('id')]
        if ids:
            return ids

    # Fallback 1 — REST /rest/query/{query_id} (paginated, server caps at 100/page)
    try:
        all_ids: list[str] = []
        start = 0
        while True:
            resp = session.get(
                f'{_HSD_BASE}/rest/query/{query_id}',
                params={'fields': 'id', 'start_at': start},
                timeout=30,
            )
            if resp.status_code != 200:
                break
            data  = resp.json()
            items = data.get('data', []) if isinstance(data, dict) else (data if isinstance(data, list) else [])
            page_ids = [str(i['id']) for i in items if isinstance(i, dict) and i.get('id')]
            all_ids.extend(page_ids)
            total = data.get('total', 0) if isinstance(data, dict) else 0
            start += len(items)
            if not items or start >= total:
                break
        if all_ids:
            return all_ids
    except Exception as exc:
        logger.warning('HSD query REST /rest/query failed: %s', exc)

    # Fallback 2 — REST /rest/article?query_id=... (paginated)
    try:
        all_ids = []
        start = 0
        while True:
            resp = session.get(
                f'{_HSD_BASE}/rest/article',
                params={'query_id': query_id, 'fields': 'id', 'start_at': start},
                timeout=30,
            )
            if resp.status_code != 200:
                break
            data  = resp.json()
            items = data.get('data', []) if isinstance(data, dict) else (data if isinstance(data, list) else [])
            page_ids = [str(i['id']) for i in items if isinstance(i, dict) and i.get('id')]
            all_ids.extend(page_ids)
            total = data.get('total', 0) if isinstance(data, dict) else 0
            start += len(items)
            if not items or start >= total:
                break
        if all_ids:
            return all_ids
    except Exception as exc:
        logger.warning('HSD query REST /rest/article?query_id failed: %s', exc)

    return None


_AR_TYPE_VALID = {'(unassigned)', 'other', 'rtl', 'study', 'sw_assessment'}


def _create_hsd_ar_child(session: requests.Session, parent_id: str,
                          title: str, type_: str, team: str,
                          exposure: str, task: str, jira_key: str):
    """
    Create a new dg_soc.ar child HSD under parent_id.
    Returns the new record ID as a string, or None on failure.

    Correct REST format (discovered via swagger + trial):
      POST /rest/article
      Body: {"tenant":"dg_soc","subject":"ar",
             "fieldValues":[{"title":...},{"parent_id":...},{"ar.type":...},...]}
      Response: {"new_id": <int>}
    """
    ar_type = 'sw_assessment'

    field_values = [
        {'title':              title},
        {'parent_id':          str(parent_id)},
        {'ar.type':            ar_type},
        {'dg_soc.ar.team':     team},
        {'dg_soc.ar.task':     task},
        {'dg_soc.ar.exposure': exposure},
        {'owner':              _ASSIGNEE},
    ]

    try:
        resp = session.post(
            f'{_HSD_BASE}/rest/article',
            json={'tenant': 'dg_soc', 'subject': 'ar', 'fieldValues': field_values},
            headers={'Content-Type': 'application/json'},
            timeout=20,
        )
        if resp.status_code in (200, 201):
            data = resp.json()
            new_id = str(data.get('new_id', ''))
            if new_id and new_id != 'None':
                logger.info('HSD AR child created: %s under %s', new_id, parent_id)
                return new_id
            logger.warning('HSD REST POST 200 but no new_id in response: %s', str(data)[:300])
        else:
            logger.warning('HSD REST POST returned %s: %s', resp.status_code, resp.text[:500])
    except Exception as exc:
        logger.error('HSD create_record REST failed: %s', exc)

    return None


def _make_jira_creator(token: str):
    """Return (post_issue, create_link) helpers bound to the given Jira token."""
    hdrs = {
        'Authorization': f'Bearer {token}',
        'Content-Type':  'application/json',
        'Accept':        'application/json',
    }

    def post_issue(payload):
        r = requests.post(f'{_JIRA_API}/issue', headers=hdrs,
                          json=payload, verify=False, timeout=20)
        if not r.ok:
            try:
                detail = r.json()
                msgs = detail.get('errorMessages', [])
                errs = detail.get('errors', {})
                raise requests.HTTPError(
                    f'Jira {r.status_code}: {msgs} {errs}', response=r)
            except (ValueError, KeyError):
                r.raise_for_status()
        return r.json()['key']

    def create_link(parent_key, child_key):
        try:
            r = requests.post(
                f'{_JIRA_API}/issueLink', headers=hdrs,
                json={'type': {'name': 'Child to Parent relations'},
                      'outwardIssue': {'key': parent_key},
                      'inwardIssue':  {'key': child_key}},
                verify=False, timeout=20,
            )
            r.raise_for_status()
        except Exception as exc:
            logger.error('Link failed %s -> %s: %s', parent_key, child_key, exc)

    return post_issue, create_link


def _norm_components(component) -> list:
    """Normalise component arg (str or list) to a Jira-ready list of dicts."""
    if isinstance(component, list):
        names = [c.strip() for c in component if str(c).strip()]
    else:
        names = [c.strip() for c in str(component or '').split(',') if c.strip()]
    return [{'name': n} for n in names] or [{'name': _COMPONENT}]


def _build_jira_issues(post_issue, create_link,
                       title: str, platform: str, component,
                       assignee: str, hsd_id: str = '',
                       child_hsd_id: str = '',
                       project_tag: str = '') -> dict:
    """
    Create Epic + 8 sub-issues for a feature.
    When hsd_id is provided, External Issue ID links are set on stories.
    Returns {'created': [...], 'errors': [...], 'epic_key': str, 'cc_story_key': str}.
    child_hsd_id: if provided, included in the External Issue ID field.
    component: str or list — one or more Jira component names.
    project_tag: if set, adds '{tag}_DCN_Query' label to the CC parent story.
    """
    created: list = []
    errors:  list = []
    _comps = _norm_components(component)
    _pf  = {'customfield_10700': [{'value': platform}]} if platform else {}
    _ext = None
    if hsd_id:
        _parent_url = f'https://hsdes.intel.com/appstore/article-one/#/{hsd_id}'
        _child_url  = f'https://hsdes.intel.com/appstore/article-one/#/{child_hsd_id}' if child_hsd_id else ''
        _ext = (f'[parent] {_parent_url}\n[child] {_child_url}'
                if _child_url else f'[parent] {_parent_url}')
    _ext_field = {_EXT_ISSUE_ID: _ext} if _ext else {}

    # 1. Epic
    epic_summary = f'[Epic] {title}'
    try:
        epic_key = post_issue({'fields': {
            'project':     {'key': _PROJECT},
            'summary':     epic_summary,
            'description': f'Auto-created from HSD {hsd_id}' if hsd_id else 'Pls add description',
            'issuetype':   {'name': 'Epic'},
            _EPIC_NAME:    epic_summary,
            'components':  _comps,
            'assignee':    {'name': assignee},
        }})
        created.append({'type': 'Epic', 'badge': 'badge-epic', 'key': epic_key,
                        'summary': epic_summary, 'note': f'HSD {hsd_id}' if hsd_id else 'Standalone feature',
                        'url': f'{_JIRA_BASE}/browse/{epic_key}', 'level': 0})
    except Exception as exc:
        return {'created': created,
                'errors': [{'summary': epic_summary, 'error': str(exc)}],
                'epic_key': None, 'cc_story_key': None}

    # 2. Code Complete Parent Story
    cc_summary    = f'[Parent][Story] [Code Complete] {title}'
    cc_parent_key = None
    _cc_labels = [f'{project_tag}_DCN_Query'] if project_tag else []
    try:
        cc_parent_key = post_issue({'fields': {
            'project':     {'key': _PROJECT},
            'summary':     cc_summary,
            'description': 'Pls add description',
            'issuetype':   {'name': 'Story'},
            _EPIC_LINK:    epic_key,
            'components':  _comps,
            'assignee':    {'name': assignee},
            **({'labels': _cc_labels} if _cc_labels else {}),
            **_ext_field,
            **_pf,
        }})
        created.append({'type': 'Story', 'badge': 'badge-story-cc', 'key': cc_parent_key,
                        'summary': cc_summary, 'note': 'Code Complete',
                        'url': f'{_JIRA_BASE}/browse/{cc_parent_key}', 'level': 1})
    except Exception as exc:
        errors.append({'summary': cc_summary, 'error': str(exc)})

    # 2a. CC Child Stories
    if cc_parent_key:
        for summary, badge, note, _child_comps in [
            (f'[Child][Story][Code Complete][Part-1] {title}', 'badge-child-ph1', 'Part-1',    _comps),
            (f'[Child][Story][Code Complete][Part-2] {title}', 'badge-child-ph2', 'Part-2',    _comps),
            (f'[Child][Story][Test][IGT] {title}',             'badge-child-igt', 'IGT Testing', [{'name': 'Test-Core'}]),
        ]:
            try:
                key = post_issue({'fields': {
                    'project': {'key': _PROJECT}, 'summary': summary,
                    'description': 'Pls add description', 'issuetype': {'name': 'Story'},
                    _EPIC_LINK: epic_key, 'components': _child_comps,
                    'assignee': {'name': assignee}, **_pf,
                }})
                create_link(cc_parent_key, key)
                created.append({'type': 'Child Story', 'badge': badge, 'key': key,
                                'summary': summary, 'note': note,
                                'url': f'{_JIRA_BASE}/browse/{key}', 'level': 2})
            except Exception as exc:
                errors.append({'summary': summary, 'error': str(exc)})

    # 3. Sim & Emu
    for summary, badge, note in [
        (f'[Parent][Story][Val][Sim] {title}', 'badge-story-sim', 'Simulation Validation'),
        (f'[Parent][Story][Val][Emu] {title}', 'badge-story-emu', 'Emulation Validation'),
    ]:
        try:
            key = post_issue({'fields': {
                'project': {'key': _PROJECT}, 'summary': summary,
                'description': 'Pls add description', 'issuetype': {'name': 'Story'},
                _EPIC_LINK: epic_key, 'components': _comps,
                'assignee': {'name': assignee}, **_ext_field, **_pf,
            }})
            created.append({'type': 'Story', 'badge': badge, 'key': key,
                            'summary': summary, 'note': note,
                            'url': f'{_JIRA_BASE}/browse/{key}', 'level': 1})
        except Exception as exc:
            errors.append({'summary': summary, 'error': str(exc)})

    # 4. Upstream Task
    upstream_summary = f'[Parent][Task][Upstream] {title}'
    upstream_key     = None
    for _fields in [
        {'project': {'key': _PROJECT}, 'summary': upstream_summary,
         'description': 'Pls add description', 'issuetype': {'name': 'Task'},
         _EPIC_LINK: epic_key, 'components': _comps,
         'assignee': {'name': assignee}, **_pf},
        {'project': {'key': _PROJECT}, 'summary': upstream_summary,
         'description': 'Pls add description', 'issuetype': {'name': 'Task'},
         _EPIC_LINK: epic_key, 'components': _comps,
         'assignee': {'name': assignee}},
        {'project': {'key': _PROJECT}, 'summary': upstream_summary,
         'description': 'Pls add description', 'issuetype': {'name': 'Task'},
         'components': _comps, 'assignee': {'name': assignee}},
    ]:
        try:
            upstream_key = post_issue({'fields': _fields})
            break
        except requests.exceptions.HTTPError as e:
            if e.response.status_code == 400:
                continue
            errors.append({'summary': upstream_summary, 'error': str(e)})
            break
        except Exception as exc:
            errors.append({'summary': upstream_summary, 'error': str(exc)})
            break

    if upstream_key:
        created.append({'type': 'Task', 'badge': 'badge-task', 'key': upstream_key,
                        'summary': upstream_summary, 'note': 'Upstream Task',
                        'url': f'{_JIRA_BASE}/browse/{upstream_key}', 'level': 1})
        ups_child = f'[Child][Story][Upstream] {title}'
        try:
            key = post_issue({'fields': {
                'project': {'key': _PROJECT}, 'summary': ups_child,
                'description': 'Pls add description', 'issuetype': {'name': 'Story'},
                _EPIC_LINK: epic_key, 'components': _norm_components(component),
                'assignee': {'name': assignee}, **_pf,
            }})
            create_link(upstream_key, key)
            created.append({'type': 'Child Story', 'badge': 'badge-child-ups', 'key': key,
                            'summary': ups_child, 'note': 'Upstream Implementation',
                            'url': f'{_JIRA_BASE}/browse/{key}', 'level': 2})
        except Exception as exc:
            errors.append({'summary': ups_child, 'error': str(exc)})

    return {'created': created, 'errors': errors,
            'epic_key': epic_key, 'cc_story_key': cc_parent_key}


# ── Shared per-HSD processing ─────────────────────────────────────────────────

def _process_one_hsd(session: requests.Session,
                      post_issue, create_link_fn,
                      hsd_id: str, info: dict,
                      platform: str = '',
                      project_tag: str = '',
                      component: str = '',
                      assignee: str = '',
                      token: str = '') -> dict:
    """
    Core per-HSD processing: detect tenant, apply rules, create Jira issues,
    write back to HSD children.
    Platform/tag are auto-detected from HSD fields; caller values are fallback.
    Returns a result dict compatible with hsd2jira_process_hsd_query.
    """
    auto_platform, auto_tag = _detect_platform_and_tag(info)
    eff_platform = platform    or auto_platform
    eff_tag      = project_tag or auto_tag
    eff_comp     = component or _COMPONENT
    eff_asgn     = assignee  or _ASSIGNEE

    title      = info.get('title', f'HSD {hsd_id}')
    jira_title = f'[{eff_tag}] {title}' if eff_tag else title

    result = {
        'hsd_id':           hsd_id,
        'title':            title,
        'tenant':           f"{info.get('tenant','')}.{info.get('subject','')}".strip('.'),
        'url':              f'https://hsdes.intel.com/appstore/article-one/#/{hsd_id}',
        'skipped':          False,
        'skip_reason':      None,
        'created_jira':     [],
        'cc_story_key':     None,
        'new_hsd_child_id': None,
        'hsd_writeback_ok': False,
        'writeback_field':  '',
        'errors':           [],
        'platform_used':    eff_platform,
        'family_affected':  info.get('family_affected', ''),
    }

    # ── server.feature ────────────────────────────────────────────────────────
    if _is_server_feature(info):
        ar_raw    = _get_server_ar_children(session, hsd_id)
        ar_parsed = [_parse_server_ar_record(r) for r in ar_raw]
        arch_reviews = [
            c for c in ar_parsed
            if c['title'].strip().lower() == 'sw arch review'
            and c['component'].lower() in _SERVER_AR_VALID_COMPONENTS
            and c['exposure'].lower() not in ('none', 'to_be_assigned', '')
        ]
        if not arch_reviews:
            result['skipped']     = True
            result['skip_reason'] = (
                f'No server.ar SW Arch Review children with component in '
                f'{_SERVER_AR_VALID_COMPONENTS} and valid exposure '
                f'({len(ar_parsed)} AR records checked)')
            return result

        for arch_child in arch_reviews:
            comp     = arch_child['component']
            exposure = arch_child['exposure']
            jira_cfg = _server_ar_jira_config(comp)
            dev_review = next(
                (c for c in ar_parsed
                 if c['title'].strip().lower() == 'sw development review'
                 and c['component'].lower() == comp.lower()),
                None
            )
            if dev_review:
                dev_exp = dev_review.get('exposure', '').strip()
                if dev_exp.lower() in ('none', 'to_be_assigned', ''):
                    result['created_jira'].append({
                        'type': 'Skipped', 'key': '', 'level': -1, 'badge': 'badge-skip',
                        'summary': f'[{comp}] SW Development Review skipped',
                        'note': f"exposure='{dev_exp}' — no KMD action required", 'url': '',
                    })
                    continue
                if dev_review['jira_key']:
                    result['created_jira'].append({
                        'type': 'Skipped', 'key': '', 'level': -1, 'badge': 'badge-skip',
                        'summary': f'[{comp}] SW Development Review already processed',
                        'note': f"jira_key={dev_review['jira_key']}", 'url': '',
                    })
                    continue

            child_hsd_id = dev_review['id'] if dev_review else ''
            jira = _build_jira_issues(
                post_issue, create_link_fn, jira_title,
                eff_platform, jira_cfg['jira_component'], jira_cfg['assignee'],
                hsd_id, child_hsd_id=child_hsd_id,
                project_tag=eff_tag,
            )
            result['created_jira'].extend(jira.get('created', []))
            result['errors'].extend(jira.get('errors', []))
            cc_key = jira.get('cc_story_key')
            result['cc_story_key'] = cc_key

            dev_id = dev_review['id'] if dev_review else _create_server_ar_dev_child(
                session, hsd_id, comp, exposure)
            if not dev_review:
                result['new_hsd_child_id'] = dev_id
                if not dev_id:
                    result['errors'].append({
                        'summary': f'Create server.ar SW Development Review [{comp}]',
                        'error':   f'Failed to create server.ar dev child under HSD {hsd_id}',
                    })
            if dev_id and cc_key:
                result['hsd_writeback_ok'] = _writeback_server_ar_jira_key(
                    session, dev_id, cc_key)
                result['writeback_field'] = 'server.ar.jira_key'
        return result

    # ── server_platf.feature ──────────────────────────────────────────────────
    if _is_server_platf_feature(info):
        ar_raw    = _get_server_platf_ar_children(session, hsd_id)
        ar_parsed = [_parse_server_platf_ar(r) for r in ar_raw]
        kmd_list  = [c for c in ar_parsed if _title_has_kmd(c['title'])]
        if not kmd_list:
            result['skipped']     = True
            result['skip_reason'] = (
                f'No server_platf.ar child with KMD in title found '
                f'({len(ar_parsed)} AR records checked)')
            return result

        matched_ar   = kmd_list[0]
        child_hsd_id = matched_ar['id']
        existing_tag = matched_ar.get('tag', '').strip()
        if existing_tag:
            result['skipped']          = True
            result['skip_reason']      = (
                f'server_platf.ar child ({child_hsd_id}) already has tag={existing_tag}')
            result['new_hsd_child_id'] = child_hsd_id
            return result

        result['release'] = info.get('release', '')
        result['new_hsd_child_id'] = child_hsd_id
        jira = _build_jira_issues(post_issue, create_link_fn, jira_title,
                                  eff_platform, eff_comp, eff_asgn,
                                  hsd_id, child_hsd_id=child_hsd_id,
                                  project_tag=eff_tag)
        result['created_jira'] = jira.get('created', [])
        result['errors'].extend(jira.get('errors', []))
        cc_key = jira.get('cc_story_key')
        result['cc_story_key'] = cc_key
        if child_hsd_id and cc_key:
            result['hsd_writeback_ok'] = _writeback_tag_field(session, child_hsd_id, cc_key)
            result['writeback_field'] = 'server_platf.ar.tag'
        return result

    # ── ip_hw_graphics.bugeco (SW WA) ─────────────────────────────────────────
    if _is_sw_wa_hsd(info):
        sw_raw    = _get_sw_impact_children(session, hsd_id)
        sw_parsed = [_parse_sw_impact_record(r) for r in sw_raw]
        kmd_wa    = [c for c in sw_parsed
                     if 'i915_kmd' in c['sw_component'].lower()
                     and c['func_impact'].lower() == 'wa_needed']
        if not kmd_wa:
            result['skipped']     = True
            result['skip_reason'] = (
                f'No sw_impact child with sw_component=i915_kmd and func_impact=wa_needed '
                f'({len(sw_parsed)} sw_impact records checked)')
            return result

        matched_sw         = kmd_wa[0]
        child_hsd_id       = matched_sw['id']
        existing_sw_record = matched_sw.get('sw_record', '').strip()
        if existing_sw_record:
            result['skipped']          = True
            result['skip_reason']      = (
                f'sw_impact child ({child_hsd_id}) already has sw_record={existing_sw_record}')
            result['new_hsd_child_id'] = child_hsd_id
            return result

        result['release'] = info.get('release', '')
        result['new_hsd_child_id'] = child_hsd_id
        jira = _build_jira_issues(post_issue, create_link_fn, jira_title,
                                  eff_platform, eff_comp, eff_asgn,
                                  hsd_id, child_hsd_id=child_hsd_id,
                                  project_tag=eff_tag)
        result['created_jira'] = jira.get('created', [])
        result['errors'].extend(jira.get('errors', []))
        cc_key = jira.get('cc_story_key')
        result['cc_story_key'] = cc_key
        if child_hsd_id and cc_key:
            result['hsd_writeback_ok'] = _writeback_sw_record(session, child_hsd_id, cc_key)
            result['writeback_field'] = 'sw_impact.sw_record'
        return result

    # ── ip_hw_graphics.feature ────────────────────────────────────────────────
    if _is_ip_hw_feature_hsd(info):
        sw_raw    = _get_sw_impact_children(session, hsd_id)
        sw_parsed = [_parse_sw_impact_record(r) for r in sw_raw]
        core_arch = [c for c in sw_parsed
                     if c['sw_component'].lower() == 'core'
                     and c['os'].lower() == 'common'
                     and c['sw_task'].lower() == 'architecture'
                     and c['sw_exposure'].lower() not in ('to_be_assigned', 'none', '')]
        if not core_arch:
            result['skipped']     = True
            result['skip_reason'] = (
                f'No sw_impact child with sw_component=core, os=common, '
                f'sw_task=architecture and valid exposure '
                f'({len(sw_parsed)} sw_impact records checked)')
            return result

        arch_exposure = core_arch[0]['sw_exposure']
        kmd_dev = [c for c in sw_parsed
                   if c['sw_component'].lower() == 'i915_kmd'
                   and c['os'].lower() == 'linux'
                   and c['sw_task'].lower() == 'development']

        result['release'] = info.get('release', '')

        if kmd_dev:
            kmd_child          = kmd_dev[0]
            child_hsd_id       = kmd_child['id']
            existing_sw_record = kmd_child.get('sw_record', '').strip()
            if existing_sw_record:
                result['new_hsd_child_id'] = child_hsd_id
                # If exposure is none — no KMD work required, skip TREND/DONE entirely
                dev_exposure = kmd_child.get('sw_exposure', '').strip()
                if dev_exposure.lower() == 'none':
                    result['skipped']     = True
                    result['skip_reason'] = (
                        f'kmd dev child ({child_hsd_id}) has sw_exposure=none — skipping TREND/DONE sync')
                    return result
                # If already marked Done — nothing left to do
                if kmd_child.get('done', '').strip().lower() == 'yes':
                    result['skipped']     = True
                    result['skip_reason'] = (
                        f'kmd dev child ({child_hsd_id}) is marked Done=Yes — no sync needed')
                    return result
                # Fetch Jira status only (trend updates handled exclusively by the weekly cron)
                jira_status, _ = _fetch_jira_status_and_trend(token, existing_sw_record)
                if jira_status.lower() == 'closed':
                    # Jira is closed → mark Done=Yes, skip TREND update
                    _update_sw_impact_field(session, child_hsd_id,
                                           'ip_hw_graphics.sw_impact.done', 'Yes')
                    result['skipped']          = True
                    result['skip_reason']      = (
                        f'Jira {existing_sw_record} is Closed — set Done=Yes on child {child_hsd_id}')
                    result['hsd_writeback_ok'] = True
                    result['writeback_field']  = 'sw_impact.done=Yes'
                    return result
                # Jira still open — skip (trend sync handled by weekly cron)
                result['skipped']          = True
                result['skip_reason']      = (
                    f'kmd dev child ({child_hsd_id}) already has sw_record={existing_sw_record}'
                    ' — trend updates handled by weekly cron')
                result['hsd_writeback_ok'] = False
                result['writeback_field']  = ''
                return result
            # sw_record is empty — check exposure before creating Jira
            dev_exposure = kmd_child.get('sw_exposure', '').strip()
            if dev_exposure.lower() in ('none', 'to_be_assigned', ''):
                result['skipped']          = True
                result['skip_reason']      = (
                    f"kmd dev child ({child_hsd_id}) has sw_exposure='{dev_exposure}' and no sw_record — no KMD action needed")
                result['new_hsd_child_id'] = child_hsd_id
                return result
            result['new_hsd_child_id'] = child_hsd_id
        else:
            child_hsd_id = _create_sw_impact_child(
                session, parent_id=hsd_id, sw_exposure=arch_exposure, owner=_ASSIGNEE)
            result['new_hsd_child_id'] = child_hsd_id
            if not child_hsd_id:
                result['errors'].append({
                    'summary': 'New ip_hw_graphics.sw_impact child creation',
                    'error':   f'Failed to create sw_impact child under HSD {hsd_id}',
                })

        jira = _build_jira_issues(post_issue, create_link_fn, jira_title,
                                  eff_platform, eff_comp, eff_asgn,
                                  hsd_id, child_hsd_id=child_hsd_id or '',
                                  project_tag=eff_tag)
        result['created_jira'] = jira.get('created', [])
        result['errors'].extend(jira.get('errors', []))
        cc_key = jira.get('cc_story_key')
        result['cc_story_key'] = cc_key
        if child_hsd_id and cc_key:
            result['hsd_writeback_ok'] = _writeback_sw_record(session, child_hsd_id, cc_key)
            result['writeback_field'] = 'sw_impact.sw_record'
        return result

    # ── dg_soc.feature ────────────────────────────────────────────────────────
    if not _is_feature_hsd(info):
        result['skipped']     = True
        result['skip_reason'] = (
            f"Tenant/subject '{info.get('tenant','')}.{info.get('subject','')}' is not supported")
        return result

    ar_raw      = _get_ar_children(session, hsd_id)
    ar_children = [_parse_ar_record_extended(r) for r in ar_raw]
    arch_list    = [c for c in ar_children
                    if c['team'].lower() == 'i915_kmd'
                    and c['task'].lower() == 'architecture']
    kmd_dev_list = [c for c in ar_children
                    if c['team'].lower() == 'i915_kmd'
                    and c['task'].lower() == 'development']
    ual_dev_list = [c for c in ar_children
                    if c['team'].lower() == 'sw.ual_kmd'
                    and c['task'].lower() == 'development']

    if not arch_list and not kmd_dev_list and not ual_dev_list:
        result['skipped']     = True
        result['skip_reason'] = (f'No AR child with team=i915_kmd or sw.ual_kmd found '
                                 f'({len(ar_children)} AR records checked)')
        return result

    result['release'] = info.get('release', '')
    kmd_acted = False

    # ── i915_kmd processing ───────────────────────────────────────────────────
    if arch_list or kmd_dev_list:
        if arch_list:
            arch          = arch_list[0]
            exposure      = arch.get('exposure', '').strip()
            arch_title    = arch.get('title') or title
            arch_exposure = arch.get('exposure', '')
        else:
            dev_direct    = kmd_dev_list[0]
            exposure      = dev_direct.get('exposure', '').strip()
            arch_title    = dev_direct.get('title') or title
            arch_exposure = dev_direct.get('exposure', '')

        if exposure.lower() in ('none', 'to_be_assigned'):
            result['created_jira'].append({
                'type': 'Skipped', 'key': '', 'level': -1, 'badge': 'badge-skip',
                'summary': f"[i915_kmd] skipped — exposure='{exposure}'",
                'note': 'no KMD action required', 'url': '',
            })
        else:
            existing_dev_id = None
            if kmd_dev_list:
                existing_dev      = kmd_dev_list[0]
                existing_dev_id   = existing_dev['id']
                dev_exposure      = existing_dev.get('exposure', '').strip()
                existing_jira_key = existing_dev.get('jira_key', '').strip()
                if dev_exposure.lower() == 'none':
                    result['created_jira'].append({
                        'type': 'Skipped', 'key': '', 'level': -1, 'badge': 'badge-skip',
                        'summary': f"[i915_kmd] skipped — dev exposure='none'",
                        'note': 'no KMD action required', 'url': '',
                    })
                    existing_dev_id = None  # don't create Jira
                elif existing_jira_key:
                    if dev_exposure.lower() in ('to_be_assigned', '') and arch_exposure:
                        try:
                            resp = session.put(
                                f'{_HSD_REST}/{existing_dev_id}',
                                json={'tenant': 'dg_soc', 'subject': 'ar',
                                      'fieldValues': [{'dg_soc.ar.exposure': arch_exposure}]},
                                headers={'Content-Type': 'application/json'}, timeout=20,
                            )
                            if resp.status_code not in (200, 201, 204):
                                logger.warning('Exposure update returned %s for child %s',
                                               resp.status_code, existing_dev_id)
                        except Exception as exc:
                            logger.error('Exposure update failed for child %s: %s', existing_dev_id, exc)
                    result['created_jira'].append({
                        'type': 'Skipped', 'key': '', 'level': -1, 'badge': 'badge-skip',
                        'summary': f'[i915_kmd] already has jira_key={existing_jira_key}',
                        'note': 'skipping', 'url': f'{_JIRA_BASE}/browse/{existing_jira_key}',
                    })
                    result['new_hsd_child_id'] = existing_dev_id
                    existing_dev_id = None  # don't create another Jira

            if existing_dev_id is not None or not kmd_dev_list:
                # create Jira and writeback
                new_hsd_id = existing_dev_id if existing_dev_id else None
                if not new_hsd_id and kmd_dev_list:
                    new_hsd_id = kmd_dev_list[0]['id']

                if not new_hsd_id:
                    if arch_exposure:
                        try:
                            resp = session.put(
                                f'{_HSD_REST}/{existing_dev_id}',
                                json={'tenant': 'dg_soc', 'subject': 'ar',
                                      'fieldValues': [{'dg_soc.ar.exposure': arch_exposure}]},
                                headers={'Content-Type': 'application/json'}, timeout=20,
                            )
                        except Exception:
                            pass
                    new_hsd_id = _create_hsd_ar_child(
                        session, parent_id=hsd_id, title=arch_title,
                        type_='sw_assessment', team='i915_kmd',
                        exposure=arch_exposure, task='development', jira_key='',
                    )
                    result['new_hsd_child_id'] = new_hsd_id
                    if not new_hsd_id:
                        result['errors'].append({
                            'summary': 'New HSD child creation',
                            'error':   f'Failed to create dg_soc.ar child under HSD {hsd_id}',
                        })
                else:
                    result['new_hsd_child_id'] = new_hsd_id
                    if arch_exposure:
                        try:
                            session.put(
                                f'{_HSD_REST}/{new_hsd_id}',
                                json={'tenant': 'dg_soc', 'subject': 'ar',
                                      'fieldValues': [{'dg_soc.ar.exposure': arch_exposure}]},
                                headers={'Content-Type': 'application/json'}, timeout=20,
                            )
                        except Exception as exc:
                            logger.error('Exposure update failed for child %s: %s', new_hsd_id, exc)

                jira = _build_jira_issues(post_issue, create_link_fn, jira_title,
                                          eff_platform, eff_comp, eff_asgn,
                                          hsd_id, child_hsd_id=new_hsd_id or '',
                                          project_tag=eff_tag)
                result['created_jira'].extend(jira.get('created', []))
                result['errors'].extend(jira.get('errors', []))
                cc_key = jira.get('cc_story_key')
                result['cc_story_key'] = cc_key
                if new_hsd_id and cc_key:
                    result['hsd_writeback_ok'] = _writeback_jira_key(session, new_hsd_id, cc_key)
                    result['writeback_field']   = 'dg_soc.ar.jira_key'
                kmd_acted = True

    # ── sw.ual_kmd processing ─────────────────────────────────────────────────
    if ual_dev_list:
        ual_dev      = ual_dev_list[0]
        ual_exp      = ual_dev.get('exposure', '').strip()
        ual_jira_key = ual_dev.get('jira_key', '').strip()
        ual_dev_id   = ual_dev['id']
        if ual_jira_key:
            result['created_jira'].append({
                'type': 'Skipped', 'key': '', 'level': -1, 'badge': 'badge-skip',
                'summary': '[sw.ual_kmd] development already has jira_key',
                'note': f'jira_key={ual_jira_key}', 'url': '',
            })
        elif ual_exp.lower() in ('none', 'to_be_assigned', ''):
            result['created_jira'].append({
                'type': 'Skipped', 'key': '', 'level': -1, 'badge': 'badge-skip',
                'summary': '[sw.ual_kmd] development skipped',
                'note': f"exposure='{ual_exp}' — no UAL action required", 'url': '',
            })
        else:
            ual_jira = _build_jira_issues(
                post_issue, create_link_fn, jira_title,
                eff_platform, 'Kernel - UAL', 'jzuo2',
                hsd_id, child_hsd_id=ual_dev_id,
                project_tag=eff_tag,
            )
            result['created_jira'].extend(ual_jira.get('created', []))
            result['errors'].extend(ual_jira.get('errors', []))
            ual_cc_key = ual_jira.get('cc_story_key')
            if ual_dev_id and ual_cc_key:
                _writeback_jira_key(session, ual_dev_id, ual_cc_key)
            kmd_acted = True

    if not kmd_acted:
        result['skipped']     = True
        result['skip_reason'] = 'All AR children skipped (exposure=none/to_be_assigned or jira_key already set)'

    return result


# ── Cron job infrastructure ────────────────────────────────────────────────────

_CRON_CONFIG_PATH  = _BASE_DIR / 'hsd2jira_cron.json'
_CRON_LOG_PATH     = _BASE_DIR / 'hsd2jira_cron_log.json'
_CRON_LOG_MAX_RUNS = 5
_CRON_DEFAULT_QUERIES = [
    # ip_hw_graphics.feature — GT queries (CRI / JGS / TGS)
    'https://hsdes.intel.com/appstore/generalapps/#/pages/community/query?queryId=14024748599',
    'https://hsdes.intel.com/appstore/community_legacy/#/1206341217?queryId=16026920807&articleId=14024704478',
    'https://hsdes.intel.com/appstore/generalapps/#/pages/community/1206341217?queryId=14026130398',
    # server.feature — JGS SOC
    'https://hsdes.intel.com/appstore/generalapps/#/pages/community/22020180134?queryId=14028056822',
    # dg_soc.feature — CRI / TGS SOC
    'https://hsdes.intel.com/appstore/generalapps/#/pages/community/14024672201?queryId=14024897996',
    'https://hsdes.intel.com/appstore/generalapps/#/pages/community/14026922422?queryId=14027898530',
    # server_platf.feature — CRI CCB
    'https://hsdes.intel.com/appstore/generalapps/#/pages/community/query?queryId=13013902803',
    # ip_hw_graphics.bugeco — CRI WA
    'https://hsdes.intel.com/appstore/generalapps/#/pages/community/1206341217?queryId=16029952357',
]

# Maps a queryId → (platform, project_tag) for ip_hw_graphics.feature queries.
# When a query is listed here the cron uses the mapped values directly and
# skips the family_affected / release_affected field detection entirely.
_QUERY_PROJECT_MAP: dict[str, tuple[str, str]] = {
    '14024748599': ('crescent_island', 'CRI'),   # ip_hw_graphics.feature CRI
    '16026920807': ('Jaguar Shores',   'JGS'),   # ip_hw_graphics.feature JGS
    '14026130398': ('tiger_shores',    'TGS'),   # ip_hw_graphics.feature TGS
    '16029952357': ('crescent_island', 'CRI'),   # ip_hw_graphics.bugeco  CRI
}
_cron_stop_event  = threading.Event()
_cron_thread_ref  = [None]   # mutable container so inner functions can update
_cron_lock        = threading.Lock()


def _load_cron_config() -> dict:
    try:
        if _CRON_CONFIG_PATH.exists():
            with open(_CRON_CONFIG_PATH) as f:
                return json.load(f)
    except Exception:
        pass
    return {
        'enabled':          False,
        'schedule':         'weekly_friday',  # 'weekly_friday' | 'interval'
        'interval_minutes': 30,               # used only when schedule='interval'
        'query_urls':       _CRON_DEFAULT_QUERIES[:],
        'last_run':         None,
        'last_result':      None,
    }


def _save_cron_config(cfg: dict) -> None:
    try:
        with open(_CRON_CONFIG_PATH, 'w') as f:
            json.dump(cfg, f, indent=2)
    except Exception as exc:
        logger.error('Failed to save cron config: %s', exc)


def _append_cron_log(run_time: str, summary: dict, details: list) -> None:
    """Prepend a new run entry to the rolling cron log (keeps last _CRON_LOG_MAX_RUNS)."""
    try:
        log = []
        if _CRON_LOG_PATH.exists():
            with open(_CRON_LOG_PATH) as f:
                log = json.load(f)
    except Exception:
        log = []
    log.insert(0, {'run_time': run_time, 'summary': summary, 'details': details})
    log = log[:_CRON_LOG_MAX_RUNS]
    try:
        with open(_CRON_LOG_PATH, 'w') as f:
            json.dump(log, f, indent=2)
    except Exception as exc:
        logger.error('Failed to write cron log: %s', exc)


def _cron_run_pass(cfg: dict) -> dict:
    """Execute one cron pass over all configured query URLs.
    Returns {'summary': {...}, 'details': [per-HSD result dicts]}.

    For queries listed in _QUERY_PROJECT_MAP the platform/tag is taken directly
    from the map — no family_affected / release_affected detection needed.
    For all other queries the existing _detect_platform_and_tag() fallback is used.
    """
    query_urls = cfg.get('query_urls', _CRON_DEFAULT_QUERIES)
    token, err = _token_response()
    if err:
        return {'summary': {'error': 'token read failed', 'processed': 0, 'skipped': 0, 'errors': 1, 'total': 0}, 'details': []}

    post_issue, create_link_fn = _make_jira_creator(token)
    session = _new_hsd_session()

    # Build a mapping: hsd_id → (platform, tag) using query origin where known.
    # If a record appears in multiple queries (edge case) the last mapping wins.
    id_to_project: dict[str, tuple[str, str]] = {}
    for qurl in query_urls:
        m = re.search(r'queryId=(\d+)', qurl)
        if not m:
            continue
        query_id = m.group(1)
        forced   = _QUERY_PROJECT_MAP.get(query_id)   # (platform, tag) or None
        ids = _execute_hsd_query(session, query_id) or []
        for hid in ids:
            if forced:
                id_to_project[hid] = forced
            elif hid not in id_to_project:     # don't overwrite a forced mapping
                id_to_project[hid] = ('', '')  # will be resolved via field detection below

    processed = skipped = error_count = 0
    details: list = []

    for hsd_id, (platform, tag) in id_to_project.items():
        try:
            info = _fetch_hsd_info(session, hsd_id)
            if not platform:
                # No query-based mapping — fall back to field detection
                platform, tag = _detect_platform_and_tag(info)
            if not platform:
                skipped += 1
                details.append({
                    'hsd_id':      hsd_id,
                    'title':       info.get('title', f'HSD {hsd_id}'),
                    'tenant':      f"{info.get('tenant','')}.{info.get('subject','')}".strip('.'),
                    'url':         f'https://hsdes.intel.com/appstore/article-one/#/{hsd_id}',
                    'skipped':     True,
                    'skip_reason': 'Platform/tag could not be determined from query map or HSD fields',
                    'created_jira': [], 'errors': [],
                })
                continue
            r = _process_one_hsd(session, post_issue, create_link_fn,
                                  hsd_id, info,
                                  platform=platform, project_tag=tag,
                                  token=token)
            details.append(r)
            if r.get('skipped'):
                skipped += 1
            else:
                processed += 1
                error_count += len(r.get('errors', []))
        except Exception as exc:
            logger.error('Cron: error processing HSD %s: %s', hsd_id, exc)
            error_count += 1
            details.append({
                'hsd_id': hsd_id, 'title': f'HSD {hsd_id}',
                'skipped': False, 'skip_reason': None,
                'created_jira': [],
                'errors': [{'summary': 'Cron exception', 'error': str(exc)}],
            })

    summary = {'processed': processed, 'skipped': skipped, 'errors': error_count, 'total': len(id_to_project)}
    return {'summary': summary, 'details': details}


def _cron_worker(stop_event: threading.Event) -> None:
    logger.info('HSD2Jira cron worker started')
    while not stop_event.is_set():
        cfg = _load_cron_config()

        if not cfg.get('enabled'):
            stop_event.wait(60)
            continue

        now       = datetime.utcnow()
        last_run  = cfg.get('last_run', '') or ''
        today_str = now.strftime('%Y-%m-%d')
        if now.weekday() != 4:              # 4 = Friday
            stop_event.wait(3600)           # not Friday — check again in 1 h
            continue
        if last_run.startswith(today_str):  # already ran this Friday
            stop_event.wait(3600)
            continue

            logger.info('Cron: starting processing pass')
        try:
            result   = _cron_run_pass(cfg)
            run_time = datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')
            summary  = result.get('summary', result)
            details  = result.get('details', [])
            _append_cron_log(run_time, summary, details)
            cfg = _load_cron_config()
            cfg['last_run']    = run_time
            cfg['last_result'] = summary
            _save_cron_config(cfg)
            for item in details:
                tenant   = item.get('tenant', 'unknown')
                hid      = item.get('hsd_id', '?')
                wb_field = item.get('writeback_field', '')
                jiras    = ' '.join(j.get('key', '') for j in item.get('created_jira', []) if j.get('key'))
                child    = item.get('new_hsd_child_id', '')
                errs     = item.get('errors', [])
                # Only emit a log line when something actually changed or failed
                if not (jiras or child or item.get('hsd_writeback_ok') or errs):
                    continue
                parts = []
                if jiras:    parts.append(f'Jira: {jiras}')
                if child:    parts.append(f'child: {child}')
                if wb_field: parts.append(f'wrote {wb_field}')
                elif item.get('hsd_writeback_ok'): parts.append('writeback OK')
                if errs:     parts.append('ERRORS: ' + '; '.join((e.get('error') or e.get('summary', '?'))[:80] for e in errs[:3]))
                logger.info('Cron [%s] HSD %s -> %s', tenant, hid, ' | '.join(parts))
            logger.info('Cron: pass complete — %s', summary)
        except Exception as exc:
            logger.exception('Cron: pass failed: %s', exc)

        stop_event.wait(3600)   # done for today; check again in 1 h (blocked by last_run check)
    logger.info('HSD2Jira cron worker stopped')


def _ensure_cron_thread(enabled: bool) -> None:
    with _cron_lock:
        if enabled:
            if _cron_thread_ref[0] is None or not _cron_thread_ref[0].is_alive():
                _cron_stop_event.clear()
                t = threading.Thread(target=_cron_worker, args=(_cron_stop_event,),
                                     daemon=True, name='hsd2jira-cron')
                t.start()
                _cron_thread_ref[0] = t
                logger.info('Cron thread started')
        else:
            _cron_stop_event.set()
            logger.info('Cron thread stop requested')


# ── Routes ─────────────────────────────────────────────────────────────────────

@bp.route('/hsd2jira')
def hsd2jira():
    return render_template('hsd2jira.html')


@bp.route('/hsd2jira/explain')
def hsd2jira_explain():
    return render_template('hsd2jira_explain.html')


@bp.route('/hsd2jira/explain/slides')
def hsd2jira_explain_slides():
    return render_template('hsd2jira_explain_slides.html')


@bp.route('/hsd2jira/explain/download-pptx')
def hsd2jira_explain_pptx():
    """Generate and return a .pptx explaining the HSD2Jira decision logic."""
    import io, base64, json as _json
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    BG       = RGBColor(0x0f, 0x0f, 0x1a)
    PURPLE   = RGBColor(0xa7, 0x8b, 0xfa)
    PURPLE_D = RGBColor(0x7c, 0x3a, 0xed)
    WHITE    = RGBColor(0xff, 0xff, 0xff)
    DIM      = RGBColor(0x94, 0x91, 0xb8)
    GREEN    = RGBColor(0x86, 0xef, 0xac)
    RED      = RGBColor(0xfc, 0xa5, 0xa5)
    BLUE     = RGBColor(0x93, 0xc5, 0xfd)

    W, H = Inches(13.33), Inches(7.5)   # 16:9 widescreen

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    blank_layout = prs.slide_layouts[6]

    def _bg(slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = BG

    def _txbox(slide, text, l, t, w, h, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        p  = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        return tb

    def _fetch_mermaid_png(diagram: str) -> bytes | None:
        """Fetch PNG from mermaid.ink for the given diagram code."""
        import os as _os
        _proxies = {}
        for _k in ('https_proxy', 'HTTPS_PROXY', 'http_proxy', 'HTTP_PROXY'):
            if _os.environ.get(_k):
                _proxies = {'http': _os.environ.get('http_proxy') or _os.environ.get('HTTP_PROXY'),
                            'https': _os.environ.get('https_proxy') or _os.environ.get('HTTPS_PROXY')}
                break
        if not _proxies:
            _proxies = {'http': 'http://proxy-chain.intel.com:911',
                        'https': 'http://proxy-chain.intel.com:912'}
        try:
            payload = base64.urlsafe_b64encode(
                _json.dumps({'code': diagram, 'mermaid': {'theme': 'default'}}).encode()
            ).decode().rstrip('=')
            r = requests.get(
                f'https://mermaid.ink/img/{payload}',
                timeout=20, verify=False,
                headers={'User-Agent': 'HSD2Jira/1.0'},
                proxies=_proxies,
            )
            if r.ok and r.content:
                return r.content
        except Exception as exc:
            logger.warning('mermaid.ink fetch failed: %s', exc)
        return None

    def _table(slide, rows_data, l, t, w, h):
        nrows = len(rows_data) + 1
        tbl   = slide.shapes.add_table(nrows, 3, l, t, w, h).table
        for i, cw in enumerate([Inches(2.0), Inches(4.0), Inches(1.8)]):
            tbl.columns[i].width = cw

        def _cell(r, c, text, hdr=False, result=False):
            cell = tbl.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                RGBColor(0x2d, 0x1e, 0x5e) if hdr else
                (RGBColor(0x14, 0x12, 0x2a) if r % 2 == 0 else RGBColor(0x1a, 0x17, 0x33))
            )
            tf  = cell.text_frame
            tf.word_wrap = True
            run = tf.paragraphs[0].add_run()
            run.text = text
            run.font.size = Pt(8) if hdr else Pt(9)
            run.font.bold = hdr
            if hdr:
                run.font.color.rgb = PURPLE
            elif result:
                lo = text.lower()
                run.font.color.rgb = RED if 'skip' in lo else (GREEN if 'create' in lo else BLUE)
            else:
                run.font.color.rgb = WHITE

        for c, h_txt in enumerate(('CHECK', 'CONDITION', 'RESULT')):
            _cell(0, c, h_txt, hdr=True)
        for r, (chk, cond, res) in enumerate(rows_data, start=1):
            _cell(r, 0, chk); _cell(r, 1, cond); _cell(r, 2, res, result=True)

    # ── Slide data ────────────────────────────────────────────────────────────
    slides_data = [
        {
            'label':   'JGS / CRI / TGS GT queries',
            'title':   'ip_hw_graphics.feature',
            'queries': 'CRI 14024748599  ·  JGS 16026920807  ·  TGS 14026130398',
            'write':   'Writeback: sw_impact.sw_record  on  i915_kmd / linux / development  child',
            'diagram': '''flowchart TD
    A([HSD from query]) --> B{ip_hw_graphics\\n.feature?}
    B -- No --> S1([⏭ Skip])
    B -- Yes --> C{core/common/arch\\nchild + valid exposure?}
    C -- No --> S2([⏭ Skip])
    C -- Yes --> D{i915_kmd/linux/\\ndevelopment child?}
    D -- No --> E[Create child] --> J
    D -- Yes --> F{sw_exposure =\\nnone / TBA?}
    F -- Yes --> S3([⏭ Skip])
    F -- No --> G{sw_record set?}
    G -- No --> J[🔨 Create Jira]
    G -- Yes --> H{done = Yes?}
    H -- Yes --> S4([⏭ Skip])
    H -- No --> I{Jira Closed?}
    I -- Yes --> W1([✏ Write done=Yes])
    I -- No --> W2([✏ Write trend=WW])
    J --> W3([✏ CC key → sw_record])''',
            'rows': [
                ('HSD type',      'tenant=ip_hw_graphics, subject=feature',       'SKIP if not'),
                ('Gate child',    'core/common/arch sw_impact, valid exposure',    'SKIP if missing'),
                ('KMD exposure',  'i915_kmd/linux/dev: exposure=none/TBA',         'SKIP'),
                ('Already done',  'sw_record set AND done=Yes',                    'SKIP'),
                ('Jira closed',   'sw_record set, Jira is Closed',                 'Write done=Yes'),
                ('Trend sync',    'sw_record set, Jira open + Actual Trend WW',    'Write trend=WW'),
                ('Needs Jira',    'sw_record empty, exposure valid',               'CREATE Jira'),
                ('No KMD child',  'i915_kmd/linux/dev child missing',              'CREATE child + Jira'),
            ],
        },
        {
            'label':   'CRI WA',
            'title':   'ip_hw_graphics.bugeco',
            'queries': 'CRI 16029952357',
            'write':   'Writeback: sw_impact.sw_record  on  i915_kmd / wa_needed  child',
            'diagram': '''flowchart TD
    A([HSD from query]) --> B{ip_hw_graphics\\n.bugeco?}
    B -- No --> S1([⏭ Skip])
    B -- Yes --> C{sw_impact child\\ni915_kmd / wa_needed?}
    C -- No --> S2([⏭ Skip])
    C -- Yes --> D{sw_record\\nalready set?}
    D -- Yes --> S3([⏭ Skip])
    D -- No --> E[🔨 Create Jira]
    E --> W([✏ CC key → sw_record])''',
            'rows': [
                ('HSD type',          'tenant=ip_hw_graphics, subject=bugeco', 'SKIP if not'),
                ('WA child',          'sw_impact: i915_kmd + wa_needed',       'SKIP if missing'),
                ('Already processed', 'sw_record already set',                  'SKIP'),
                ('Ready',             'sw_record empty',                        'CREATE Jira'),
            ],
        },
        {
            'label':   'JGS SOC query',
            'title':   'server.feature',
            'queries': 'JGS 14028056822',
            'write':   'Writeback: server.ar.jira_key  on  SW Development Review  child',
            'diagram': '''flowchart TD
    A([HSD from query]) --> B{server\\n.feature?}
    B -- No --> S1([⏭ Skip])
    B -- Yes --> C{SW Arch Review\\nchild + valid exposure?}
    C -- No --> S2([⏭ Skip])
    C -- Yes --> D[For each Arch\\nReview child]
    D --> E{SW Dev Review\\nchild exists?}
    E -- No --> F[Create dev child] --> J
    E -- Yes --> G{dev exposure =\\nnone / TBA?}
    G -- Yes --> S3([⏭ Skip component])
    G -- No --> H{jira_key set?}
    H -- Yes --> S4([⏭ Skip component])
    H -- No --> J[🔨 Create Jira]
    J --> W([✏ CC key → jira_key])''',
            'rows': [
                ('HSD type',          'tenant=server, subject=feature',              'SKIP if not'),
                ('Gate child',        'SW Arch Review, valid component + exposure',  'SKIP if missing'),
                ('Dev exposure',      'SW Dev Review: exposure=none/TBA',             'SKIP component'),
                ('Already processed', 'jira_key set on SW Dev Review child',          'SKIP component'),
                ('No dev child',      'SW Development Review missing',               'CREATE child + Jira'),
                ('Ready',             'dev child exists, jira_key empty',            'CREATE Jira'),
            ],
        },
        {
            'label':   'CRI / TGS SOC query',
            'title':   'dg_soc.feature',
            'queries': 'CRI 14024897996  ·  TGS 14027898530',
            'write':   'Writeback: dg_soc.ar.jira_key  on  i915_kmd / development  AR child',
            'diagram': '''flowchart TD
    A([HSD from query]) --> B{dg_soc\\n.feature?}
    B -- No --> S1([⏭ Skip])
    B -- Yes --> C{i915_kmd/arch\\nAR child?}
    C -- No --> S2([⏭ Skip])
    C -- Yes --> D{arch exposure =\\nnone / TBA?}
    D -- Yes --> S3([⏭ Skip])
    D -- No --> E{i915_kmd/dev\\nAR child?}
    E -- No --> F[Create dev child] --> J
    E -- Yes --> G{dev exposure =\\nnone / TBA?}
    G -- Yes --> S4([⏭ Skip])
    G -- No --> H{jira_key set?}
    H -- Yes --> S5([⏭ Skip])
    H -- No --> J[🔨 Create Jira]
    J --> W([✏ CC key → jira_key])''',
            'rows': [
                ('HSD type',          'tenant=dg_soc, subject=feature',             'SKIP if not'),
                ('Gate child',        'AR child: i915_kmd / architecture',           'SKIP if missing'),
                ('Arch exposure',     'exposure=none or to_be_assigned',             'SKIP'),
                ('Dev exposure',      'i915_kmd/dev AR child exposure=none/TBA',     'SKIP'),
                ('Already processed', 'jira_key set on dev AR child',                'SKIP'),
                ('No dev child',      'i915_kmd/dev AR child missing',              'CREATE child + Jira'),
                ('Ready',             'dev child exists, jira_key empty',            'CREATE Jira'),
            ],
        },
        {
            'label':   'CRI CCB',
            'title':   'server_platf.feature',
            'queries': 'CRI 13013902803',
            'write':   'Writeback: server_platf.ar.tag  on  KMD AR child',
            'diagram': '''flowchart TD
    A([HSD from query]) --> B{server_platf\\n.feature?}
    B -- No --> S1([⏭ Skip])
    B -- Yes --> C{server_platf.ar\\nchild with KMD in title?}
    C -- No --> S2([⏭ Skip])
    C -- Yes --> D{tag field\\nalready set?}
    D -- Yes --> S3([⏭ Skip])
    D -- No --> E[🔨 Create Jira]
    E --> W([✏ CC key → tag])''',
            'rows': [
                ('HSD type',          'tenant=server_platf, subject=feature',    'SKIP if not'),
                ('KMD child',         'server_platf.ar child with KMD in title', 'SKIP if missing'),
                ('Already processed', 'tag field already set',                    'SKIP'),
                ('Ready',             'tag empty',                                'CREATE Jira'),
            ],
        },
    ]

    # Fetch diagrams sequentially (proxy throttles concurrent connections)
    diagram_pngs = [_fetch_mermaid_png(sd['diagram']) for sd in slides_data]

    # ── Title slide ───────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank_layout)
    _bg(sl)
    bar = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.08), H)
    bar.fill.solid(); bar.fill.fore_color.rgb = PURPLE_D; bar.line.fill.background()
    _txbox(sl, 'HSD2Jira', Inches(0.35), Inches(1.8), Inches(12), Inches(1.0),
           size=44, bold=True)
    _txbox(sl, 'Jira Creation — Decision Logic', Inches(0.35), Inches(2.9),
           Inches(12), Inches(0.6), size=24, color=PURPLE)
    _txbox(sl, 'For each HSD returned by the cron queries, the tool decides whether to\n'
               'create a 9-issue Jira hierarchy or skip, based on child record state.',
           Inches(0.35), Inches(3.7), Inches(10), Inches(0.9), size=14, color=DIM)
    _txbox(sl, '5 tenants  ·  8 queries  ·  Idempotent — safe to re-run at any time',
           Inches(0.35), Inches(4.8), Inches(10), Inches(0.5), size=12, color=DIM)

    # ── Tenant slides ─────────────────────────────────────────────────────────
    for sd, diagram_png in zip(slides_data, diagram_pngs):
        sl = prs.slides.add_slide(blank_layout)
        _bg(sl)

        bar = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.08), H)
        bar.fill.solid(); bar.fill.fore_color.rgb = PURPLE_D; bar.line.fill.background()

        hdr = sl.shapes.add_shape(1, Inches(0.08), Inches(0), W - Inches(0.08), Inches(1.55))
        hdr.fill.solid(); hdr.fill.fore_color.rgb = RGBColor(0x18, 0x14, 0x30)
        hdr.line.fill.background()

        _txbox(sl, sd['label'].upper(), Inches(0.25), Inches(0.1),
               Inches(9), Inches(0.28), size=8, color=DIM)
        _txbox(sl, sd['title'], Inches(0.25), Inches(0.36),
               Inches(9), Inches(0.65), size=26, bold=True)
        _txbox(sl, sd['queries'], Inches(0.25), Inches(1.0),
               Inches(7), Inches(0.35), size=11, color=PURPLE)
        _txbox(sl, sd['write'], Inches(0.25), Inches(1.28),
               Inches(9), Inches(0.28), size=9, color=DIM)

        # ── Flowchart diagram (left) ──────────────────────────────────────────
        if diagram_png:
            img_buf = io.BytesIO(diagram_png)
            sl.shapes.add_picture(img_buf, Inches(0.2), Inches(1.65),
                                  width=Inches(5.8), height=Inches(5.65))
            _txbox(sl, 'FLOWCHART', Inches(0.25), Inches(1.58),
                   Inches(2.5), Inches(0.2), size=7, bold=True,
                   color=RGBColor(0x55, 0x4a, 0x80))
        else:
            _txbox(sl, '⚠ Diagram unavailable (mermaid.ink unreachable)',
                   Inches(0.25), Inches(1.75), Inches(5.6), Inches(0.5),
                   size=10, color=RED)

        # ── Conditions table (right) ──────────────────────────────────────────
        _txbox(sl, 'CONDITIONS', Inches(6.2), Inches(1.58),
               Inches(2.5), Inches(0.2), size=7, bold=True,
               color=RGBColor(0x55, 0x4a, 0x80))
        tbl_h = Inches(0.38 * (len(sd['rows']) + 1))
        _table(sl, sd['rows'], Inches(6.2), Inches(1.82), Inches(6.9), tbl_h)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return send_file(
        buf,
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
        as_attachment=True,
        download_name='HSD2Jira_Explain.pptx',
    )



# ══════════════════════════════════════════════════════════════════════════════
# xe5 DCN Enabling Status — Excel Generator
# Mirrors the JGS "GT DCNs" + "DO NOT MODIFY" two-sheet structure.
# ══════════════════════════════════════════════════════════════════════════════

_XE5_DCN_QUERY_ID  = '14028312865'
_HSD_ARTICLE_URL   = 'https://hsdes.intel.com/appstore/article-one/#/article/'
_JIRA_BROWSE_URL   = 'https://jira.devtools.intel.com/browse/'

# Cache for /xe5-dcn-status/data.json  (avoids re-fetching all HSD data on every refresh)
_XE5_JSON_CACHE: dict = {'rows': None, 'ts': 0.0, 'lock': None}
_XE5_JSON_TTL   = 600   # seconds — 10 minutes

def _xe5_cache_lock():
    import threading
    if _XE5_JSON_CACHE['lock'] is None:
        _XE5_JSON_CACHE['lock'] = threading.Lock()
    return _XE5_JSON_CACHE['lock']

# (sw_component keyword, column key, display label)
_XE5_DEV_COMPONENTS = [
    ('compute',        'compute', 'Compute'),
    ('igc_compute',    'igc',     'IGC Compute'),
    ('i915_kmd',       'kmd',     'XeKMD'),
    ('pisa_finalizer', 'pisa',    'PISA Finalizer'),
    ('uc_global',      'guc',     'GUC'),
]

# DATA sheet column map: key → (jira_col, status_col, trend_col, done_col)
_XE5_DATA_COL = {
    'compute': ('E', 'F', 'G', 'H'),
    'igc':     ('I', 'J', 'K', 'L'),
    'kmd':     ('M', 'N', 'O', 'P'),
    'pisa':    ('Q', 'R', 'S', 'T'),
    'guc':     ('U', 'V', 'W', 'X'),
}

# Architecture sw_component keyword that gates each dev column key.
# If arch exposure is valid but dev exposure is to_be_assigned → "Pending engineering scoping".
# If arch exposure is to_be_assigned/none/empty itself → "Pending arch scoping".
_XE5_DEV_ARCH_KW = {
    'compute': ('compute',),
    'igc':     ('compiler',),
    'kmd':     ('core',),
    'pisa':    ('compiler', 'pisa_finalizer'),  # either arch component gates pisa_finalizer/dev
    'guc':     ('guc',),
}
_XE5_NO_SCOPE = ('', 'none', 'to_be_assigned')  # exposures that mean "not yet scoped"


def _fetch_xe5_parent_row(token: str, hsd_id: str) -> dict:
    """Fetch one parent DCN: info + sw_impact children + Jira status per component."""
    session = _new_hsd_session()
    try:
        info  = _fetch_hsd_info(session, hsd_id)
        title = info.get('title', f'HSD {hsd_id}')
        sw_raw    = _get_sw_impact_children(session, hsd_id)
        sw_parsed = [_parse_sw_impact_record(r) for r in sw_raw]

        comp_data: dict[str, dict] = {}
        for sw_kw, key, _ in _XE5_DEV_COMPONENTS:
            arch_kw = _XE5_DEV_ARCH_KW[key]

            # Dev children with valid exposure OR already marked done
            dev_valid = [
                c for c in sw_parsed
                if c.get('sw_component', '').lower() == sw_kw
                and c.get('sw_task', '').lower() == 'development'
                and (c.get('sw_exposure', '').lower() not in _XE5_NO_SCOPE
                     or c.get('done', '').lower() in ('yes', 'true', '1'))
            ]

            if dev_valid:
                # Normal path: engineering has scoped, Jira may or may not exist
                sw_record = (dev_valid[0].get('sw_record') or '').strip()
                if not sw_record:
                    sw_done_hsd = dev_valid[0].get('done', '').lower() in ('yes', 'true', '1')
                    hsd_trend   = dev_valid[0].get('trend', '').strip()
                    comp_data[key] = {'jira': '(blank)',
                                      'status': f"No Jira ({dev_valid[0].get('sw_exposure','')})",
                                      'trend': hsd_trend or '(blank)',
                                      'done': 'yes' if sw_done_hsd else '(blank)'}
                else:
                    jira_status, trend_ww = _fetch_jira_status_and_trend(token, sw_record)
                    sw_done   = dev_valid[0].get('done', '').lower() in ('yes', 'true', '1')
                    jira_done = jira_status.lower() in ('closed', 'done', 'resolved', 'implemented')
                    done = 'yes' if (sw_done or jira_done) else 'no'
                    hsd_trend = dev_valid[0].get('trend', '').strip()
                    comp_data[key] = {
                        'jira':   sw_record,
                        'status': jira_status or '(blank)',
                        'trend':  trend_ww or hsd_trend or '(blank)',
                        'done':   done,
                    }
                continue

            # No valid dev child — check architecture gate
            arch_children = [
                c for c in sw_parsed
                if c.get('sw_component', '').lower() in arch_kw
                and c.get('sw_task', '').lower() == 'architecture'
            ]
            if not arch_children:
                comp_data[key] = {'jira': '(blank)', 'status': '(blank)',
                                  'trend': '(blank)', 'done': '(blank)'}
                continue

            arch_exp = arch_children[0].get('sw_exposure', '').lower()
            if arch_exp == 'to_be_assigned':
                # Arch gate explicitly not yet assessed
                comp_data[key] = {'jira': '(blank)', 'status': 'Pending arch scoping',
                                  'trend': '(blank)', 'done': '(blank)'}
            elif arch_exp in ('', 'none'):
                # Arch explicitly said not involved -- show ---
                comp_data[key] = {'jira': '(blank)', 'status': '(blank)',
                                  'trend': '(blank)', 'done': '(blank)'}
            else:
                # Arch has valid exposure — check dev child state
                arch_trend = arch_children[0].get('trend', '').strip()
                any_dev = [c for c in sw_parsed
                           if c.get('sw_component', '').lower() == sw_kw
                           and c.get('sw_task', '').lower() == 'development']
                dev_pending = [c for c in any_dev
                               if c.get('sw_exposure', '').lower() == 'to_be_assigned']
                if not any_dev or dev_pending:
                    # No dev child yet, or dev child still unassigned — show pending
                    comp_data[key] = {'jira': '(blank)', 'status': 'Pending engineering scoping',
                                      'trend': arch_trend or '(blank)', 'done': '(blank)'}
                else:
                    # Dev child exists with none/empty exposure — explicitly no work needed
                    comp_data[key] = {'jira': '(blank)', 'status': '(blank)',
                                      'trend': '(blank)', 'done': '(blank)'}

        has_kmd = any(
            'i915_kmd' in c.get('sw_component', '').lower()
            and c.get('sw_task', '').lower() == 'development'
            and c.get('sw_exposure', '').lower() not in _XE5_NO_SCOPE
            for c in sw_parsed)
        # Compute involved = yes if igc_compute, pisa_finalizer, or compute
        # has dev/valcontent work (GUC/KMD-only changes don't count)
        _COMPUTE_COMPS = frozenset({'compute', 'igc_compute', 'pisa_finalizer'})
        has_compute = any(
            c.get('sw_component', '').lower() in _COMPUTE_COMPS
            and c.get('sw_task', '').lower() in ('development', 'valcontent')
            and c.get('sw_exposure', '').lower() not in _XE5_NO_SCOPE
            for c in sw_parsed)

        # Val-content for compute
        val_compute_children = [
            c for c in sw_parsed
            if c.get('sw_component', '').lower() == 'compute'
            and c.get('sw_task', '').lower() == 'valcontent'
            and c.get('sw_exposure', '').lower() not in _XE5_NO_SCOPE
        ]
        _blank_vc = {'jira': '(blank)', 'status': '(blank)', 'trend': '(blank)', 'done': '(blank)'}
        if val_compute_children:
            _vc_rec = val_compute_children[0]
            _vc_sw  = (_vc_rec.get('sw_record') or '').strip()
            if _vc_sw:
                _vc_jira_status, _vc_trend = _fetch_jira_status_and_trend(token, _vc_sw)
                _vc_sw_done   = _vc_rec.get('done', '').lower() in ('yes', 'true', '1')
                _vc_jira_done = _vc_jira_status.lower() in ('closed', 'done', 'resolved', 'implemented')
                val_compute_data = {
                    'jira':   _vc_sw,
                    'status': _vc_jira_status or '(blank)',
                    'trend':  _vc_trend       or '(blank)',
                    'done':   'yes' if (_vc_sw_done or _vc_jira_done) else 'no',
                }
            else:
                val_compute_data = {'jira': '(blank)',
                                    'status': f"No Jira ({_vc_rec.get('sw_exposure','')})",
                                    'trend': '(blank)', 'done': '(blank)'}
        else:
            val_compute_data = _blank_vc

        _no_date = 'no date yet'
        return {'id': hsd_id, 'title': title,
                'hsd_status': info.get('status', ''), 'hsd_state': info.get('state', ''),
                'comp_data': comp_data,
                'kmd_dcn': 'yes' if has_kmd else 'no',
                'compute_dcn': 'yes' if has_compute else 'no',
                'sim_rdy':         info.get('fulsim_rdy', '')       or _no_date,
                'emu_rdy':         info.get('emu_rdy', '')          or _no_date,
                'rtl_trend_rdy':   info.get('rtl_trend_rdy', '')    or _no_date,
                'turnin_trend_ww': info.get('turnin_trend_ww', '')  or _no_date,
                'study_priority':  info.get('study_priority', ''),
                'val_compute':     val_compute_data}
    except Exception as exc:
        logger.error('xe5 DCN: error on HSD %s: %s', hsd_id, exc)
        return {'id': hsd_id, 'title': f'HSD {hsd_id}', 'hsd_status': '', 'hsd_state': '',
                'comp_data': {}, 'kmd_dcn': 'no', 'compute_dcn': 'no',
                'sim_rdy': 'no date yet', 'emu_rdy': 'no date yet', 'rtl_trend_rdy': 'no date yet',
                'turnin_trend_ww': 'no date yet', 'study_priority': '',
                'val_compute': {'jira': '(blank)', 'status': '(blank)', 'trend': '(blank)', 'done': '(blank)'}}


def _fetch_xe5_all_rows(token: str) -> list:
    """Fetch all xe5 DCN parent rows (HSD info + Jira status per component).
    Returns a list of row dicts sorted by HSD id, ready for writing to the data sheet or CSV.
    """
    session    = _new_hsd_session()
    parent_ids = _execute_hsd_query(session, _XE5_DCN_QUERY_ID) or []
    logger.info('xe5 DCN: %d parent HSDs in query', len(parent_ids))
    rows: list[dict] = []
    with ThreadPoolExecutor(max_workers=6) as pool:
        futs = {pool.submit(_fetch_xe5_parent_row, token, hid): hid
                for hid in parent_ids}
        for fut in as_completed(futs):
            rows.append(fut.result())
    rows.sort(key=lambda x: x['id'])
    return rows


def _inject_web_query_connection(xlsx_bytes: bytes, csv_url: str) -> bytes:
    """Post-process xlsx bytes to embed Power Query (type=100, modern Excel 365 format)
    on the DO NOT MODIFY sheet, mirroring the JGS GT-DCNs setup.

    The DataMashup customXml part holds the M script so the query appears in the
    Data → Queries & Connections pane (Queries tab).  The queryTable wires the
    connection to the sheet cells so Data → Refresh All rewrites the cached rows.
    """
    import io as _io, zipfile as _zf, re as _re, struct as _st, base64 as _b64
    import uuid as _uuid, os as _os

    REL_NS   = 'http://schemas.openxmlformats.org/package/2006/relationships'
    OFF_REL  = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    CT_CONN  = ('application/vnd.openxmlformats-officedocument'
                '.spreadsheetml.connections+xml')
    CT_QT    = ('application/vnd.openxmlformats-officedocument'
                '.spreadsheetml.queryTable+xml')
    CT_XML   = 'application/xml'
    CT_PROPS = ('application/vnd.openxmlformats-officedocument'
                '.customXmlProperties+xml')

    # ── Pull PermissionList (seg1) and LocalPackageMetadataFile (seg2) from the
    #    JGS reference file.  These have known-good format that works in Excel.
    _JGS_PATH = _os.path.join(_os.path.dirname(__file__),
                              '..', 'JGS Compute and KMD DCNs enabling status.xlsx')
    _jgs_seg1 = None
    _jgs_seg2 = None
    try:
        with _zf.ZipFile(_JGS_PATH) as _jf:
            _raw16 = _jf.read('customXml/item15.xml')[2:].decode('utf-16-le')
            _b64jgs = _re.search(r'<DataMashup[^>]*>(.*?)</DataMashup>',
                                 _raw16, _re.DOTALL).group(1).strip()
            _raw = _b64.b64decode(_b64jgs)
            _pos = 4
            _segs = []
            while _pos < len(_raw) - 4:
                _slen = _st.unpack_from('<I', _raw, _pos)[0]; _pos += 4
                if _slen == 0 or _slen > len(_raw) - _pos: break
                _segs.append(_raw[_pos:_pos+_slen]); _pos += _slen
            if len(_segs) >= 3:
                _jgs_seg1 = _segs[1]   # PermissionList
                _jgs_seg2 = _segs[2]   # LocalPackageMetadataFile
    except Exception:
        pass   # fall back to our built-in segments if JGS file unavailable

    QUERY_NAME      = 'xe5_dcn_data'
    HSD_SW_QUERY_ID = '22023003945'   # xe5 sw_impact query
    HSD_DCN_QID     = '14028312865'   # xe5 parent DCN list query
    JIRA_TREND_CF   = 'customfield_34504'

    # ── Build DataMashup blob (same container format as JGS file) ─────────────
    # Format: [4-byte version=0][4-byte zip_len][zip_bytes][4-byte perm_len][perm_xml]
    # M script queries HSD + Jira directly with Intel Kerberos credentials,
    # exactly like JGS does — no Flask CSV proxy needed.
    def _make_mashup_b64() -> str:
        m_script = f"""\
section Section1;

shared {QUERY_NAME} = let
    Source   = Csv.Document(
                   Web.Contents("{csv_url}"),
                   [Delimiter=",", Encoding=65001, QuoteStyle=QuoteStyle.None]),
    #"Hdr"   = Table.PromoteHeaders(Source, [PromoteAllScalars=true])
in  #"Hdr";
"""
        # ── Graft our Section1.m into JGS's proven DataMashup blob ──────────────
        # JGS DataMashup is known to make queries appear in the pane.
        # We keep JGS's Package.xml, Content_Types, perm, and seg4 verbatim;
        # only Section1.m (our CSV query) and the LPMF change.
        jgs_path = Path(__file__).resolve().parent.parent / \
            'JGS Compute and KMD DCNs enabling status.xlsx'
        with _zf.ZipFile(str(jgs_path)) as _jz:
            _jgs_raw = _jz.read('customXml/item15.xml')
        _jgs_xml = _jgs_raw[2:].decode('utf-16-le')
        _jgs_b64 = _re.search(r'>([A-Za-z0-9+/=\s]+)<', _jgs_xml).group(1).replace('\n','').replace(' ','')
        jgs_blob = _b64.b64decode(_jgs_b64)

        # parse JGS blob segments
        _jzip_len  = _st.unpack_from('<I', jgs_blob, 4)[0]
        _jzip_data = jgs_blob[8:8+_jzip_len]
        _pos = 8 + _jzip_len
        _jperm_len = _st.unpack_from('<I', jgs_blob, _pos)[0]
        jgs_perm   = jgs_blob[_pos+4:_pos+4+_jperm_len]
        _pos += 4 + _jperm_len
        _jlpmf_len = _st.unpack_from('<I', jgs_blob, _pos)[0]
        _pos += 4 + _jlpmf_len
        _jseg4_len = _st.unpack_from('<I', jgs_blob, _pos)[0]
        jgs_seg4   = jgs_blob[_pos+4:_pos+4+_jseg4_len]

        # rebuild inner zip: keep Package.xml + Content_Types from JGS, swap Section1.m
        with _zf.ZipFile(_io.BytesIO(_jzip_data)) as _jinner:
            _inner_entries = {n: _jinner.read(n) for n in _jinner.namelist()}
        _inner_entries['Formulas/Section1.m'] = m_script.encode('utf-8')
        zip_buf = _io.BytesIO()
        with _zf.ZipFile(zip_buf, 'w', _zf.ZIP_DEFLATED) as _zo:
            for _n, _d in _inner_entries.items():
                _zo.writestr(_n, _d)
        zip_bytes = zip_buf.getvalue()

        perm = jgs_perm   # use JGS's PermissionList verbatim

        # LPMF: single entry for our one query
        grp = str(_uuid.uuid4())
        def _item(path, fill_enabled, fill_obj_type, is_private='l0'):
            qid = str(_uuid.uuid4())
            return (
                f'<Item><ItemLocation><ItemType>Formula</ItemType>'
                f'<ItemPath>Section1/{path}</ItemPath></ItemLocation>'
                f'<StableEntries>'
                f'<Entry Type="BufferNextRefresh" Value="l1" />'
                f'<Entry Type="FillEnabled" Value="{fill_enabled}" />'
                f'<Entry Type="FilledCompleteResultToWorksheet" Value="l0" />'
                f'<Entry Type="FillToDataModelEnabled" Value="l0" />'
                f'<Entry Type="IsPrivate" Value="{is_private}" />'
                f'<Entry Type="QueryGroupID" Value="s{grp}" />'
                f'<Entry Type="QueryID" Value="s{qid}" />'
                f'<Entry Type="ResultType" Value="sTable" />'
                f'<Entry Type="NavigationStepName" Value="sNavigation" />'
                f'<Entry Type="FillObjectType" Value="{fill_obj_type}" />'
                f'</StableEntries></Item>'
            )
        lpmf_xml = (
            '<?xml version="1.0" encoding="utf-8"?>'
            '<LocalPackageMetadataFile '
            'xmlns:xsd="http://www.w3.org/2001/XMLSchema" '
            'xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance">'
            '<Items>'
            + _item(QUERY_NAME, 'l1', 'sTable')
            + '</Items></LocalPackageMetadataFile>'
        )
        lpmf_bytes = b'\xef\xbb\xbf' + lpmf_xml.encode('utf-8')
        lpmf_seg = (
            _st.pack('<I', 0)
            + _st.pack('<I', len(lpmf_bytes))
            + lpmf_bytes
            + _st.pack('<I', 22)
            + b'PK\x05\x06' + b'\x00' * 18
        )

        blob = (_st.pack('<I', 0)
                + _st.pack('<I', len(zip_bytes)) + zip_bytes
                + _st.pack('<I', len(perm))      + perm
                + _st.pack('<I', len(lpmf_seg))  + lpmf_seg
                + _st.pack('<I', len(jgs_seg4))  + jgs_seg4)
        return _b64.b64encode(blob).decode('ascii')

    mashup_b64 = _make_mashup_b64()
    item_guid  = '{' + str(_uuid.uuid4()).upper() + '}'
    sqm_guid   = str(_uuid.uuid4())

    # customXml/item1.xml — DataMashup element stored as UTF-16-LE (matching JGS)
    item1_xml = (
        '<?xml version="1.0" encoding="UTF-16"?>'
        f'<DataMashup sqmid="{sqm_guid}" '
        'xmlns="http://schemas.microsoft.com/DataMashup">'
        f'{mashup_b64}'
        '</DataMashup>'
    )
    item1_bytes = b'\xff\xfe' + item1_xml.encode('utf-16-le')

    itemProps1_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="no"?>'
        f'<ds:datastoreItem ds:itemID="{item_guid}" '
        'xmlns:ds="http://schemas.openxmlformats.org/officeDocument/2006/customXml">'
        '<ds:schemaRefs>'
        '<ds:schemaRef ds:uri="http://schemas.microsoft.com/DataMashup"/>'
        '</ds:schemaRefs></ds:datastoreItem>'
    )
    item1_rels = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        f'<Relationships xmlns="{REL_NS}">'
        '<Relationship Id="rId1" '
        f'Type="{OFF_REL}/customXmlProps" '
        'Target="itemProps1.xml"/>'
        '</Relationships>'
    )

    # ── Read existing files ───────────────────────────────────────────────────
    with _zf.ZipFile(_io.BytesIO(xlsx_bytes), 'r') as zin:
        files = {name: zin.read(name) for name in zin.namelist()}

    # ── Resolve which sheet file is "DO NOT MODIFY" ───────────────────────────
    wb_xml = files['xl/workbook.xml'].decode('utf-8')
    m_rid  = (_re.search(r'name="DO NOT MODIFY"[^/]*?r:id="(rId\d+)"', wb_xml)
              or _re.search(r'r:id="(rId\d+)"[^/]*?name="DO NOT MODIFY"', wb_xml))
    data_rid = m_rid.group(1) if m_rid else 'rId2'
    wb_rels_xml = files['xl/_rels/workbook.xml.rels'].decode('utf-8')
    # Attribute order in openpyxl rels: Target comes before Id, so check both orderings
    m_tgt = (_re.search(rf'Id="{_re.escape(data_rid)}"[^>]*Target="([^"]+)"', wb_rels_xml)
             or _re.search(rf'Target="([^"]+)"[^>]*Id="{_re.escape(data_rid)}"', wb_rels_xml))
    if m_tgt:
        tgt = m_tgt.group(1)
        # Absolute path: "/xl/worksheets/sheet1.xml" → strip leading /
        # Relative path: "worksheets/sheet1.xml"     → prepend "xl/"
        data_sheet_path = tgt.lstrip('/') if tgt.startswith('/') else f'xl/{tgt}'
    else:
        data_sheet_path = 'xl/worksheets/sheet1.xml'
    data_rels_path  = data_sheet_path.replace(
        'xl/worksheets/', 'xl/worksheets/_rels/') + '.rels'

    # 1. connections.xml  (type=100 — modern Power Query format, same as JGS)
    uid_xr  = str(_uuid.uuid4()).upper()
    uid_x15 = str(_uuid.uuid4())
    files['xl/connections.xml'] = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<connections xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" '
        'xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006" '
        'mc:Ignorable="xr16" '
        'xmlns:xr16="http://schemas.microsoft.com/office/spreadsheetml/2017/revision16">'
        f'<connection id="1" xr16:uid="{{{uid_xr}}}" name="Query - {QUERY_NAME}" '
        f'description="Connection to the \'{QUERY_NAME}\' query in the workbook." '
        f'type="100" refreshedVersion="8" minRefreshableVersion="5">'
        f'<extLst><ext uri="{{DE250136-89BD-433C-8126-D09CA5730AF9}}" '
        f'xmlns:x15="http://schemas.microsoft.com/office/spreadsheetml/2010/11/main">'
        f'<x15:connection id="{uid_x15}"/>'
        f'</ext></extLst></connection>'
        '</connections>'
    ).encode('utf-8')

    # 2. queryTables/queryTable1.xml
    files['xl/queryTables/queryTable1.xml'] = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<queryTable xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" '
        f'name="{QUERY_NAME}" headers="1" rowNumbers="0" disableRefresh="0" '
        'backgroundRefresh="0" firstBackgroundRefresh="0" refreshOnLoad="1" '
        'growShrinkType="insertDelete" fillFormulas="0" removeDataOnSave="0" '
        'editPage="" connectionId="1">'
        '<queryTableRefresh preserveSortFilterLayout="0" fieldIdWrapped="0">'
        '<queryTableFields count="0"/>'
        '</queryTableRefresh>'
        '</queryTable>'
    ).encode('utf-8')

    # 3. customXml files
    files['customXml/item1.xml']             = item1_bytes
    files['customXml/itemProps1.xml']        = itemProps1_xml.encode('utf-8')
    files['customXml/_rels/item1.xml.rels']  = item1_rels.encode('utf-8')

    # 4. workbook.xml.rels: connections + customXml relationships
    existing_rids = {int(m) for m in _re.findall(r'Id="rId(\d+)"', wb_rels_xml)}
    next_id  = max(existing_rids, default=0) + 1
    conn_rid = f'rId{next_id}'
    cxml_rid = f'rId{next_id + 1}'
    wb_rels_xml = wb_rels_xml.replace(
        '</Relationships>',
        f'<Relationship Id="{conn_rid}" Type="{OFF_REL}/connections" Target="connections.xml"/>'
        f'<Relationship Id="{cxml_rid}" Type="{OFF_REL}/customXml" Target="../customXml/item1.xml"/>'
        '</Relationships>'
    )
    files['xl/_rels/workbook.xml.rels'] = wb_rels_xml.encode('utf-8')

    # 5. data-sheet _rels: queryTable relationship
    if data_rels_path in files:
        ds_rels = files[data_rels_path].decode('utf-8')
        ds_rids = {int(m) for m in _re.findall(r'Id="rId(\d+)"', ds_rels)}
        qt_rid  = f'rId{max(ds_rids, default=0) + 1}'
        ds_rels = ds_rels.replace(
            '</Relationships>',
            f'<Relationship Id="{qt_rid}" Type="{OFF_REL}/queryTable" '
            'Target="../queryTables/queryTable1.xml"/></Relationships>'
        )
    else:
        qt_rid  = 'rId1'
        ds_rels = (
            f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            f'<Relationships xmlns="{REL_NS}">'
            f'<Relationship Id="{qt_rid}" Type="{OFF_REL}/queryTable" '
            'Target="../queryTables/queryTable1.xml"/>'
            '</Relationships>'
        )
    files[data_rels_path] = ds_rels.encode('utf-8')

    # 6. data-sheet XML: ensure xmlns:r + queryTableParts before </worksheet>
    ds_xml = files[data_sheet_path].decode('utf-8')
    if 'xmlns:r=' not in ds_xml:
        ds_xml = ds_xml.replace(
            'xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main"',
            'xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" '
            f'xmlns:r="{OFF_REL}"',
            1,
        )
    ds_xml = ds_xml.replace(
        '</worksheet>',
        f'<queryTableParts count="1"><queryTablePart r:id="{qt_rid}"/></queryTableParts>'
        '</worksheet>'
    )
    files[data_sheet_path] = ds_xml.encode('utf-8')

    # 7. [Content_Types].xml
    ct_xml = files['[Content_Types].xml'].decode('utf-8')
    ct_xml = ct_xml.replace(
        '</Types>',
        f'<Override PartName="/xl/connections.xml" ContentType="{CT_CONN}"/>'
        f'<Override PartName="/xl/queryTables/queryTable1.xml" ContentType="{CT_QT}"/>'
        f'<Override PartName="/customXml/item1.xml" ContentType="{CT_XML}"/>'
        f'<Override PartName="/customXml/itemProps1.xml" ContentType="{CT_PROPS}"/>'
        '</Types>'
    )
    files['[Content_Types].xml'] = ct_xml.encode('utf-8')

    # Reassemble zip
    out_buf = _io.BytesIO()
    with _zf.ZipFile(out_buf, 'w', _zf.ZIP_DEFLATED) as zout:
        for name, data in files.items():
            zout.writestr(name, data)
    return out_buf.getvalue()


def _build_from_jgs_template() -> bytes:
    """Use the JGS file as the base, swap the HSD/Jira query IDs in the Setup
    sheet, and patch the M script so it works with xe5 DCN records:
      - HSD_feature: filter changed from subject="feature" to subject="dcn";
        rtl_actual_rdy_for_drv_ww added to the field expansion.
      - Shown_DCNs: replaced with a direct reference to HSD_feature so ALL
        DCNs from the query are shown without a manual maintenance list.
    The DataMashup keeps JGS's Data-Model / pivot-table structure, so refresh
    goes directly to HSD/Jira via Kerberos — no Flask server needed."""
    import io as _io, zipfile as _zf, base64 as _b64, struct as _st, re as _re

    jgs_path = Path(__file__).resolve().parent.parent / \
        'JGS Compute and KMD DCNs enabling status.xlsx'
    with open(jgs_path, 'rb') as fh:
        jgs_bytes = fh.read()

    entries = {}
    with _zf.ZipFile(_io.BytesIO(jgs_bytes), 'r') as zin:
        for name in zin.namelist():
            entries[name] = zin.read(name)

    # ── 1. Setup sheet: swap HSD query IDs ───────────────────────────────────
    setup_xml = entries['xl/worksheets/sheet18.xml'].decode('utf-8')
    setup_xml = re.sub(
        r'(<c r="C10"[^>]*>)<v>[^<]*</v>',
        r'\g<1><v>14028312865</v>',
        setup_xml,
    )
    setup_xml = re.sub(
        r'(<c r="C11"[^>]*>)<v>[^<]*</v>',
        r'\g<1><v>22023003945</v>',
        setup_xml,
    )
    entries['xl/worksheets/sheet18.xml'] = setup_xml.encode('utf-8')

    # ── 2. DataMashup: patch M script ────────────────────────────────────────
    try:
        raw16    = entries['customXml/item15.xml'][2:].decode('utf-16-le')
        b64      = _re.search(r'<DataMashup[^>]*>(.*?)</DataMashup>',
                              raw16, _re.DOTALL).group(1).strip()
        raw      = _b64.b64decode(b64)
        pos      = 4
        zip_len  = _st.unpack_from('<I', raw, pos)[0]; pos += 4
        zip_data = raw[pos:pos + zip_len]; pos += zip_len
        rest     = raw[pos:]   # perm + lpmf + seg4 stay unchanged

        with _zf.ZipFile(_io.BytesIO(zip_data)) as zin:
            inner = {n: zin.read(n) for n in zin.namelist()}

        m = inner['Formulas/Section1.m'].decode('utf-8')

        # a) HSD_feature: change subject filter to "dcn"
        m = m.replace('each ([subject] = "feature")',
                      'each ([subject] = "dcn")')

        # b) HSD_feature: add rtl_actual_rdy_for_drv_ww to ExpandRecordColumn
        #    Source-field list ends with: ..., "rtl_trend_rdy_for_drv_ww"}, {
        #    Column-name list ends with:  ..., "rtl_trend_rdy_for_drv_ww"})
        m = m.replace(
            '"state", "subject", "rtl_trend_rdy_for_drv_ww"}, {"parent_id"',
            '"state", "subject", "rtl_trend_rdy_for_drv_ww", "rtl_actual_rdy_for_drv_ww"}, {"parent_id"',
        )
        m = m.replace(
            '"Status", "subject", "rtl_trend_rdy_for_drv_ww"})',
            '"Status", "subject", "rtl_trend_rdy_for_drv_ww", "Emu readiness"})',
        )

        # c) Shown_DCNs: drop the Manual_DCN_list join and show all HSD_feature
        #    rows directly so no manual worksheet maintenance is needed.
        m = _re.sub(
            r'shared Shown_DCNs = let\b.*?in\s+#"Filtered Rows";',
            ('shared Shown_DCNs = let\n'
             '    Source = HSD_feature,\n'
             '    #"Replaced Value" = Table.ReplaceValue(Source,null,'
             '"<DCN removed>",Replacer.ReplaceValue,{"Title"}),\n'
             '    #"Added Index" = Table.AddIndexColumn(#"Replaced Value",'
             '"ord_num",1,1,Int64.Type)\n'
             'in\n'
             '    #"Added Index";'),
            m,
            flags=_re.DOTALL,
        )

        # d) Remove Jira_sw_DCN — JGS-specific Jira query (JQL uses JGS labels),
        #    not needed for xe5 DCN tracking.
        for _nl in ('\r\nshared Jira_sw_DCN', '\nshared Jira_sw_DCN'):
            _ji = m.find(_nl)
            if _ji != -1:
                m = m[:_ji]
                break

        inner['Formulas/Section1.m'] = m.encode('utf-8')

        new_zip_buf = _io.BytesIO()
        with _zf.ZipFile(new_zip_buf, 'w', _zf.ZIP_DEFLATED) as zout:
            for n, d in inner.items():
                zout.writestr(n, d)
        new_zip = new_zip_buf.getvalue()

        new_raw  = raw[:4] + _st.pack('<I', len(new_zip)) + new_zip + rest
        new_b64  = _b64.b64encode(new_raw).decode('ascii')
        new_raw16 = _re.sub(
            r'(<DataMashup[^>]*>)[^<]*(</DataMashup>)',
            lambda mo: mo.group(1) + new_b64 + mo.group(2),
            raw16,
            flags=_re.DOTALL,
        )
        entries['customXml/item15.xml'] = b'\xff\xfe' + new_raw16.encode('utf-16-le')
        logger.info('xe5 JGS template: M script patched OK')
    except Exception as exc:
        logger.warning('xe5 JGS template: M script patch failed (%s) — '
                       'serving JGS M script unchanged', exc)

    # ── 3. Write output ───────────────────────────────────────────────────────
    out_buf = _io.BytesIO()
    with _zf.ZipFile(out_buf, 'w', _zf.ZIP_DEFLATED) as zout:
        for name, data in entries.items():
            zout.writestr(name, data)
    return out_buf.getvalue()


def _build_xe5_dcn_excel(token: str) -> 'io.BytesIO':
    """Fetch xe5 DCN data and build the Excel matching the JGS GT-DCNs structure."""
    import io as _io
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment
    from openpyxl.utils import get_column_letter
    from openpyxl.formatting.rule import FormulaRule

    rows = _fetch_xe5_all_rows(token)

    wb = Workbook()

    # ── Sheet 1: DO NOT MODIFY (raw data) ─────────────────────────────────────
    data_ws = wb.active
    data_ws.title = 'DO NOT MODIFY'

    DATA_HEADERS = [
        'ord_num', 'id', 'Title', 'hsd_status',
        'dev_compute_jira', 'dev_compute_status', 'dev_compute_trend', 'dev_compute_done',
        'dev_igc_jira',     'dev_igc_status',     'dev_igc_trend',     'dev_igc_done',
        'dev_kmd_jira',     'dev_kmd_status',     'dev_kmd_trend',     'dev_kmd_done',
        'dev_pisa_jira',    'dev_pisa_status',    'dev_pisa_trend',    'dev_pisa_done',
        'dev_guc_jira',     'dev_guc_status',     'dev_guc_trend',     'dev_guc_done',
        'KMD DCN', 'Compute DCN', 'hsd_state',
        'sim_rdy', 'emu_rdy', 'rtl_trend_rdy',
        'val_compute_jira', 'val_compute_status', 'val_compute_trend', 'val_compute_done',
    ]
    _hdr_fill  = PatternFill('solid', fgColor='2C4770')
    _hdr_font  = Font(color='FFFFFF', bold=True, size=9)
    _ctr_align = Alignment(horizontal='center', vertical='center', wrap_text=True)

    for c, h in enumerate(DATA_HEADERS, 1):
        cell = data_ws.cell(1, c, h)
        cell.fill = _hdr_fill; cell.font = _hdr_font; cell.alignment = _ctr_align

    DATA_COL_IDX = {
        'compute': (5, 6, 7, 8),
        'igc':     (9, 10, 11, 12),
        'kmd':     (13, 14, 15, 16),
        'pisa':    (17, 18, 19, 20),
        'guc':     (21, 22, 23, 24),
    }
    _blank = {'jira': '(blank)', 'status': '(blank)', 'trend': '(blank)', 'done': '(blank)'}

    for i, row in enumerate(rows, 1):
        dr = i + 1   # data row (row 1 = headers, row 2+ = data)
        cd = row['comp_data']
        data_ws.cell(dr, 1, i)
        data_ws.cell(dr, 2, row['id'])
        data_ws.cell(dr, 3, row['title'])
        data_ws.cell(dr, 4, row.get('hsd_status', ''))
        for key, (jc, sc, tc, dc) in DATA_COL_IDX.items():
            d = cd.get(key, _blank)
            data_ws.cell(dr, jc, d['jira'])
            data_ws.cell(dr, sc, d['status'])
            data_ws.cell(dr, tc, d['trend'])
            data_ws.cell(dr, dc, d['done'])
        data_ws.cell(dr, 25, row['kmd_dcn'])
        data_ws.cell(dr, 26, row['compute_dcn'])
        data_ws.cell(dr, 27, row.get('hsd_state', ''))    # AA
        data_ws.cell(dr, 28, row.get('sim_rdy', 'no date yet'))   # AB
        data_ws.cell(dr, 29, row.get('emu_rdy', 'no date yet'))   # AC
        data_ws.cell(dr, 30, row.get('rtl_trend_rdy', 'no date yet'))  # AD
        _vc = row.get('val_compute', _blank)
        data_ws.cell(dr, 31, _vc['jira'])    # AE
        data_ws.cell(dr, 32, _vc['status'])  # AF
        data_ws.cell(dr, 33, _vc['trend'])   # AG
        data_ws.cell(dr, 34, _vc['done'])    # AH

    data_ws.column_dimensions['A'].width = 8
    data_ws.column_dimensions['B'].width = 14
    data_ws.column_dimensions['C'].width = 52
    data_ws.column_dimensions['D'].width = 14
    for c in range(5, 35):
        data_ws.column_dimensions[get_column_letter(c)].width = 15
    data_ws.freeze_panes = 'C2'

    # ── Sheet 2: xe5 GT DCNs (main view with formulas) ────────────────────────
    # Main sheet column layout:
    #  A:ID  B:Title  C:Enabling(manual)  D:Compute involved  E:KMD involved
    #  F:Compute  G:IGC Compute  H:XeKMD  I:PISA  J:GUC  ← dev readiness
    #  K:Compute  L:IGC  M:XeKMD  N:PISA  O:GUC  P:Impl.Trend ← trends
    #  Q:Comments
    main_ws = wb.create_sheet('xe5 GT DCNs', 0)

    DN = "'DO NOT MODIFY'"   # data sheet reference in formulas

    # Row 1: title banner
    main_ws.cell(1, 1, f'xe5 DCN Enabling Status — {len(rows)} DCNs — '
                        f'Generated {datetime.utcnow().strftime("%Y-%m-%d")}')
    main_ws.merge_cells('A1:X1')
    main_ws['A1'].fill  = PatternFill('solid', fgColor='1F3864')
    main_ws['A1'].font  = Font(color='FFFFFF', bold=True, size=12)
    main_ws['A1'].alignment = Alignment(horizontal='center', vertical='center')
    main_ws.row_dimensions[1].height = 26

    # Row 2: group headers
    for rng, label, fg in [
        ('A2:G2', 'xe5 DCN Info',          'C9D7F0'),
        ('H2:L2', 'Development Readiness', 'D6E4BC'),
        ('M2:R2', 'Implementation Trends', 'FCE4D6'),
        ('S2:S2', 'Comments',              'F2F2F2'),
        ('T2:V2', 'HSD Readiness',         'EAF1FB'),
        ('W2:X2', 'Test Development',       'F4E6FA'),
    ]:
        main_ws.merge_cells(rng)
        c = main_ws[rng.split(':')[0]]
        c.value = label
        c.fill  = PatternFill('solid', fgColor=fg)
        c.font  = Font(bold=True, size=10)
        c.alignment = Alignment(horizontal='center', vertical='center')
    main_ws.row_dimensions[2].height = 20

    # Row 3: column headers
    COL_HEADERS = [
        'ID', 'Title', 'Enabling', 'Compute involved', 'KMD involved',
        'HSD Status', 'HSD State',
        'Compute', 'IGC Compute', 'XeKMD', 'PISA', 'GUC',
        'Compute', 'IGC Compute', 'XeKMD', 'PISA', 'GUC', 'Impl. Trend',
        'Comments',
        'Sim Readiness', 'Emu Readiness', 'RTL Trend WW',
        'Test Dev Status', 'Compute Test Dev',
    ]
    _col_hdr_fill = PatternFill('solid', fgColor='4472C4')
    for c, h in enumerate(COL_HEADERS, 1):
        cell = main_ws.cell(3, c, h)
        cell.fill = _col_hdr_fill
        cell.font = Font(color='FFFFFF', bold=True, size=9)
        cell.alignment = _ctr_align
    main_ws.row_dimensions[3].height = 30

    # Helper functions for formulas
    def _dev_formula(jc, sc, dc, mr, dr):
        """Dev readiness: scope-pending text, Done/Status hyperlink, or In Analysis."""
        # When jira=(blank): show scope-status text if set, else show ―
        return (
            f"=IF({DN}!{jc}{dr}<>\"\","
            f"IF({DN}!{jc}{dr}=\"(blank)\","
            f"IF({DN}!{sc}{dr}=\"(blank)\",\"―\",{DN}!{sc}{dr}),"
            f"IF({DN}!{dc}{dr}=\"yes\","
            f"HYPERLINK(\"{_JIRA_BROWSE_URL}\"&{DN}!{jc}{dr},\"Done\"),"
            f"HYPERLINK(\"{_JIRA_BROWSE_URL}\"&{DN}!{jc}{dr},{DN}!{sc}{dr}))),"
            f"IF(A{mr}<>\"\",\"In Analysis\",\"\"))"
        )

    def _trend_formula(tc, mr, dr):
        return (
            f"=IF({DN}!{tc}{dr}<>\"\","
            f"IF({DN}!{tc}{dr}=\"(blank)\",\"―\",{DN}!{tc}{dr}),"
            f"IF(A{mr}<>\"\",\"In Analysis\",\"\"))"
        )

    def _impl_trend_formula(mr):
        """Max WW across trend cols M-Q; if non-open state and all involved components done → 'Done'."""
        dr = mr - 2  # data row
        # (jira_col, done_col) pairs in DATA sheet for each component
        comp_jc_dc = [('E', 'H'), ('I', 'L'), ('M', 'P'), ('Q', 'T'), ('U', 'X')]
        any_inv   = 'OR(' + ','.join(f"{DN}!{jc}{dr}<>\"(blank)\"" for jc, dc in comp_jc_dc) + ')'
        all_done  = 'AND(' + ','.join(f'OR({DN}!{jc}{dr}="(blank)",{DN}!{dc}{dr}="yes")' for jc, dc in comp_jc_dc) + ')'
        nonopen   = f'AND({DN}!AA{dr}<>"",LOWER({DN}!AA{dr})<>"open")'
        done_cond = f'AND({nonopen},{any_inv},{all_done})'

        def parse_ww(col):
            v = f"{col}{mr}"
            return (
                f"IF(ISNUMBER(VALUE(RIGHT({v},2))),"
                f"VALUE(RIGHT(_xlfn.TEXTBEFORE({v},\"ww\"),2)&_xlfn.TEXTAFTER({v},\"ww\")),0)"
            )
        parts   = ','.join(parse_ww(c) for c in ['M', 'N', 'O', 'P', 'Q'])
        max_exp = f'MAX({parts})'
        return (
            f"=IF({done_cond},\"Done\","
            f"IF({max_exp}=0,"
            f"IF(A{mr}<>\"\",IF(B{mr}<>\"<DCN removed>\",\"In Analysis\",\"―\"),\"\"),"
            f"LEFT({max_exp},2)&\"ww\"&RIGHT({max_exp},2)))"
        )

    # Data rows
    _id_font = Font(color='0563C1', underline='single')
    for i in range(len(rows)):
        mr = i + 4   # main row (4, 5, 6 …)
        dr = i + 2   # data row  (2, 3, 4 …)

        # A: ID hyperlink
        cell = main_ws.cell(mr, 1,
            f"=IF({DN}!B{dr}<>\"\","
            f"HYPERLINK(\"{_HSD_ARTICLE_URL}\"&{DN}!B{dr},{DN}!B{dr}),\"\")")
        cell.font = _id_font; cell.alignment = _ctr_align

        # B: Title
        main_ws.cell(mr, 2, f"=IF({DN}!C{dr}<>\"\",{DN}!C{dr},\"\")")
        main_ws.cell(mr, 2).alignment = Alignment(vertical='center', wrap_text=True)

        # C: Enabling — intentionally left empty (manual)

        # D: Compute involved
        main_ws.cell(mr, 4, f"=IF({DN}!Z{dr}<>\"\",{DN}!Z{dr},\"\")")
        main_ws.cell(mr, 4).alignment = _ctr_align

        # E: KMD involved
        main_ws.cell(mr, 5, f"=IF({DN}!Y{dr}<>\"\",{DN}!Y{dr},\"\")")
        main_ws.cell(mr, 5).alignment = _ctr_align

        # F: HSD Status (DATA col D)  G: HSD State (DATA col AA)
        main_ws.cell(mr, 6, f"=IF({DN}!D{dr}<>\"\",{DN}!D{dr},\"\")").alignment = _ctr_align
        main_ws.cell(mr, 7, f"=IF({DN}!AA{dr}<>\"\",{DN}!AA{dr},\"\")").alignment = _ctr_align

        # H-L: Development readiness (Compute, IGC, KMD, PISA, GUC)
        dev_srcs = [
            (_XE5_DATA_COL['compute'][0], _XE5_DATA_COL['compute'][1], _XE5_DATA_COL['compute'][3]),
            (_XE5_DATA_COL['igc'][0],     _XE5_DATA_COL['igc'][1],     _XE5_DATA_COL['igc'][3]),
            (_XE5_DATA_COL['kmd'][0],     _XE5_DATA_COL['kmd'][1],     _XE5_DATA_COL['kmd'][3]),
            (_XE5_DATA_COL['pisa'][0],    _XE5_DATA_COL['pisa'][1],    _XE5_DATA_COL['pisa'][3]),
            (_XE5_DATA_COL['guc'][0],     _XE5_DATA_COL['guc'][1],     _XE5_DATA_COL['guc'][3]),
        ]
        for col_offset, (jc, sc, dc) in enumerate(dev_srcs):
            c = main_ws.cell(mr, 8 + col_offset, _dev_formula(jc, sc, dc, mr, dr))
            c.alignment = _ctr_align

        # M-Q: Trend values
        trend_cols = [
            _XE5_DATA_COL['compute'][2], _XE5_DATA_COL['igc'][2],
            _XE5_DATA_COL['kmd'][2],     _XE5_DATA_COL['pisa'][2],
            _XE5_DATA_COL['guc'][2],
        ]
        for col_offset, tc in enumerate(trend_cols):
            c = main_ws.cell(mr, 13 + col_offset, _trend_formula(tc, mr, dr))
            c.alignment = _ctr_align

        # R: Impl. Trend
        main_ws.cell(mr, 18, _impl_trend_formula(mr)).alignment = _ctr_align

        # S: Comments — intentionally empty (manual)

        # T: Sim Readiness (fulsim_rdy_for_drv_ww)
        main_ws.cell(mr, 20, f"=IF({DN}!AB{dr}<>\"\",{DN}!AB{dr},\"\")").alignment = _ctr_align
        # U: Emu Readiness (rtl_actual_rdy_for_drv_ww)
        main_ws.cell(mr, 21, f"=IF({DN}!AC{dr}<>\"\",{DN}!AC{dr},\"\")").alignment = _ctr_align
        # V: RTL Trend WW (rtl_trend_rdy_for_drv_ww)
        main_ws.cell(mr, 22, f"=IF({DN}!AD{dr}<>\"\",{DN}!AD{dr},\"\")").alignment = _ctr_align
        # W: Test Dev Status (val compute jira/status)
        main_ws.cell(mr, 23, _dev_formula('AE', 'AF', 'AH', mr, dr)).alignment = _ctr_align
        # X: Compute Test Dev trend
        main_ws.cell(mr, 24, _trend_formula('AG', mr, dr)).alignment = _ctr_align
        main_ws.row_dimensions[mr].height = 18

    # Column widths (19 cols)
    for c, w in enumerate([12, 58, 10, 14, 12, 14, 12, 12, 14, 10, 14, 10, 12, 14, 10, 14, 10, 12, 35, 16, 16, 16, 16, 16], 1):
        main_ws.column_dimensions[get_column_letter(c)].width = w

    # Freeze row 1-3 + col A-B
    main_ws.freeze_panes = 'C4'

    # Auto-filter on header row
    main_ws.auto_filter.ref = f'A3:{get_column_letter(24)}3'

    # ── Conditional formatting ────────────────────────────────────────────────
    n_rows   = len(rows)
    _dev_rng = f'H4:L{3 + n_rows}'
    # Done = green cell background  (DXF fills use bgColor for rendered background)
    main_ws.conditional_formatting.add(_dev_rng, FormulaRule(
        formula=['H4="Done"'],
        fill=PatternFill(bgColor='C6EFCE'),
        font=Font(color='276221', bold=True)))
    # Pending arch scoping = amber
    main_ws.conditional_formatting.add(_dev_rng, FormulaRule(
        formula=['ISNUMBER(SEARCH("arch scoping",H4))'],
        fill=PatternFill(bgColor='FFEB9C'),
        font=Font(color='9C6500', bold=True)))
    # Pending engineering scoping = light orange
    main_ws.conditional_formatting.add(_dev_rng, FormulaRule(
        formula=['ISNUMBER(SEARCH("engineering scoping",H4))'],
        fill=PatternFill(bgColor='FCE4D6'),
        font=Font(color='843C0C', bold=True)))

    buf = _io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    # Inject legacy web-query connection so Refresh All re-fetches the DO NOT MODIFY data
    _XE5_CSV_URL = (
        'http://10.88.27.190:8889/xe5-dcn-status/data.csv'
    )
    enhanced = _inject_web_query_connection(buf.read(), _XE5_CSV_URL)
    return _io.BytesIO(enhanced)


@bp.route('/xe5-dcn-status')
def xe5_dcn_status():
    return render_template('xe5_dcn_status.html')


@bp.route('/xe5-dcn-plan')
def xe5_dcn_plan():
    return render_template('xe5_dcn_plan.html')


@bp.route('/xe5-dcn-status/data.json')
def xe5_dcn_data_json():
    """Return all xe5 DCN rows as JSON for the DataTables web page.
    Results are cached for _XE5_JSON_TTL seconds; pass ?refresh=1 to force a live fetch.
    """
    import time as _time
    force = request.args.get('refresh', '0') == '1'
    now   = _time.time()

    # Serve from cache if still fresh and no force-refresh requested
    with _xe5_cache_lock():
        cached_rows = _XE5_JSON_CACHE['rows']
        cache_age   = now - _XE5_JSON_CACHE['ts']
        if cached_rows is not None and not force and cache_age < _XE5_JSON_TTL:
            age_s = int(cache_age)
            logger.info('xe5 JSON: serving from cache (age %ds)', age_s)
            return jsonify({'rows': cached_rows, 'cached': True, 'cache_age_s': age_s})

    token, err = _token_response()
    if err:
        return err
    try:
        raw_rows = _fetch_xe5_all_rows(token)
    except Exception as exc:
        logger.error('xe5 JSON: fetch failed: %s', exc)
        return jsonify({'error': str(exc)}), 500

    _blank = {'jira': '(blank)', 'status': '(blank)', 'trend': '(blank)', 'done': '(blank)'}

    def _fmt_comp(row, key):
        d = (row.get('comp_data') or {}).get(key, _blank)
        return {
            'jira':   d.get('jira', '(blank)'),
            'status': d.get('status', '(blank)'),
            'trend':  d.get('trend', '(blank)'),
            'done':   d.get('done', '(blank)'),
        }

    rows_out = []
    for i, row in enumerate(raw_rows, 1):
        vc = row.get('val_compute') or _blank
        rows_out.append({
            'ord':              i,
            'id':               str(row.get('id', '')),
            'title':            row.get('title', ''),
            'sim_rdy':          row.get('sim_rdy', 'no date yet') or 'no date yet',
            'emu_rdy':          row.get('emu_rdy', 'no date yet') or 'no date yet',
            'rtl_trend':        row.get('rtl_trend_rdy', 'no date yet') or 'no date yet',
            'turnin_trend_ww':  row.get('turnin_trend_ww', 'no date yet') or 'no date yet',
            'study_priority':   row.get('study_priority', ''),
            'hsd_status':       row.get('hsd_status', ''),
            'hsd_state':        row.get('hsd_state', ''),
            'compute_involved': row.get('compute_dcn', 'no'),
            'kmd_involved':     row.get('kmd_dcn', 'no'),
            'compute':          _fmt_comp(row, 'compute'),
            'igc':              _fmt_comp(row, 'igc'),
            'kmd':              _fmt_comp(row, 'kmd'),
            'pisa':             _fmt_comp(row, 'pisa'),
            'guc':              _fmt_comp(row, 'guc'),
            'val_status':       vc.get('status', '(blank)'),
            'val_trend':        vc.get('trend', '(blank)'),
            'val_done':         vc.get('done', '(blank)'),
            'val_jira':         vc.get('jira', '(blank)'),
        })

    # Store in cache only if we got real data (don't cache empty results from transient failures)
    import time as _time2
    if rows_out:
        with _xe5_cache_lock():
            _XE5_JSON_CACHE['rows'] = rows_out
            _XE5_JSON_CACHE['ts']   = _time2.time()
    elif _XE5_JSON_CACHE['rows']:
        # Fetch returned 0 rows — likely a transient HSD failure; serve stale cache with a warning
        logger.warning('xe5 JSON: live fetch returned 0 rows — serving stale cache instead')
        stale_age = int(_time2.time() - _XE5_JSON_CACHE['ts'])
        return jsonify({'rows': _XE5_JSON_CACHE['rows'],
                        'cached': True, 'cache_age_s': stale_age, 'stale': True})

    return jsonify({'rows': rows_out, 'cached': False, 'cache_age_s': 0})


@bp.route('/xe5-dcn-status/data.csv')
def xe5_dcn_data_csv():
    """Serve the raw DCN data as CSV — consumed by the Excel web-query connection.
    Hitting Data → Refresh All in Excel re-fetches this URL and updates the DO NOT MODIFY
    sheet, which cascades through all the main-sheet HYPERLINK/IF formulas automatically.
    """
    import csv as _csv
    token, err = _token_response()
    if err:
        return err
    try:
        rows = _fetch_xe5_all_rows(token)
    except Exception as exc:
        logger.error('xe5 CSV: data fetch failed: %s', exc)
        return f'# error fetching data: {exc}\n', 500, {'Content-Type': 'text/plain'}

    import io as _csv_io
    buf = _csv_io.StringIO()
    w   = _csv.writer(buf)
    w.writerow([
        'ord_num', 'id', 'Title', 'hsd_status',
        'dev_compute_jira', 'dev_compute_status', 'dev_compute_trend', 'dev_compute_done',
        'dev_igc_jira',     'dev_igc_status',     'dev_igc_trend',     'dev_igc_done',
        'dev_kmd_jira',     'dev_kmd_status',     'dev_kmd_trend',     'dev_kmd_done',
        'dev_pisa_jira',    'dev_pisa_status',    'dev_pisa_trend',    'dev_pisa_done',
        'dev_guc_jira',     'dev_guc_status',     'dev_guc_trend',     'dev_guc_done',
        'KMD DCN', 'Compute DCN', 'hsd_state',
        'sim_rdy', 'emu_rdy', 'rtl_trend_rdy',
        'val_compute_jira', 'val_compute_status', 'val_compute_trend', 'val_compute_done',
    ])
    _blank_d = {'jira': '(blank)', 'status': '(blank)', 'trend': '(blank)', 'done': '(blank)'}
    for i, row in enumerate(rows, 1):
        cd  = row.get('comp_data', {})
        rec = [i, row['id'], row.get('title', ''), row.get('hsd_status', '')]
        for key in ('compute', 'igc', 'kmd', 'pisa', 'guc'):
            d = cd.get(key, _blank_d)
            rec.extend([d['jira'], d['status'], d['trend'], d['done']])
        _vc = row.get('val_compute', _blank_d)
        rec.extend([row.get('kmd_dcn', 'no'), row.get('compute_dcn', 'no'),
                    row.get('hsd_state', ''),
                    row.get('sim_rdy', 'no date yet'),
                    row.get('emu_rdy', 'no date yet'),
                    row.get('rtl_trend_rdy', 'no date yet'),
                    _vc['jira'], _vc['status'], _vc['trend'], _vc['done']])
        w.writerow(rec)

    return buf.getvalue(), 200, {
        'Content-Type':        'text/csv; charset=utf-8',
        'Content-Disposition': 'inline; filename="xe5_dcn_data.csv"',
    }


@bp.route('/xe5-dcn-status/download-jgs-raw')
def xe5_dcn_download_jgs_raw():
    """Serve the JGS file completely unmodified — for debugging only."""
    jgs_path = Path(__file__).resolve().parent.parent / 'JGS Compute and KMD DCNs enabling status.xlsx'
    return send_file(str(jgs_path), as_attachment=True, download_name='jgs_raw_test.xlsx',
                     mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')


@bp.route('/xe5-dcn-status/download')
def xe5_dcn_download():
    """Generate and return the xe5 DCN enabling status Excel.
    Serves the JGS template with xe5 query IDs and a patched M script so
    Data → Refresh All queries HSD/Jira directly via Kerberos — no Flask
    server involvement required after download."""
    try:
        xlsx = _build_from_jgs_template()
    except Exception as exc:
        logger.error('xe5 DCN: Excel generation failed: %s', exc)
        return jsonify({'error': str(exc)}), 500
    import io
    fname = f'xe5_DCN_enabling_status_{datetime.utcnow().strftime("%Y%m%d")}.xlsx'
    return send_file(
        io.BytesIO(xlsx),
        mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        as_attachment=True,
        download_name=fname,
    )


@bp.route('/hsd2jira-flow')
def hsd2jira_flow():
    return render_template('hsd2jira_flow.html')


@bp.route('/jira-tracking-design')
def jira_tracking_design():
    return render_template('jira_tracking_design.html')


@bp.route('/jira-tracking-rules')
def jira_tracking_rules():
    return render_template('jira_tracking_rules.html')


@bp.route('/hsd2jira/create', methods=['POST'])
def hsd2jira_create():
    body          = request.get_json(force=True)
    platform      = (body.get('platform') or '').strip()
    feature_title = (body.get('feature_title') or '').strip()
    _comp_raw     = body.get('components') or body.get('component') or ''
    component     = _comp_raw if isinstance(_comp_raw, list) else _comp_raw.strip()
    assignee      = (body.get('assignee')      or '').strip()
    project_tag   = (body.get('project_tag')   or '').strip()

    _comp_list    = _norm_components(component)
    if not platform or not feature_title or not _comp_list or not assignee:
        return jsonify({'error': 'platform, feature_title, component, and assignee are required'}), 400

    token, err = _token_response()
    if err:
        return err

    post_issue, create_link = _make_jira_creator(token)
    result = _build_jira_issues(post_issue, create_link,
                                feature_title, platform, component, assignee,
                                project_tag=project_tag)
    if result.get('epic_key') is None:
        return jsonify({'error': 'Epic creation failed',
                        'created': result['created']}), 500
    return jsonify({'created': result['created'],
                    'errors':  result['errors'],
                    'epic_key': result['epic_key']})


# ── HSD-link routes ───────────────────────────────────────────────────────────

@bp.route('/hsd2jira/debug-ar/<hsd_id>')
def hsd2jira_debug_ar(hsd_id):
    """
    Debug endpoint — dumps raw ESService response for AR children of hsd_id.
    GET /hsd2jira/debug-ar/16029853828
    Shows every key present in the raw records so we can find correct field names.
    """
    session = _new_hsd_session()

    # Try several commands to find children
    out = {'hsd_id': hsd_id, 'attempts': {}}

    # Step 1: raw related records (no filter) to discover what's linked
    all_links = _esservice_request(session, 'get_related_records', {'id': str(hsd_id)}) or []
    ar_links  = [
        r for r in all_links
        if isinstance(r, dict)
        and str(r.get('subject', '')).lower() == 'ar'
        and str(r.get('tenant',  '')).lower() in ('dg_soc', 'dg.soc')
    ]
    out['raw_link_count']    = len(all_links)
    out['ar_link_ids']       = [str(r['id']) for r in ar_links if r.get('id')]
    out['all_link_subjects'] = sorted({str(r.get('subject','')) for r in all_links if isinstance(r, dict)})
    out['all_link_tenants']  = sorted({str(r.get('tenant', '')) for r in all_links if isinstance(r, dict)})

    # Step 2: fetch full AR articles and show ALL their keys
    full_ar_records = []
    for r in ar_links:   # fetch all AR records
        ar_id  = str(r.get('id', ''))
        result = _esservice_request(session, 'get_record_by_id', {'id': ar_id})
        if result and isinstance(result, list):
            full_ar_records.extend(result)
        elif isinstance(result, dict):
            full_ar_records.append(result)
    out['full_ar_records']  = full_ar_records
    out['full_ar_all_keys'] = sorted({k for rec in full_ar_records if isinstance(rec, dict) for k in rec.keys()})

    for cmd, args in [
        # legacy attempts kept for comparison
        ('get_related_records', {'id': str(hsd_id), 'tenant': 'ip_hw_graphics', 'subject': 'feature'}),
    ]:
        label   = f'{cmd}:{args.get("tenant")}/{args.get("subject")}'
        records = _esservice_request(session, cmd, args)
        if records:
            # Return first 3 records with ALL their keys so we can inspect field names
            sample = records[:3]
            out['attempts'][label] = {
                'count':   len(records),
                'sample':  sample,
                'all_keys': sorted({k for r in records if isinstance(r, dict) for k in r.keys()}),
            }
        else:
            out['attempts'][label] = {'count': 0, 'sample': [], 'all_keys': []}

    return jsonify(out)



@bp.route('/hsd2jira/create-from-hsd-link', methods=['POST'])
def hsd2jira_create_from_hsd_link():
    """
    Full flow:
      1. Extract HSD ID from the supplied URL.
      2. Fetch the HSD article title.
      3. Find the AR child with team=i915_kmd & task=development.
      4. Create Epic + 3 Stories + 1 Task in Jira (title = HSD title).
      5. Write the Code Complete story key back to dg.soc.ar.jira_key of the matched child.
    """
    body      = request.get_json(force=True)
    hsd_url   = (body.get('hsd_url')   or '').strip()
    component = (body.get('component') or _COMPONENT).strip()
    assignee  = (body.get('assignee')  or _ASSIGNEE).strip()
    platform  = (body.get('platform')  or '').strip()
    project_tag = (body.get('project_tag') or '').strip()

    if not hsd_url:
        return jsonify({'error': 'hsd_url is required'}), 400

    m = re.search(r'(\d{10,})', hsd_url)
    if not m:
        return jsonify({'error': 'Could not extract an HSD article ID from the URL'}), 400

    hsd_id  = m.group(1)
    session = _new_hsd_session()

    # ── Fetch HSD title ───────────────────────────────────────────────────────
    title = _fetch_hsd_title(session, hsd_id)

    # ── Find matching AR child ────────────────────────────────────────────────
    ar_raw      = _get_ar_children(session, hsd_id)
    ar_children = [_parse_ar_record(r) for r in ar_raw]
    matched     = [c for c in ar_children
                   if c['team'].lower() == 'i915_kmd'
                   and c['task'].lower() == 'development']

    if not matched:
        return jsonify({
            'error':       'No AR child found with team=i915_kmd and task=development',
            'ar_children': ar_children,
        }), 400

    child_hsd_id = matched[0]['id']

    # ── Load Jira token ───────────────────────────────────────────────────────
    token, err = _token_response()
    if err:
        return err

    post_issue, create_link = _make_jira_creator(token)
    result       = _build_jira_issues(post_issue, create_link,
                                      title, platform, component, assignee,
                                      hsd_id, child_hsd_id=child_hsd_id,
                                      project_tag=project_tag)
    created      = result['created']
    errors       = result['errors']
    epic_key     = result.get('epic_key')
    cc_story_key = result.get('cc_story_key')

    if epic_key is None:
        return jsonify({'error': 'Epic creation failed', 'created': created}), 500

    # ── Write Code Complete parent story key back to HSD AR child ─────────────
    writeback_key = cc_story_key or epic_key
    writeback_ok  = _writeback_jira_key(session, child_hsd_id, writeback_key)

    return jsonify({
        'created':       created,
        'errors':        errors,
        'epic_key':      epic_key,
        'cc_story_key':  cc_story_key,
        'writeback_key': writeback_key,
        'hsd_id':        hsd_id,
        'child_hsd_id':  child_hsd_id,
        'hsd_title':     title,
        'writeback_ok':  writeback_ok,
    })


@bp.route('/hsd2jira/create-from-ar-owner-hsd', methods=['POST'])
def hsd2jira_create_from_ar_owner_hsd():
    """
        Owner-based helper for all currently supported tenants.
        Scans the tenant's child records, matches owner, creates Jira using the
        parent HSD title + child title, and writes back using the tenant-specific field:
            - dg_soc.ar.jira_key
            - server.ar.jira_key
            - server_platf.ar.tag
            - ip_hw_graphics.sw_impact.sw_record
    """
    body      = request.get_json(force=True)
    hsd_url   = (body.get('hsd_url') or '').strip()
    owner     = (body.get('owner') or '').strip()
    platform  = (body.get('platform') or '').strip()
    component = (body.get('component') or _COMPONENT).strip()
    component2 = (body.get('component2') or '').strip()
    project_tag = (body.get('project_tag') or '').strip()

    if not hsd_url:
        return jsonify({'error': 'hsd_url is required'}), 400
    if not owner:
        return jsonify({'error': 'owner is required'}), 400

    m = re.search(r'(\d{10,})', hsd_url)
    if not m:
        return jsonify({'error': 'Could not extract an HSD article ID from the URL'}), 400

    hsd_id  = m.group(1)
    session = _new_hsd_session()
    info    = _fetch_hsd_info(session, hsd_id)
    parent_title = (info.get('title') or f'HSD {hsd_id}').strip()
    tenant_subject = f"{info.get('tenant','')}.{info.get('subject','')}".strip('.')

    if not platform:
        platform, _ = _detect_platform_and_tag(info)
    if not platform:
        return jsonify({'error': 'Could not auto-detect platform; please provide platform'}), 400

    token, err = _token_response()
    if err:
        return err

    post_issue, create_link = _make_jira_creator(token)
    components = [c for c in (component, component2) if c]

    def _extract_owner(raw: dict, keys: tuple[str, ...]) -> str:
        for k in keys:
            v = str(raw.get(k) or '').strip()
            if v and v.lower() not in ('none', 'null'):
                return v
        return ''

    # Build tenant-specific child targets, then apply owner matching only.
    targets: list[dict] = []

    if _is_server_feature(info):
        ar_raw = _get_server_ar_children(session, hsd_id)
        ar_by_id = {str(r.get('id', '')): r for r in ar_raw if isinstance(r, dict)}
        parsed = [_parse_server_ar_record(r) for r in ar_raw]
        for c in parsed:
            raw = ar_by_id.get(c['id'], {})
            targets.append({
                'child_hsd_id': c['id'],
                'title': c.get('title', ''),
                'record_owner': _extract_owner(raw, ('owner', 'server.ar.owner')),
                'writeback_existing': c.get('jira_key', '').strip(),
                'writeback_field': 'server.ar.jira_key',
                'writeback_kind': 'server',
            })

    elif _is_feature_hsd(info):
        ar_raw = _get_ar_children(session, hsd_id)
        ar_by_id = {str(r.get('id', '')): r for r in ar_raw if isinstance(r, dict)}
        parsed = [_parse_ar_record_extended(r) for r in ar_raw]
        for c in parsed:
            raw = ar_by_id.get(c['id'], {})
            targets.append({
                'child_hsd_id': c['id'],
                'title': c.get('title', '') or 'i915_kmd development',
                'record_owner': _extract_owner(raw, ('owner', 'dg_soc.ar.owner', 'dg.soc.ar.owner')),
                'writeback_existing': c.get('jira_key', '').strip(),
                'writeback_field': 'dg_soc.ar.jira_key',
                'writeback_kind': 'dg_soc',
            })

    elif _is_server_platf_feature(info):
        ar_raw = _get_server_platf_ar_children(session, hsd_id)
        ar_by_id = {str(r.get('id', '')): r for r in ar_raw if isinstance(r, dict)}
        parsed = [_parse_server_platf_ar(r) for r in ar_raw]
        for c in parsed:
            raw = ar_by_id.get(c['id'], {})
            targets.append({
                'child_hsd_id': c['id'],
                'title': c.get('title', ''),
                'record_owner': _extract_owner(raw, ('owner', 'server_platf.ar.owner')),
                'writeback_existing': c.get('tag', '').strip(),
                'writeback_field': 'server_platf.ar.tag',
                'writeback_kind': 'server_platf',
            })

    elif _is_sw_wa_hsd(info) or _is_ip_hw_feature_hsd(info):
        sw_raw = _get_sw_impact_children(session, hsd_id)
        sw_by_id = {str(r.get('id', '')): r for r in sw_raw if isinstance(r, dict)}
        parsed = [_parse_sw_impact_record(r) for r in sw_raw]
        for c in parsed:
            raw = sw_by_id.get(c['id'], {})
            child_title = (c.get('title') or '').strip()
            if not child_title:
                child_title = f"{c.get('sw_component', '').strip()} / {c.get('sw_task', '').strip()}".strip(' /')
            targets.append({
                'child_hsd_id': c['id'],
                'title': child_title,
                'record_owner': _extract_owner(raw, ('owner', 'ip_hw_graphics.sw_impact.owner')),
                'writeback_existing': c.get('sw_record', '').strip(),
                'writeback_field': 'ip_hw_graphics.sw_impact.sw_record',
                'writeback_kind': 'sw_impact',
            })

    else:
        return jsonify({
            'error': 'Unsupported tenant/subject for owner-based flow',
            'hsd_id': hsd_id,
            'tenant': tenant_subject,
        }), 400

    if not targets:
        return jsonify({
            'error': 'No eligible child records found for this tenant workflow',
            'hsd_id': hsd_id,
            'tenant': tenant_subject,
        }), 400

    results = []
    processed = skipped = total_jira_issues = 0

    for target in targets:
        child_id = target.get('child_hsd_id', '')
        ar_owner = str(target.get('record_owner') or '').strip()
        title    = str(target.get('title') or '').strip()
        existing = str(target.get('writeback_existing') or '').strip()
        jira_title = f'{parent_title} - {title}'

        item = {
            'child_hsd_id': child_id,
            'parent_title': parent_title,
            'title': title,
            'jira_title': jira_title,
            'ar_owner': ar_owner,
            'writeback_field': target.get('writeback_field', ''),
            'writeback_before': existing,
            'skipped': False,
            'skip_reason': '',
            'created_jira': [],
            'errors': [],
            'cc_story_key': None,
            'writeback_ok': False,
        }

        if ar_owner.lower() != owner.lower():
            item['skipped'] = True
            item['skip_reason'] = f"owner mismatch ({ar_owner or 'empty'})"
            skipped += 1
            results.append(item)
            continue

        if not ar_owner:
            item['skipped'] = True
            item['skip_reason'] = 'owner is empty on AR child'
            skipped += 1
            results.append(item)
            continue

        if not title:
            item['skipped'] = True
            item['skip_reason'] = 'title is empty on AR child'
            skipped += 1
            results.append(item)
            continue

        if existing:
            item['skipped'] = True
            item['skip_reason'] = f"{item['writeback_field']} already set ({existing})"
            skipped += 1
            results.append(item)
            continue

        jira = _build_jira_issues(
            post_issue,
            create_link,
            jira_title,
            platform,
            components,
            ar_owner,
            hsd_id,
            child_hsd_id=child_id,
            project_tag=project_tag,
        )
        item['created_jira']  = jira.get('created', [])
        item['errors']        = jira.get('errors', [])
        item['cc_story_key']  = jira.get('cc_story_key')
        total_jira_issues    += len(item['created_jira'])

        if item['cc_story_key']:
            kind = target.get('writeback_kind')
            if kind == 'server':
                item['writeback_ok'] = _writeback_server_ar_jira_key(session, child_id, item['cc_story_key'])
            elif kind == 'dg_soc':
                item['writeback_ok'] = _writeback_jira_key(session, child_id, item['cc_story_key'])
            elif kind == 'server_platf':
                item['writeback_ok'] = _writeback_tag_field(session, child_id, item['cc_story_key'])
            elif kind == 'sw_impact':
                item['writeback_ok'] = _writeback_sw_record(session, child_id, item['cc_story_key'])

        processed += 1
        results.append(item)

    return jsonify({
        'hsd_id': hsd_id,
        'tenant': tenant_subject,
        'owner_filter': owner,
        'platform_used': platform,
        'component_used': components,
        'processed_count': processed,
        'skipped_count': skipped,
        'total_jira_issues': total_jira_issues,
        'results': results,
    })


# ── HSD Query processing routes ─────────────────────────────────────────────

@bp.route('/hsd2jira/debug-feature/<hsd_id>')
def hsd2jira_debug_feature(hsd_id):
    """
    Debug: dump raw AR children of a dg_soc.feature HSD with extended field info.
    GET /hsd2jira/debug-feature/16029853828
    """
    session     = _new_hsd_session()
    info        = _fetch_hsd_info(session, hsd_id)
    ar_raw      = _get_ar_children(session, hsd_id)
    ar_extended = [_parse_ar_record_extended(r) for r in ar_raw]
    return jsonify({
        'hsd_id':      hsd_id,
        'info':        info,
        'ar_raw_count': len(ar_raw),
        'ar_extended': ar_extended,
        'ar_all_keys': sorted({k for r in ar_raw if isinstance(r, dict) for k in r.keys()}),
    })


@bp.route('/hsd2jira/debug-query/<query_id>')
def hsd2jira_debug_query(query_id):
    """
    Debug: execute a saved HSD query and return IDs found.
    GET /hsd2jira/debug-query/14027898530
    """
    session = _new_hsd_session()
    ids     = _execute_hsd_query(session, query_id)
    return jsonify({'query_id': query_id, 'ids': ids, 'count': len(ids) if ids else 0})


@bp.route('/hsd2jira/debug-server-platf/<hsd_id>')
def hsd2jira_debug_server_platf(hsd_id):
    """
    Debug: dump all fields of server_platf.ar children of hsd_id.
    GET /hsd2jira/debug-server-platf/16030623500
    """
    session  = _new_hsd_session()
    info     = _fetch_hsd_info(session, hsd_id)
    ar_raw   = _get_server_platf_ar_children(session, hsd_id)
    parsed   = [_parse_server_platf_ar(r) for r in ar_raw]
    all_keys = sorted({k for r in ar_raw if isinstance(r, dict) for k in r.keys()})
    # Also show fields that look like they could be tag/jira related
    tag_like = {k: ar_raw[0].get(k) for k in all_keys
                if ar_raw and any(x in k.lower() for x in ('tag', 'jira', 'key', 'link', 'ref'))}
    return jsonify({
        'hsd_id':    hsd_id,
        'info':      info,
        'ar_count':  len(ar_raw),
        'parsed':    parsed,
        'all_keys':  all_keys,
        'tag_like_fields': tag_like,
        'raw_first': ar_raw[0] if ar_raw else None,
    })



@bp.route('/hsd2jira/preview-hsd-query', methods=['POST'])
def hsd2jira_preview_hsd_query():
    """
    Preview what will happen for a query or single HSD URL — no Jira/HSD creation.
    Accepts either:
      • A query URL: https://hsdes.intel.com/...?queryId=14027898530
      • A single article URL: https://hsdes.intel.com/appstore/article-one/#/16029853828
    """
    body      = request.get_json(force=True)
    query_url = (body.get('query_url') or '').strip()
    if not query_url:
        return jsonify({'error': 'query_url is required'}), 400

    session = _new_hsd_session()

    query_id_match   = re.search(r'queryId=(\d+)', query_url)
    article_id_match = re.search(r'(\d{10,})', query_url)

    if query_id_match:
        query_id = query_id_match.group(1)
        hsd_ids  = _execute_hsd_query(session, query_id)
        if hsd_ids is None:
            return jsonify({'error': f'Failed to execute query {query_id}. '
                                      'Check that the query ID is valid and accessible.'}), 500
        if not hsd_ids:
            return jsonify({'error': 'Query returned no results', 'query_id': query_id}), 400
    elif article_id_match:
        query_id = None
        hsd_ids  = [article_id_match.group(1)]
    else:
        return jsonify({'error': 'Could not extract a queryId or HSD article ID from the URL'}), 400

    # Pre-fetch all HSD infos in parallel to avoid sequential round-trips
    with ThreadPoolExecutor(max_workers=16) as pool:
        info_map = dict(zip(hsd_ids, pool.map(lambda h: _fetch_hsd_info(session, h), hsd_ids)))

    items = []
    for hsd_id in hsd_ids:
        info   = info_map[hsd_id]
        item   = {
            'hsd_id':      hsd_id,
            'title':       info.get('title', f'HSD {hsd_id}'),
            'tenant':      f"{info.get('tenant','')} / {info.get('subject','')}".strip(' /'),
            'url':         f'https://hsdes.intel.com/appstore/article-one/#/{hsd_id}',
            'actionable':  False,
            'skip_reason': None,
            'core_ar':     None,
        }

        if _is_server_feature(info):
            # ── server.feature branch ─────────────────────────────────────────
            ar_raw    = _get_server_ar_children(session, hsd_id)
            ar_parsed = [_parse_server_ar_record(r) for r in ar_raw]

            # SW Arch Review children with valid component and exposure
            arch_reviews = [
                c for c in ar_parsed
                if c['title'].strip().lower() == 'sw arch review'
                and c['component'].lower() in _SERVER_AR_VALID_COMPONENTS
                and c['exposure'].lower() not in ('none', 'to_be_assigned', '')
            ]

            if not arch_reviews:
                item['skip_reason'] = (
                    f'No server.ar SW Arch Review children with component in '
                    f'{_SERVER_AR_VALID_COMPONENTS} and valid exposure found '
                    f'({len(ar_parsed)} AR records checked)')
                items.append(item)
                continue

            checks = []
            for arch_child in arch_reviews:
                comp = arch_child['component']
                dev_review = next(
                    (c for c in ar_parsed
                     if c['title'].strip().lower() == 'sw development review'
                     and c['component'].lower() == comp.lower()),
                    None
                )
                check = {
                    'arch_id':       arch_child['id'],
                    'component':     comp,
                    'exposure':      arch_child['exposure'],
                    'dev_review_id': dev_review['id']       if dev_review else None,
                    'dev_jira_key':  dev_review['jira_key'] if dev_review else None,
                    'dev_exposure':  dev_review['exposure'] if dev_review else None,
                    'will_process':  False,
                    'skip_reason':   None,
                }
                dev_exp = (dev_review.get('exposure', '') if dev_review else '').lower()
                if dev_review and dev_exp in ('none', 'to_be_assigned', ''):
                    check['skip_reason'] = (
                        f"SW Development Review ({dev_review['id']}) exposure="
                        f"'{dev_review['exposure']}' — no KMD action required")
                elif dev_review and dev_review['jira_key']:
                    check['skip_reason'] = (
                        f"SW Development Review ({dev_review['id']}) already has "
                        f"jira_key={dev_review['jira_key']}")
                else:
                    check['will_process'] = True
                    item['actionable'] = True
                checks.append(check)

            item['server_ar_checks'] = checks
            if not item['actionable']:
                item['skip_reason'] = 'All SW Development Review children already have jira_key set'
            items.append(item)
            continue

        if _is_server_platf_feature(info):
            # ── server_platf.feature branch ──────────────────────────────────
            ar_raw    = _get_server_platf_ar_children(session, hsd_id)
            ar_parsed = [_parse_server_platf_ar(r) for r in ar_raw]
            kmd_list  = [c for c in ar_parsed if _title_has_kmd(c['title'])]

            if not kmd_list:
                item['skip_reason'] = (
                    f'No server_platf.ar child with KMD in title found '
                    f'({len(ar_parsed)} AR records checked)')
            else:
                matched_ar = kmd_list[0]
                existing_tag = matched_ar.get('tag', '').strip()
                item['kmd_ar'] = {
                    'id':    matched_ar['id'],
                    'title': matched_ar['title'],
                    'tag':   existing_tag,
                }
                if existing_tag:
                    item['skip_reason'] = (
                        f"server_platf.ar child ({matched_ar['id']}) already has "
                        f"tag={existing_tag} — will skip in process")
                else:
                    item['actionable'] = True
            items.append(item)
            continue

        if _is_sw_wa_hsd(info):
            # ── ip_hw_graphics.bugeco branch (SW WA) ─────────────────────────
            sw_raw    = _get_sw_impact_children(session, hsd_id)
            sw_parsed = [_parse_sw_impact_record(r) for r in sw_raw]
            kmd_wa    = [c for c in sw_parsed
                         if 'i915_kmd' in c['sw_component'].lower()
                         and c['func_impact'].lower() == 'wa_needed']

            if not kmd_wa:
                item['skip_reason'] = (
                    f'No ip_hw_graphics.sw_impact child with sw_component=i915_kmd and '
                    f'func_impact=wa_needed ({len(sw_parsed)} sw_impact records checked)')
            else:
                matched_sw = kmd_wa[0]
                item['sw_impact_child'] = {
                    'id':           matched_sw['id'],
                    'sw_component': matched_sw['sw_component'],
                    'func_impact':  matched_sw['func_impact'],
                    'sw_record':    matched_sw['sw_record'],
                }
                if matched_sw['sw_record']:
                    item['skip_reason'] = (
                        f"sw_impact child ({matched_sw['id']}) already has "
                        f"sw_record={matched_sw['sw_record']} — will skip in process")
                else:
                    item['actionable'] = True
            items.append(item)
            continue

        if _is_ip_hw_feature_hsd(info):
            # ── ip_hw_graphics.feature branch ────────────────────────────────
            sw_raw    = _get_sw_impact_children(session, hsd_id)
            sw_parsed = [_parse_sw_impact_record(r) for r in sw_raw]

            core_arch = [c for c in sw_parsed
                         if c['sw_component'].lower() == 'core'
                         and c['os'].lower() == 'common'
                         and c['sw_task'].lower() == 'architecture'
                         and c['sw_exposure'].lower() not in ('to_be_assigned', 'none', '')]

            if not core_arch:
                item['skip_reason'] = (
                    f'No sw_impact child with sw_component=core, os=common, '
                    f'sw_task=architecture and valid exposure '
                    f'({len(sw_parsed)} sw_impact records checked)')
                items.append(item)
                continue

            core_child    = core_arch[0]
            arch_exposure = core_child['sw_exposure']
            item['core_arch'] = {'id': core_child['id'], 'sw_exposure': arch_exposure}

            kmd_dev = [c for c in sw_parsed
                       if c['sw_component'].lower() == 'i915_kmd'
                       and c['os'].lower() == 'linux'
                       and c['sw_task'].lower() == 'development']

            if kmd_dev:
                kmd_child    = kmd_dev[0]
                dev_exposure = kmd_child.get('sw_exposure', '').strip()
                item['kmd_dev_child'] = {
                    'id':          kmd_child['id'],
                    'sw_record':   kmd_child['sw_record'],
                    'sw_exposure': dev_exposure,
                }
                if dev_exposure.lower() in ('none', 'to_be_assigned', ''):
                    item['skip_reason'] = (
                        f"kmd dev child ({kmd_child['id']}) has sw_exposure='{dev_exposure}' "
                        f"— no KMD action required")
                elif kmd_child['sw_record']:
                    item['skip_reason'] = (
                        f"kmd dev child ({kmd_child['id']}) already has "
                        f"sw_record={kmd_child['sw_record']} — will skip in process")
                else:
                    item['actionable'] = True
            else:
                item['actionable']    = True
                item['kmd_dev_child'] = None  # will be created during process
            items.append(item)
            continue

        if not _is_feature_hsd(info):
            item['skip_reason'] = (f"Tenant/subject '{info.get('tenant','')}.{info.get('subject','')}'"
                                   " is not a supported type")
            items.append(item)
            continue

        ar_raw      = _get_ar_children(session, hsd_id)
        ar_children = [_parse_ar_record_extended(r) for r in ar_raw]
        arch_list    = [c for c in ar_children
                        if c['team'].lower() == 'i915_kmd'
                        and c['task'].lower() == 'architecture']
        kmd_dev_list = [c for c in ar_children
                        if c['team'].lower() == 'i915_kmd'
                        and c['task'].lower() == 'development']
        ual_dev_list = [c for c in ar_children
                        if c['team'].lower() == 'sw.ual_kmd'
                        and c['task'].lower() == 'development']

        if not arch_list and not kmd_dev_list and not ual_dev_list:
            item['skip_reason'] = (f'No AR child with team=i915_kmd or sw.ual_kmd found '
                                   f'({len(ar_children)} AR records checked)')
            items.append(item)
            continue

        item['release']       = info.get('release', '')
        item['auto_platform'] = _platform_from_release(info.get('release', ''))

        # ── i915_kmd check ────────────────────────────────────────────────────
        if arch_list or kmd_dev_list:
            if arch_list:
                arch     = arch_list[0]
                exposure = arch.get('exposure', '').strip()
                item['arch_ar'] = {
                    'id': arch['id'], 'team': arch['team'], 'task': arch['task'],
                    'exposure': exposure, 'type': arch.get('type_', ''), 'title': arch.get('title', ''),
                }
            else:
                dev_direct = kmd_dev_list[0]
                exposure   = dev_direct.get('exposure', '').strip()
                item['arch_ar'] = {
                    'id': dev_direct['id'], 'team': dev_direct['team'], 'task': dev_direct['task'],
                    'exposure': exposure, 'type': dev_direct.get('type_', ''), 'title': dev_direct.get('title', ''),
                    'note': 'no architecture child — using development child directly',
                }

            if exposure.lower() not in ('none', 'to_be_assigned'):
                if kmd_dev_list:
                    existing_dev      = kmd_dev_list[0]
                    existing_jira_key = existing_dev.get('jira_key', '').strip()
                    dev_exposure      = existing_dev.get('exposure', '').strip()
                    item['dev_ar'] = {'id': existing_dev['id'], 'jira_key': existing_jira_key,
                                      'exposure': dev_exposure}
                    if dev_exposure.lower() != 'none' and not existing_jira_key:
                        item['actionable'] = True
                else:
                    item['actionable'] = True

        # ── sw.ual_kmd check ──────────────────────────────────────────────────
        if ual_dev_list:
            ual_dev  = ual_dev_list[0]
            ual_exp  = ual_dev.get('exposure', '').strip()
            ual_jira = ual_dev.get('jira_key', '').strip()
            item['ual_dev_ar'] = {'id': ual_dev['id'], 'exposure': ual_exp, 'jira_key': ual_jira}
            if ual_exp.lower() not in ('none', 'to_be_assigned', '') and not ual_jira:
                item['actionable'] = True

        if not item['actionable'] and not item.get('skip_reason'):
            item['skip_reason'] = 'All AR children skipped (exposure=none/to_be_assigned or jira_key already set)'

        items.append(item)

    actionable_count = sum(1 for i in items if i['actionable'])
    return jsonify({
        'query_id':   query_id,
        'total':      len(items),
        'actionable': actionable_count,
        'skipped':    len(items) - actionable_count,
        'items':      items,
    })


@bp.route('/hsd2jira/process-hsd-query', methods=['POST'])
def hsd2jira_process_hsd_query():
    """
    Main query-processing endpoint.
    For each dg_soc.feature HSD in the query (or single article URL):
      1. Find AR child with team=i915_kmd and task=architecture
      2. If exposure is high/medium/low → Create 9 Jira issues
      3. Create new dg_soc.ar child HSD (team=i915_kmd, task=development) with CC story key written to jira_key
    """
    body        = request.get_json(force=True)
    query_url   = (body.get('query_url')    or '').strip()
    platform    = (body.get('platform')     or '').strip()
    _comp_raw   = body.get('components') or body.get('component') or _COMPONENT
    component   = _comp_raw if isinstance(_comp_raw, list) else _comp_raw.strip()
    assignee    = (body.get('assignee')     or _ASSIGNEE).strip()
    project_tag = (body.get('project_tag')  or '').strip()

    if not query_url:
        return jsonify({'error': 'query_url is required'}), 400

    session = _new_hsd_session()

    query_id_match   = re.search(r'queryId=(\d+)', query_url)
    article_id_match = re.search(r'(\d{10,})', query_url)

    if query_id_match:
        query_id = query_id_match.group(1)
        hsd_ids  = _execute_hsd_query(session, query_id)
        if hsd_ids is None:
            return jsonify({'error': f'Failed to execute query {query_id}'}), 500
        if not hsd_ids:
            return jsonify({'error': 'Query returned no results', 'query_id': query_id}), 400
    elif article_id_match:
        query_id = None
        hsd_ids  = [article_id_match.group(1)]
    else:
        return jsonify({'error': 'Could not extract a queryId or HSD article ID from the URL'}), 400

    token, err = _token_response()
    if err:
        return err

    post_issue, create_link_fn = _make_jira_creator(token)
    all_results: list = []

    for hsd_id in hsd_ids:
        info   = _fetch_hsd_info(session, hsd_id)
        result = _process_one_hsd(session, post_issue, create_link_fn,
                                   hsd_id, info,
                                   platform=platform, project_tag=project_tag,
                                   component=component, assignee=assignee,
                                   token=token)
        all_results.append(result)

    processed = [r for r in all_results if not r.get('skipped')]
    total_jira = sum(len(r.get('created_jira', [])) for r in processed)
    return jsonify({
        'query_id':          query_id,
        'total_hsd':         len(all_results),
        'processed_count':   len(processed),
        'skipped_count':     len(all_results) - len(processed),
        'total_jira_issues': total_jira,
        'results':           all_results,
    })


# ── Cron API routes ────────────────────────────────────────────────────────────

@bp.route('/hsd2jira/cron/status')
def hsd2jira_cron_status():
    """Return current cron configuration and status."""
    cfg = _load_cron_config()
    alive = bool(_cron_thread_ref[0] and _cron_thread_ref[0].is_alive())
    return jsonify({
        'enabled':          cfg.get('enabled', False),
        'thread_alive':     alive,
        'interval_minutes': cfg.get('interval_minutes', 30),
        'query_count':      len(cfg.get('query_urls', _CRON_DEFAULT_QUERIES)),
        'last_run':         cfg.get('last_run'),
        'last_result':      cfg.get('last_result'),
    })


@bp.route('/hsd2jira/cron/toggle', methods=['POST'])
def hsd2jira_cron_toggle():
    """Enable or disable the cron job."""
    body    = request.get_json(force=True) or {}
    cfg     = _load_cron_config()
    enabled = body.get('enabled', not cfg.get('enabled', False))
    cfg['enabled'] = bool(enabled)
    _save_cron_config(cfg)
    _ensure_cron_thread(cfg['enabled'])
    alive = bool(_cron_thread_ref[0] and _cron_thread_ref[0].is_alive())
    return jsonify({'enabled': cfg['enabled'], 'thread_alive': alive})


@bp.route('/hsd2jira/cron/update-queries', methods=['POST'])
def hsd2jira_cron_update_queries():
    """Replace the list of query URLs the cron will process."""
    body       = request.get_json(force=True) or {}
    query_urls = body.get('query_urls', [])
    if not isinstance(query_urls, list):
        return jsonify({'error': 'query_urls must be a list'}), 400
    query_urls = [u.strip() for u in query_urls if str(u).strip()]
    cfg = _load_cron_config()
    cfg['query_urls'] = query_urls
    _save_cron_config(cfg)
    return jsonify({'query_count': len(query_urls)})


@bp.route('/hsd2jira/cron/update-interval', methods=['POST'])
def hsd2jira_cron_update_interval():
    """Update cron interval in minutes."""
    body    = request.get_json(force=True) or {}
    minutes = int(body.get('interval_minutes', 30))
    if minutes < 5:
        return jsonify({'error': 'interval_minutes must be >= 5'}), 400
    cfg = _load_cron_config()
    cfg['interval_minutes'] = minutes
    _save_cron_config(cfg)
    return jsonify({'interval_minutes': minutes})


@bp.route('/hsd2jira/cron/log')
def hsd2jira_cron_log():
    """Return per-HSD detail log for the last _CRON_LOG_MAX_RUNS cron passes."""
    try:
        if _CRON_LOG_PATH.exists():
            with open(_CRON_LOG_PATH) as f:
                return jsonify(json.load(f))
    except Exception as exc:
        return jsonify({'error': str(exc)}), 500
    return jsonify([])


# Start cron thread on module load if config says enabled
_ensure_cron_thread(_load_cron_config().get('enabled', False))


# ── TREND/DONE weekly sync cron ───────────────────────────────────────────────

_TREND_CRON_CONFIG_PATH = _BASE_DIR / 'hsd2jira_trend_cron.json'
_TREND_CRON_LOG_PATH    = _BASE_DIR / 'hsd2jira_trend_cron_log.json'
_TREND_CRON_LOG_MAX     = 10

_TREND_GT_QUERIES: dict[str, str] = {
    '14024748599': 'CRI',
    '16026920807': 'JGS',
    '14026130398': 'TGS',
}

# server.feature queries for weekly TREND sync.
_TREND_SERVER_QUERIES: dict[str, str] = {
    '14028056822': 'JGS',   # server.feature Jaguar Shores
}

# dg_soc.feature queries for weekly TREND sync.
_TREND_DG_SOC_QUERIES: dict[str, str] = {
    '14024897996': 'CRI',   # dg_soc.feature Crescent Island
    '14027898530': 'TGS',   # dg_soc.feature Tiger Shores
}

# server_platf.feature queries for weekly TREND sync.
_TREND_SERVER_PLATF_QUERIES: dict[str, str] = {
    '13013902803': 'CRI',   # server_platf.feature Crescent Island
}

# ip_hw_graphics.bugeco queries for weekly TREND sync.
_TREND_BUGECO_QUERIES: dict[str, str] = {
    '16029952357': 'CRI',   # ip_hw_graphics.bugeco Crescent Island
}

_trend_cron_stop_event = threading.Event()
_trend_cron_thread_ref = [None]
_trend_cron_lock       = threading.Lock()


def _load_trend_cron_config() -> dict:
    try:
        if _TREND_CRON_CONFIG_PATH.exists():
            with open(_TREND_CRON_CONFIG_PATH) as f:
                return json.load(f)
    except Exception:
        pass
    return {'enabled': False, 'last_run': None, 'last_result': None}


def _save_trend_cron_config(cfg: dict) -> None:
    try:
        with open(_TREND_CRON_CONFIG_PATH, 'w') as f:
            json.dump(cfg, f, indent=2)
    except Exception as exc:
        logger.error('Failed to save trend cron config: %s', exc)


def _append_trend_cron_log(run_time: str, summary: dict, details: list) -> None:
    try:
        log = []
        if _TREND_CRON_LOG_PATH.exists():
            with open(_TREND_CRON_LOG_PATH) as f:
                log = json.load(f)
    except Exception:
        log = []
    log.insert(0, {'run_time': run_time, 'summary': summary, 'details': details})
    log = log[:_TREND_CRON_LOG_MAX]
    try:
        with open(_TREND_CRON_LOG_PATH, 'w') as f:
            json.dump(log, f, indent=2)
    except Exception as exc:
        logger.error('Failed to write trend cron log: %s', exc)


def _trend_sync_pass() -> dict:
    """Weekly TREND/DONE sync for ip_hw_graphics.feature and server.feature.
    No Jira creation — only syncs Jira Actual Trend WW / status back to HSD.
    Parallelises per-HSD work across 8 threads for fast execution.
    """
    token, err = _token_response()
    if err:
        return {'summary': {'error': 'token read failed', 'updated': 0, 'skipped': 0,
                            'errors': 1, 'total': 0}, 'details': []}

    # Enumerate IDs from all queries (sequential — just ID lists, fast)
    enum_session = _new_hsd_session()
    ip_hw_items = [
        (hsd_id, tag, 'ip_hw', query_id)
        for query_id, tag in _TREND_GT_QUERIES.items()
        for hsd_id in (_execute_hsd_query(enum_session, query_id) or [])
    ]
    server_items = [
        (hsd_id, tag, 'server', query_id)
        for query_id, tag in _TREND_SERVER_QUERIES.items()
        for hsd_id in (_execute_hsd_query(enum_session, query_id) or [])
    ]
    dg_soc_items = [
        (hsd_id, tag, 'dg_soc', query_id)
        for query_id, tag in _TREND_DG_SOC_QUERIES.items()
        for hsd_id in (_execute_hsd_query(enum_session, query_id) or [])
    ]
    server_platf_items = [
        (hsd_id, tag, 'server_platf', query_id)
        for query_id, tag in _TREND_SERVER_PLATF_QUERIES.items()
        for hsd_id in (_execute_hsd_query(enum_session, query_id) or [])
    ]
    bugeco_items = [
        (hsd_id, tag, 'bugeco', query_id)
        for query_id, tag in _TREND_BUGECO_QUERIES.items()
        for hsd_id in (_execute_hsd_query(enum_session, query_id) or [])
    ]
    all_items = ip_hw_items + server_items + dg_soc_items + server_platf_items + bugeco_items

    def _process_ip_hw(hsd_id: str, tag: str, query_id: str = '') -> list[dict]:
        """ip_hw_graphics.feature: sync TREND/DONE on i915_kmd/linux/development child."""
        session = _new_hsd_session()
        try:
            title   = _fetch_hsd_info(session, hsd_id).get('title', f'HSD {hsd_id}')[:70]
            kmd_dev = [
                c for c in
                (_parse_sw_impact_record(r) for r in _get_sw_impact_children(session, hsd_id))
                if  c['sw_component'].lower() == 'i915_kmd'
                and c['os'].lower()           == 'linux'
                and c['sw_task'].lower()      == 'development'
            ]
            if not kmd_dev:
                return []

            child    = kmd_dev[0]
            child_id = child['id']
            sw_rec   = child.get('sw_record', '').strip()
            exposure = child.get('sw_exposure', '').strip()
            done     = child.get('done', '').strip()
            base     = {'hsd_id': hsd_id, 'child_id': child_id, 'tag': tag,
                        'title': title, 'sw_record': sw_rec, 'query_id': query_id}

            if not sw_rec:
                return [{**base, 'action': 'skipped', 'reason': 'no sw_record'}]
            if exposure.lower() == 'none':
                return [{**base, 'action': 'skipped', 'reason': 'sw_exposure=none'}]
            if done.lower() == 'yes':
                return [{**base, 'action': 'skipped', 'reason': 'already Done=Yes'}]

            jira_status, trend_ww = _fetch_jira_status_and_trend(token, sw_rec)
            if jira_status.lower() == 'closed':
                ok = _update_sw_impact_field(session, child_id,
                                             'ip_hw_graphics.sw_impact.done', 'Yes')
                return [{**base, 'action': 'set_done', 'value': 'Yes',
                         'prev_value': child.get('done', ''), 'ok': ok}]
            if trend_ww:
                ok = _update_sw_impact_field(session, child_id,
                                             'ip_hw_graphics.sw_impact.trend', trend_ww)
                return [{**base, 'action': 'update_trend', 'value': trend_ww,
                         'prev_value': child.get('trend', ''), 'ok': ok}]
            return [{**base, 'action': 'skipped', 'reason': 'no Actual Trend WW on Jira'}]

        except Exception as exc:
            logger.error('Trend cron ip_hw: error on HSD %s: %s', hsd_id, exc)
            return [{'hsd_id': hsd_id, 'child_id': '', 'tag': tag,
                     'title': f'HSD {hsd_id}', 'sw_record': '',
                     'action': 'error', 'reason': str(exc)}]

    def _process_server(hsd_id: str, tag: str, query_id: str = '') -> list[dict]:
        """server.feature: sync new_projected_dates / reason on SW Development Review children."""
        session = _new_hsd_session()
        try:
            info  = _fetch_hsd_info(session, hsd_id)
            if not _is_server_feature(info):
                return []
            title     = info.get('title', f'HSD {hsd_id}')[:70]
            ar_parsed = [_parse_server_ar_record(r)
                         for r in _get_server_ar_children(session, hsd_id)]

            # Only SW Development Review children with valid component, Jira key, and valid exposure
            dev_reviews = [
                c for c in ar_parsed
                if  c['title'].strip().lower() == 'sw development review'
                and c['component'].lower() in _SERVER_AR_VALID_COMPONENTS
                and c.get('jira_key', '').strip()
                and c.get('exposure', '').lower() not in ('none', 'to_be_assigned', '')
            ]
            if not dev_reviews:
                return []

            entries: list[dict] = []
            for dev in dev_reviews:
                child_id = dev['id']
                jira_key = dev['jira_key']
                base = {'hsd_id': hsd_id, 'child_id': child_id,
                        'tag': dev.get('component', tag),
                        'title': title, 'sw_record': jira_key, 'query_id': query_id}

                jira_status, trend_ww = _fetch_jira_status_and_trend(token, jira_key)
                if jira_status.lower() == 'closed':
                    ok = _update_server_ar_field(session, child_id,
                                                 'server.ar.reason', 'verified')
                    entries.append({**base, 'action': 'set_verified',
                                    'value': 'verified',
                                    'prev_value': dev.get('reason', ''), 'ok': ok})
                elif trend_ww:
                    ok = _update_server_ar_field(session, child_id,
                                                 'server.ar.new_projected_dates', trend_ww)
                    entries.append({**base, 'action': 'update_projected_date',
                                    'value': trend_ww,
                                    'prev_value': dev.get('new_projected_dates', ''), 'ok': ok})
                else:
                    entries.append({**base, 'action': 'skipped',
                                    'reason': 'no Actual Trend WW on Jira'})
            return entries

        except Exception as exc:
            logger.error('Trend cron server: error on HSD %s: %s', hsd_id, exc)
            return [{'hsd_id': hsd_id, 'child_id': '', 'tag': tag,
                     'title': f'HSD {hsd_id}', 'sw_record': '',
                     'action': 'error', 'reason': str(exc)}]

    def _process_dg_soc(hsd_id: str, tag: str, query_id: str = '') -> list[dict]:
        """dg_soc.feature: write Actual Trend WW to phase, or phase=Done if Jira closed."""
        session = _new_hsd_session()
        try:
            info = _fetch_hsd_info(session, hsd_id)
            if not _is_feature_hsd(info):
                return []
            title      = info.get('title', f'HSD {hsd_id}')[:70]
            ar_children = [_parse_ar_record_extended(r)
                           for r in _get_ar_children(session, hsd_id)]

            # i915_kmd/development children with jira_key set and valid exposure
            dev_list = [
                c for c in ar_children
                if  c['team'].lower() == 'i915_kmd'
                and c['task'].lower() == 'development'
                and c.get('jira_key', '').strip()
                and c.get('exposure', '').lower() not in ('none', 'to_be_assigned', '')
            ]
            if not dev_list:
                return []

            entries: list[dict] = []
            for dev in dev_list:
                child_id = dev['id']
                jira_key = dev['jira_key']
                base = {'hsd_id': hsd_id, 'child_id': child_id, 'tag': tag,
                        'title': title, 'sw_record': jira_key, 'query_id': query_id}

                jira_status, trend_ww = _fetch_jira_status_and_trend(token, jira_key)
                if jira_status.lower() == 'closed':
                    ok = _update_dg_soc_ar_field(session, child_id,
                                                 'dg_soc.ar.phase', 'Done')
                    entries.append({**base, 'action': 'set_phase_done',
                                    'value': 'Done',
                                    'prev_value': dev.get('phase', ''), 'ok': ok})
                elif trend_ww:
                    ok = _update_dg_soc_ar_field(session, child_id,
                                                 'dg_soc.ar.phase', trend_ww)
                    entries.append({**base, 'action': 'update_eta',
                                    'value': trend_ww,
                                    'prev_value': dev.get('phase', ''), 'ok': ok})
                else:
                    entries.append({**base, 'action': 'skipped',
                                    'reason': 'no Actual Trend WW on Jira'})
            return entries

        except Exception as exc:
            logger.error('Trend cron dg_soc: error on HSD %s: %s', hsd_id, exc)
            return [{'hsd_id': hsd_id, 'child_id': '', 'tag': tag,
                     'title': f'HSD {hsd_id}', 'sw_record': '',
                     'action': 'error', 'reason': str(exc)}]

    def _process_server_platf(hsd_id: str, tag: str, query_id: str = '') -> list[dict]:
        """server_platf.feature: write Actual Trend WW to server_platf.ar.trend on KMD child."""
        session = _new_hsd_session()
        try:
            title    = _fetch_hsd_info(session, hsd_id).get('title', f'HSD {hsd_id}')[:70]
            ar_raw   = _get_server_platf_ar_children(session, hsd_id)
            ar_list  = [_parse_server_platf_ar(r) for r in ar_raw]

            # KMD children that are signed off and have a VLK Jira key in the tag field
            kmd_list = [
                c for c in ar_list
                if _title_has_kmd(c['title'])
                and c.get('status', '').lower() == 'sign_off'
                and re.search(r'VLK-\d+', c.get('tag', ''))
            ]
            if not kmd_list:
                return []

            entries: list[dict] = []
            for child in kmd_list:
                child_id = child['id']
                jira_key = re.search(r'VLK-\d+', child['tag']).group(0)
                base = {'hsd_id': hsd_id, 'child_id': child_id, 'tag': tag,
                        'title': title, 'sw_record': jira_key, 'query_id': query_id}

                jira_status, trend_ww = _fetch_jira_status_and_trend(token, jira_key)
                if jira_status.lower() == 'closed':
                    ok = _update_server_platf_ar_field(session, child_id,
                                                       'server_platf.ar.trend', 'Done')
                    entries.append({**base, 'action': 'set_done',
                                    'value': 'Done',
                                    'prev_value': child.get('trend', ''), 'ok': ok})
                elif trend_ww:
                    ok = _update_server_platf_ar_field(session, child_id,
                                                       'server_platf.ar.trend', trend_ww)
                    entries.append({**base, 'action': 'update_trend',
                                    'value': trend_ww,
                                    'prev_value': child.get('trend', ''), 'ok': ok})
                else:
                    entries.append({**base, 'action': 'skipped',
                                    'reason': 'no Actual Trend WW on Jira'})
            return entries

        except Exception as exc:
            logger.error('Trend cron server_platf: error on HSD %s: %s', hsd_id, exc)
            return [{'hsd_id': hsd_id, 'child_id': '', 'tag': tag,
                     'title': f'HSD {hsd_id}', 'sw_record': '',
                     'action': 'error', 'reason': str(exc)}]

    def _process_bugeco(hsd_id: str, tag: str, query_id: str = '') -> list[dict]:
        """ip_hw_graphics.bugeco: sync TREND/DONE on i915_kmd sw_impact child with wa_needed."""
        session = _new_hsd_session()
        try:
            info = _fetch_hsd_info(session, hsd_id)
            if not _is_sw_wa_hsd(info):
                return []
            title = info.get('title', f'HSD {hsd_id}')[:70]
            kmd_wa = [
                c for c in
                (_parse_sw_impact_record(r) for r in _get_sw_impact_children(session, hsd_id))
                if  c['sw_component'].lower() == 'i915_kmd'
                and c['func_impact'].lower()  == 'wa_needed'
                and re.search(r'VLK-\d+', c.get('sw_record', ''))
            ]
            if not kmd_wa:
                return []

            entries: list[dict] = []
            for child in kmd_wa:
                child_id = child['id']
                sw_rec   = re.search(r'VLK-\d+', child['sw_record']).group(0)
                done     = child.get('done', '').strip()
                base     = {'hsd_id': hsd_id, 'child_id': child_id, 'tag': tag,
                            'title': title, 'sw_record': sw_rec, 'query_id': query_id}

                if done.lower() == 'yes':
                    entries.append({**base, 'action': 'skipped', 'reason': 'already Done=Yes'})
                    continue

                jira_status, trend_ww = _fetch_jira_status_and_trend(token, sw_rec)
                if jira_status.lower() == 'closed':
                    ok = _update_sw_impact_field(session, child_id,
                                                 'ip_hw_graphics.sw_impact.done', 'Yes')
                    entries.append({**base, 'action': 'set_done', 'value': 'Yes',
                                    'prev_value': child.get('done', ''), 'ok': ok})
                elif trend_ww:
                    ok = _update_sw_impact_field(session, child_id,
                                                 'ip_hw_graphics.sw_impact.trend', trend_ww)
                    entries.append({**base, 'action': 'update_trend', 'value': trend_ww,
                                    'prev_value': child.get('trend', ''), 'ok': ok})
                else:
                    entries.append({**base, 'action': 'skipped',
                                    'reason': 'no Actual Trend WW on Jira'})
            return entries

        except Exception as exc:
            logger.error('Trend cron bugeco: error on HSD %s: %s', hsd_id, exc)
            return [{'hsd_id': hsd_id, 'child_id': '', 'tag': tag,
                     'title': f'HSD {hsd_id}', 'sw_record': '',
                     'action': 'error', 'reason': str(exc)}]

    # Process all HSDs in parallel (8 workers); each worker returns a list of entries
    _dispatch = {'ip_hw': _process_ip_hw, 'server': _process_server,
                 'dg_soc': _process_dg_soc, 'server_platf': _process_server_platf,
                 'bugeco': _process_bugeco}
    details: list[dict] = []
    with ThreadPoolExecutor(max_workers=8) as pool:
        futures = {
            pool.submit(_dispatch[ttype], hsd_id, tag, query_id): (hsd_id, tag, query_id)
            for hsd_id, tag, ttype, query_id in all_items
        }
        for fut in as_completed(futures):
            try:
                entries = fut.result()
            except Exception as exc:
                hsd_id, tag, query_id = futures[fut]
                entries = [{'hsd_id': hsd_id, 'child_id': '', 'tag': tag,
                             'title': f'HSD {hsd_id}', 'sw_record': '',
                             'action': 'error', 'reason': str(exc),
                             'query_id': query_id}]
            details.extend(entries)

    _update_actions = ('set_done', 'update_trend', 'set_verified',
                       'update_projected_date', 'set_phase_done', 'update_eta')
    updated     = sum(1 for d in details if d['action'] in _update_actions)
    skipped     = sum(1 for d in details if d['action'] == 'skipped')
    error_count = sum(1 for d in details if d['action'] == 'error')
    return {
        'summary': {'updated': updated, 'skipped': skipped,
                    'errors': error_count, 'total': len(details)},
        'details': details,
    }


def _trend_cron_worker(stop_event: threading.Event) -> None:
    logger.info('Trend-sync cron worker started')
    while not stop_event.is_set():
        cfg = _load_trend_cron_config()
        if not cfg.get('enabled'):
            stop_event.wait(60)
            continue
        now       = datetime.utcnow()
        last_run  = cfg.get('last_run', '') or ''
        today_str = now.strftime('%Y-%m-%d')
        if now.weekday() != 4:          # not Friday
            stop_event.wait(3600)
            continue
        if last_run.startswith(today_str):   # already ran this Friday
            stop_event.wait(3600)
            continue
        logger.info('Trend-sync cron: starting TREND/DONE pass')
        try:
            result   = _trend_sync_pass()
            run_time = datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')
            _append_trend_cron_log(run_time, result['summary'], result['details'])
            cfg = _load_trend_cron_config()
            cfg['last_run']    = run_time
            cfg['last_result'] = result['summary']
            _save_trend_cron_config(cfg)
            # Per-query change log
            _q_names = {}
            for _qid, _qtag in _TREND_GT_QUERIES.items():
                _q_names[_qid] = f'ip_hw_graphics.feature [{_qtag}]  queryId={_qid}'
            for _qid, _qtag in _TREND_SERVER_QUERIES.items():
                _q_names[_qid] = f'server.feature [{_qtag}]  queryId={_qid}'
            for _qid, _qtag in _TREND_DG_SOC_QUERIES.items():
                _q_names[_qid] = f'dg_soc.feature [{_qtag}]  queryId={_qid}'
            for _qid, _qtag in _TREND_SERVER_PLATF_QUERIES.items():
                _q_names[_qid] = f'server_platf.feature [{_qtag}]  queryId={_qid}'
            for _qid, _qtag in _TREND_BUGECO_QUERIES.items():
                _q_names[_qid] = f'ip_hw_graphics.bugeco [{_qtag}]  queryId={_qid}'
            _update_acts = frozenset(('set_done', 'update_trend', 'set_verified',
                                      'update_projected_date', 'set_phase_done', 'update_eta'))
            _by_q: dict = {}
            for _d in result['details']:
                _by_q.setdefault(_d.get('query_id', 'unknown'), []).append(_d)
            for _qid in sorted(_by_q):
                _changed = [_d for _d in _by_q[_qid]
                            if _d.get('action') in _update_acts and _d.get('ok')]
                if not _changed:
                    continue
                logger.info('Trend-sync ── %s ── %d change(s)',
                            _q_names.get(_qid, _qid), len(_changed))
                for _d in _changed:
                    _prev = _d.get('prev_value') or '(none)'
                    _new  = _d.get('value', '')
                    logger.info('  HSD %-14s  Jira %-10s  %-26s  %s -> %s',
                                _d['hsd_id'], _d.get('sw_record', '—'),
                                _d.get('action', ''), _prev, _new)
            logger.info('Trend-sync cron: pass complete — %s', result['summary'])
        except Exception as exc:
            logger.exception('Trend-sync cron: pass failed: %s', exc)
        stop_event.wait(3600)
    logger.info('Trend-sync cron worker stopped')


def _ensure_trend_cron_thread(enabled: bool) -> None:
    with _trend_cron_lock:
        t = _trend_cron_thread_ref[0]
        if enabled:
            if not t or not t.is_alive():
                _trend_cron_stop_event.clear()
                nt = threading.Thread(target=_trend_cron_worker,
                                      args=(_trend_cron_stop_event,), daemon=True)
                nt.start()
                _trend_cron_thread_ref[0] = nt
        else:
            _trend_cron_stop_event.set()


@bp.route('/hsd2jira/trend-cron/status')
def hsd2jira_trend_cron_status():
    cfg   = _load_trend_cron_config()
    alive = bool(_trend_cron_thread_ref[0] and _trend_cron_thread_ref[0].is_alive())
    return jsonify({
        'enabled':     cfg.get('enabled', False),
        'thread_alive': alive,
        'last_run':    cfg.get('last_run'),
        'last_result': cfg.get('last_result'),
        'queries':     list(_TREND_GT_QUERIES.values()),   # ['CRI','JGS','TGS']
    })


@bp.route('/hsd2jira/trend-cron/toggle', methods=['POST'])
def hsd2jira_trend_cron_toggle():
    body    = request.get_json(force=True) or {}
    cfg     = _load_trend_cron_config()
    enabled = bool(body.get('enabled', not cfg.get('enabled', False)))
    cfg['enabled'] = enabled
    _save_trend_cron_config(cfg)
    _ensure_trend_cron_thread(enabled)
    alive = bool(_trend_cron_thread_ref[0] and _trend_cron_thread_ref[0].is_alive())
    return jsonify({'enabled': enabled, 'thread_alive': alive})


@bp.route('/hsd2jira/trend-cron/log')
def hsd2jira_trend_cron_log():
    try:
        if _TREND_CRON_LOG_PATH.exists():
            with open(_TREND_CRON_LOG_PATH) as f:
                return jsonify(json.load(f))
    except Exception as exc:
        return jsonify({'error': str(exc)}), 500
    return jsonify([])


@bp.route('/hsd2jira/trend-cron/download-log')
def hsd2jira_trend_cron_download_log():
    """Download the most recent trend sync run as a CSV file."""
    import csv, io
    try:
        if not _TREND_CRON_LOG_PATH.exists():
            return 'No log yet', 404
        with open(_TREND_CRON_LOG_PATH) as f:
            log = json.load(f)
        if not log:
            return 'No log yet', 404
        latest  = log[0]
        details = latest.get('details', [])
        buf = io.StringIO()
        w   = csv.writer(buf)
        w.writerow(['Feature HSD', 'Child HSD', 'Project', 'Title',
                    'Jira (sw_record)', 'Action', 'Value', 'OK', 'Reason'])
        for d in details:
            w.writerow([
                d.get('hsd_id', ''), d.get('child_id', ''), d.get('tag', ''),
                d.get('title', ''), d.get('sw_record', ''),
                d.get('action', ''), d.get('value', ''), d.get('ok', ''),
                d.get('reason', ''),
            ])
        run_time = latest.get('run_time', 'unknown').replace(' ', '_').replace(':', '')
        filename = f'trend_sync_{run_time}.csv'
        from flask import Response
        return Response(buf.getvalue(), mimetype='text/csv',
                        headers={'Content-Disposition': f'attachment; filename="{filename}"'})
    except Exception as exc:
        return str(exc), 500


# Start trend-sync cron thread on module load if config says enabled
_ensure_trend_cron_thread(_load_trend_cron_config().get('enabled', False))

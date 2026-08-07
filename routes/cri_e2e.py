"""
CRI E2E Plan — /cri-e2e

Loads CRI issues from SharePoint Excel (weekly_updates, KMD tabs) and PPTX
(XPUM/Sysman features) and renders the e2e_plan.html template.
GT_DCN data is frozen into a local SQLite DB — independent of the source Excel.
The SW-only / HW+SW tables are live-fetched from the SharePoint Excel with a
30-minute local cache; falls back to the previous cached copy on errors.
"""
import json
import re
import logging
import sqlite3
import time
import base64
import openpyxl
import requests as _requests
from contextlib import contextmanager
from datetime import datetime, timedelta
from io import BytesIO
from pathlib import Path

from flask import Blueprint, render_template, request, jsonify, send_file

from config import CRI_EXCEL, E2E_COMMENTS, HSD_BASE_URL
from utils.cri_helpers import refresh_pptx, extract_jiras, chips, CRI_PPTX

logger = logging.getLogger(__name__)
bp = Blueprint('cri_e2e', __name__)

# Saved Jiras from the CRI Weekly Review dashboard (read-only file access, no HTTP dep)
_DASH_COMMENTS = Path('/home/vitasta/CRI_Weekly_review/comments.json')

# SharePoint live Excel ─────────────────────────────────────────────────────
_SP_SITE       = "VTTSWSOSGC"
_SP_FILE       = "CRI Pre-Map Day work Template.xlsx"
_SP_CACHE      = Path('/home/vitasta/CRI_Weekly_review') / _SP_FILE
_TOKEN_CACHE   = Path('/home/vitasta/triage/repos/Projects/soc_automation/.graph_token_cache.json')
_SP_TTL_MIN    = 30   # minutes before re-downloading
# Drive + item IDs from Graph API (stable GUIDs — don't need search)
_SP_DRIVE_ID   = "b!DoVu96jyeUi6KrTdnfiyOkD3l4uMELNPlWLzwTS0EvRl5FXkvP32Rr3npiSn6z7D"
_SP_ITEM_ID    = "257A7CFA-3248-4C7C-AFC1-C5D9E2599C1B"

# Pending Driver Features Excel — Xpum_Sysman source
_PDF_ITEM_GUID = "2E76F07E-4149-4533-A345-47ECCF471B52"
_PDF_CACHE     = Path('/home/vitasta/CRI_Weekly_review') / 'Pending_Driver_Features.xlsx'
_PDF_TTL_MIN   = 60   # minutes before re-downloading


def _get_graph_token() -> str | None:
    """Return a valid Graph API access token, auto-refreshing via refresh_token."""
    if not _TOKEN_CACHE.exists():
        return None
    try:
        data = json.loads(_TOKEN_CACHE.read_text())
        if time.time() < data.get('expires_at', 0) - 300:
            return data['access_token']
        # Token expired or about to expire — refresh it
        resp = _requests.post(
            "https://login.microsoftonline.com/46c98d88-e344-4ed4-8496-4ed7712e255d/oauth2/v2.0/token",
            data={
                'client_id':     'd3590ed6-52b3-4102-aeff-aad2292ab01c',
                'scope':         'https://graph.microsoft.com/.default offline_access',
                'refresh_token': data['refresh_token'],
                'grant_type':    'refresh_token',
            }, timeout=15,
        )
        if resp.status_code == 200:
            tok      = resp.json()
            updated  = {
                'access_token':  tok['access_token'],
                'refresh_token': tok.get('refresh_token', data['refresh_token']),
                'expires_at':    (datetime.now() + timedelta(seconds=tok.get('expires_in', 3600))).timestamp(),
            }
            _TOKEN_CACHE.write_text(json.dumps(updated, indent=2))
            logger.info("Graph token refreshed successfully")
            return updated['access_token']
        logger.warning("Graph token refresh HTTP %d: %s", resp.status_code, resp.text[:200])
    except Exception as exc:
        logger.warning("Graph token refresh failed: %s", exc)
    return None


def _fetch_sp_excel() -> Path:
    """
    Download the SharePoint Excel to a local cache (TTL=30 min).
    Returns the path to the cached file (fresh or existing).
    Falls back to CRI_EXCEL if no cache exists and download fails.
    """
    if _SP_CACHE.exists():
        age_min = (time.time() - _SP_CACHE.stat().st_mtime) / 60
        if age_min < _SP_TTL_MIN:
            return _SP_CACHE   # cache is fresh

    token = _get_graph_token()
    if not token:
        logger.warning("No Graph token — using %s",
                       _SP_CACHE.name if _SP_CACHE.exists() else CRI_EXCEL.name)
        return _SP_CACHE if _SP_CACHE.exists() else CRI_EXCEL

    hdrs = {'Authorization': f'Bearer {token}'}

    # Direct download by stable drive + item GUIDs (no search needed)
    dl_url = (f"https://graph.microsoft.com/v1.0/drives/{_SP_DRIVE_ID}"
              f"/items/{_SP_ITEM_ID}/content")
    try:
        resp = _requests.get(dl_url, headers=hdrs, timeout=30,
                             stream=True, allow_redirects=True)
        if resp.status_code == 200:
            _SP_CACHE.parent.mkdir(parents=True, exist_ok=True)
            with open(_SP_CACHE, 'wb') as fh:
                for chunk in resp.iter_content(8192):
                    fh.write(chunk)
            logger.info("SharePoint Excel downloaded: %.1f KB",
                        _SP_CACHE.stat().st_size / 1024)
            return _SP_CACHE
        logger.warning("SharePoint download HTTP %d — using cached copy",
                       resp.status_code)
    except Exception as exc:
        logger.warning("SharePoint download error: %s — using cached copy", exc)

    return _SP_CACHE if _SP_CACHE.exists() else CRI_EXCEL


# ── GT DCN frozen store ───────────────────────────────────────────────────────
GTDCN_DB = Path(__file__).parent.parent / "cri_gtdcn.db"

@contextmanager
def _gtdcn_db():
    conn = sqlite3.connect(str(GTDCN_DB))
    conn.row_factory = sqlite3.Row
    try:
        yield conn
        conn.commit()
    finally:
        conn.close()

def _init_gtdcn_db() -> None:
    with _gtdcn_db() as conn:
        conn.execute("""
            CREATE TABLE IF NOT EXISTS gtdcn_frozen (
                hsd_id  TEXT PRIMARY KEY,
                title   TEXT NOT NULL DEFAULT '',
                jiras   TEXT NOT NULL DEFAULT '[]',
                notes   TEXT NOT NULL DEFAULT ''
            )
        """)

_init_gtdcn_db()

def _load_gtdcn_frozen() -> list[dict] | None:
    """Return frozen rows or None if the table is empty (not yet seeded)."""
    with _gtdcn_db() as conn:
        rows = conn.execute(
            "SELECT hsd_id, title, jiras, notes FROM gtdcn_frozen ORDER BY rowid"
        ).fetchall()
    if not rows:
        return None
    result = []
    for r in rows:
        jira_list = json.loads(r["jiras"]) if r["jiras"] else []
        result.append({
            "hsd_id":   r["hsd_id"],
            "hsd_link": HSD_BASE_URL + r["hsd_id"],
            "title":    r["title"],
            "chips":    chips(jira_list),
            "notes":    r["notes"],
            "_jiras":   jira_list,
        })
    return result

def _seed_gtdcn(rows: list[dict]) -> None:
    """Write rows to the frozen store, replacing any existing data."""
    with _gtdcn_db() as conn:
        conn.execute("DELETE FROM gtdcn_frozen")
        for r in rows:
            conn.execute(
                "INSERT INTO gtdcn_frozen VALUES (?,?,?,?)",
                (r["hsd_id"], r["title"],
                 json.dumps(r.get("_jiras", [])),
                 r.get("notes", ""))
            )
    logger.info("GT DCN frozen: %d rows saved", len(rows))



def _load_e2e_comments() -> dict:
    if E2E_COMMENTS.exists():
        with open(E2E_COMMENTS) as f:
            return json.load(f)
    return {}


def _save_e2e_comments(c: dict) -> None:
    with open(E2E_COMMENTS, 'w') as f:
        json.dump(c, f, indent=2)


# ── Pending Driver Features Excel loader (Xpum_Sysman tab) ─────────────────

def _fetch_pending_driver_excel() -> Path | None:
    """Download Pending Driver Features.xlsx via Graph shares API (TTL=60 min)."""
    if _PDF_CACHE.exists():
        age_min = (time.time() - _PDF_CACHE.stat().st_mtime) / 60
        if age_min < _PDF_TTL_MIN:
            return _PDF_CACHE
    token = _get_graph_token()
    if not token:
        logger.warning("No Graph token — cannot refresh Pending Driver Features Excel")
        return _PDF_CACHE if _PDF_CACHE.exists() else None
    share_url = (f"https://intel.sharepoint.com/sites/druid/_layouts/15/Doc.aspx"
                 f"?sourcedoc=%7B{_PDF_ITEM_GUID}%7D")
    encoded = base64.b64encode(share_url.encode()).decode().rstrip('=').replace('+', '-').replace('/', '_')
    share_token = f"u!{encoded}"
    hdrs = {'Authorization': f'Bearer {token}'}
    try:
        meta = _requests.get(
            f"https://graph.microsoft.com/v1.0/shares/{share_token}/driveItem",
            headers=hdrs, timeout=20)
        if meta.status_code != 200:
            logger.warning("PDF Excel meta HTTP %d", meta.status_code)
            return _PDF_CACHE if _PDF_CACHE.exists() else None
        dl_url = meta.json().get('@microsoft.graph.downloadUrl')
        if not dl_url:
            return _PDF_CACHE if _PDF_CACHE.exists() else None
        resp = _requests.get(dl_url, timeout=60)
        if resp.status_code == 200:
            _PDF_CACHE.parent.mkdir(parents=True, exist_ok=True)
            _PDF_CACHE.write_bytes(resp.content)
            logger.info("Pending Driver Features Excel downloaded: %.1f KB",
                        _PDF_CACHE.stat().st_size / 1024)
            return _PDF_CACHE
        logger.warning("PDF Excel download HTTP %d", resp.status_code)
    except Exception as exc:
        logger.warning("PDF Excel download error: %s", exc)
    return _PDF_CACHE if _PDF_CACHE.exists() else None


def _load_xpum_sysman_excel() -> list[dict]:
    """Load Xpum_Sysman tab cols A (feature) + D (jiras) from Pending Driver Features Excel."""
    path = _fetch_pending_driver_excel()
    if not path:
        return []
    try:
        wb_pdf = openpyxl.load_workbook(str(path), data_only=True)
        ws_pdf = next(
            (wb_pdf[n] for n in wb_pdf.sheetnames if n.lower() == 'xpum_sysman'),
            None
        )
        if not ws_pdf:
            logger.warning("Xpum_Sysman sheet not found; sheets: %s", wb_pdf.sheetnames)
            return []
        rows = []
        for row in ws_pdf.iter_rows(min_row=3, max_row=ws_pdf.max_row,
                                    min_col=1, max_col=4, values_only=True):
            feature = str(row[0] or '').strip()
            if not feature:
                continue
            jira_raw = str(row[3] or '').strip()
            jiras_l  = extract_jiras(jira_raw) if jira_raw else []
            safe_key = re.sub(r'[^a-z0-9]+', '_', feature.lower())[:35].strip('_')
            rows.append({
                'feature':    feature,
                'chips':      chips(jiras_l),
                'key_status': f'e2e_xsys_status_{safe_key}',
                'key_cmt':    f'e2e_xsys_cmt_{safe_key}',
            })
        return rows
    except Exception as exc:
        logger.error("Failed to load Xpum_Sysman Excel: %s", exc)
        return []


# ── PPTX XPUM loader (CRI tab) ───────────────────────────────────────────────

def _load_pptx_xpum() -> list[dict]:
    if not CRI_PPTX.exists():
        refresh_pptx()
    if not CRI_PPTX.exists():
        return []
    try:
        from pptx import Presentation
    except ImportError:
        return []

    prs = Presentation(str(CRI_PPTX))
    rows = []
    section_map = {1: 'RAS', 2: 'Telemetry', 3: 'AMC', 4: 'Firmware', 5: 'Opens'}

    for slide_idx, slide in enumerate(prs.slides):
        if slide_idx == 0:
            continue
        section = section_map.get(slide_idx, '')
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            for ri, row in enumerate(shape.table.rows):
                if ri == 0:
                    continue
                cells = [c.text.replace('\xa0', ' ').strip() for c in row.cells]
                if not any(cells):
                    continue
                feature  = cells[0] if len(cells) > 0 else ''
                jira_raw = cells[1] if len(cells) > 1 else ''
                status   = cells[2] if len(cells) > 2 else ''
                dep_raw  = cells[3] if len(cells) > 3 else ''
                if slide_idx == 5:
                    blocker_eta = overall_eta = ''
                    remarks  = cells[3] if len(cells) > 3 else ''
                    dep_raw  = cells[2] if len(cells) > 2 else ''
                else:
                    blocker_eta = cells[4] if len(cells) > 4 else ''
                    overall_eta = cells[5] if len(cells) > 5 else ''
                    remarks     = cells[6] if len(cells) > 6 else ''
                jiras      = extract_jiras(jira_raw) if jira_raw else []
                dep_jiras  = extract_jiras(dep_raw)  if dep_raw  else []
                if not jiras and not feature:
                    continue
                main_jira = jiras[0] if jiras else ''
                rows.append({
                    'section':     section,
                    'feature':     feature,
                    'chips':       chips(jiras),
                    'status':      status,
                    'dep_chips':   chips(dep_jiras),
                    'dep_raw':     dep_raw if not dep_jiras else '',
                    'blocker_eta': blocker_eta,
                    'overall_eta': overall_eta,
                    'remarks':     remarks,
                    'key_bmg':     f'e2e_xbmg_{main_jira or feature[:20]}',
                    'key_cmt':     f'e2e_xcmt_{main_jira or feature[:20]}',
                })
    return rows


# ── Excel CRI loader ──────────────────────────────────────────────────────────

def load_e2e_data() -> tuple[list, list, list, str]:
    excel_path = _fetch_sp_excel()
    # Record when the Excel file was last updated (for display in template)
    try:
        fetched_at = datetime.fromtimestamp(excel_path.stat().st_mtime).strftime('%Y-%m-%d %H:%M')
    except Exception:
        fetched_at = ''
    wb = openpyxl.load_workbook(str(excel_path), data_only=True)

    # KMD tab: jira map per HSD
    kmd_jira_map, kmd_title_map = {}, {}
    ws_kmd = wb['KMD']
    for row in ws_kmd.iter_rows(min_row=3, max_col=12, values_only=True):
        if not row[0]:
            continue
        try:
            hid = str(int(row[0]))
        except Exception:
            continue
        if row[1]:
            kmd_title_map[hid] = row[1]
        jiras = []
        for v in row[7:12]:
            jiras += extract_jiras(v)
        kmd_jira_map[hid] = list(dict.fromkeys(jiras))

    # EC_Features_SOC tab — col A: marker or HSD ID, col B: title, col C: Jira(s)
    ws    = wb['EC_Features_SOC']
    cri_rows, section = [], None
    SKIP  = {'22019854607'}
    for row in ws.iter_rows(min_row=1, max_row=ws.max_row, min_col=1, max_col=3, values_only=True):
        a = row[0]
        if str(a) in ('sw_only', 'hw+sw'):
            section = str(a)
            continue
        if a is None:
            continue
        try:
            hid = str(int(a))
        except Exception:
            continue
        if hid in SKIP:
            continue
        col_c_jiras = extract_jiras(str(row[2] or ''))
        all_jiras = list(dict.fromkeys(kmd_jira_map.get(hid, []) + col_c_jiras))
        cri_rows.append({
            'hsd_id':   hid,
            'hsd_link': HSD_BASE_URL + hid,
            'title':    str(row[1] or kmd_title_map.get(hid, ''))[:120],
            'section':  section or 'sw_only',
            'chips':    chips(all_jiras),
            'key_bmg':    f'e2e_bmg_{hid}',
            'key_cmt':    f'e2e_cmt_{hid}',
            'key_status': f'e2e_status_{hid}',
        })

    # GT_DCN — use frozen SQLite store; seed from Excel+comments on first run
    gtdcn_rows = _load_gtdcn_frozen()
    if gtdcn_rows is None:
        # First run: seed from Excel and :5050 comments, then freeze
        try:
            dash_c = json.loads(_DASH_COMMENTS.read_text()) if _DASH_COMMENTS.exists() else {}
        except Exception:
            dash_c = {}
        gtdcn_rows = []
        ws_gt = wb['GT_DCN']
        for row in ws_gt.iter_rows(min_row=2, max_row=ws_gt.max_row,
                                   min_col=1, max_col=4, values_only=True):
            if not row[0]:
                continue
            try:
                hid = str(int(row[0]))
            except Exception:
                continue
            gt_jiras = list(dict.fromkeys(extract_jiras(row[2]) + extract_jiras(row[3])))
            saved_raw = dash_c.get(f'gtdcn_jiras_{hid}') or dash_c.get(f'gtdcn_status_{hid}', '')
            saved_jiras = [j for j in re.split(r'[\s,;]+', saved_raw.strip())
                           if j and re.match(r'^[A-Z]+-\d+$', j)]
            all_jiras = list(dict.fromkeys(kmd_jira_map.get(hid, []) + gt_jiras + saved_jiras))
            gtdcn_rows.append({
                'hsd_id':   hid,
                'hsd_link': HSD_BASE_URL + hid,
                'title':    str(row[1] or '')[:120],
                'chips':    chips(all_jiras),
                'notes':    dash_c.get(f'gtdcn_{hid}', ''),
                '_jiras':   all_jiras,
            })
        _seed_gtdcn(gtdcn_rows)


    # XPUM/Sysman: load from Pending Driver Features Excel (Xpum_Sysman tab)
    xpum_rows = _load_xpum_sysman_excel()
    if not xpum_rows:
        xpum_rows = _load_pptx_xpum()

    return cri_rows, xpum_rows, gtdcn_rows, fetched_at


# ── Routes ────────────────────────────────────────────────────────────────────

@bp.route('/cri-e2e')
def cri_e2e():
    cri_rows, xpum_rows, gtdcn_rows, fetched_at = load_e2e_data()
    comments = _load_e2e_comments()
    for r in cri_rows:
        r['bmg']        = comments.get(r.get('key_bmg', ''), '')
        r['comment']    = comments.get(r.get('key_cmt', ''), '')
        r['status_val'] = comments.get(r.get('key_status', ''), '')
    for r in gtdcn_rows:
        r['key_status'] = f'e2e_gtdcn_status_{r["hsd_id"]}'
        r['status_val'] = comments.get(r['key_status'], '')
        r['key_notes']  = f'e2e_gtdcn_notes_{r["hsd_id"]}'
        r['notes']      = comments.get(r['key_notes'], r.get('notes', ''))
    for r in xpum_rows:
        r['bmg']        = comments.get(r.get('key_bmg', ''), '')
        r['comment']    = comments.get(r.get('key_cmt', ''), '')
        if 'key_status' not in r:
            r['key_status'] = re.sub(r'_cmt_', '_status_', r.get('key_cmt', ''))
        r['status_val'] = comments.get(r.get('key_status', ''), '')
    return render_template('e2e_plan.html', cri_rows=cri_rows, xpum_rows=xpum_rows,
                           gtdcn_rows=gtdcn_rows, comments=comments,
                           sp_fetched_at=fetched_at, sp_file=_SP_FILE)


@bp.route('/cri-e2e/save', methods=['POST'])
def cri_e2e_save():
    data = request.get_json()
    if not data or 'key' not in data:
        return jsonify({'status': 'error'}), 400
    c = _load_e2e_comments()
    c[data['key']] = data.get('value', '')
    _save_e2e_comments(c)
    return jsonify({'status': 'ok'})


@bp.route('/cri-e2e/gtdcn-export')
def gtdcn_export():
    """Download the frozen GT DCN data as an Excel file."""
    rows = _load_gtdcn_frozen() or []
    wb_out = openpyxl.Workbook()
    ws_out = wb_out.active
    ws_out.title = 'GT_DCN'
    ws_out.append(['HSD ID', 'Title', 'JIRAs', 'Notes'])
    for r in rows:
        ws_out.append([
            r['hsd_id'],
            r['title'],
            ', '.join(r.get('_jiras', [])),
            r.get('notes', ''),
        ])
    buf = BytesIO()
    wb_out.save(buf)
    buf.seek(0)
    return send_file(
        buf,
        download_name='GT_DCN_frozen.xlsx',
        as_attachment=True,
        mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
    )


# ── Excel export helpers ───────────────────────────────────────────────────

def _chips_str(chip_list: list) -> str:
    return ', '.join(c['id'] for c in (chip_list or []))


def _make_export_data():
    """Load all three tables with comments filled in — for export routes."""
    cri, xpum, gtdcn, _ = load_e2e_data()
    comments = _load_e2e_comments()
    for r in cri:
        r['bmg']        = comments.get(r.get('key_bmg', ''), '')
        r['comment']    = comments.get(r.get('key_cmt', ''), '')
        r['status_val'] = comments.get(r.get('key_status', ''), '')
    for r in gtdcn:
        r['status_val'] = comments.get(f'e2e_gtdcn_status_{r["hsd_id"]}', '')
        r['notes']      = comments.get(f'e2e_gtdcn_notes_{r["hsd_id"]}', r.get('notes', ''))
    for r in xpum:
        r['comment']    = comments.get(r.get('key_cmt', ''), '')
        if 'key_status' not in r:
            r['key_status'] = re.sub(r'_cmt_', '_status_', r.get('key_cmt', ''))
        r['status_val'] = comments.get(r.get('key_status', ''), '')
    return cri, gtdcn, xpum


def _style_ws(ws, col_widths: list) -> None:
    """Bold blue header, freeze row 1, colour Status cells, set column widths."""
    from openpyxl.styles import PatternFill, Font, Alignment
    from openpyxl.utils import get_column_letter
    hdr_fill = PatternFill('solid', fgColor='1F4E79')
    hdr_font = Font(bold=True, color='FFFFFF', size=10)
    for cell in ws[1]:
        cell.fill      = hdr_fill
        cell.font      = hdr_font
        cell.alignment = Alignment(horizontal='center', vertical='center')
    ws.freeze_panes = 'A2'
    ws.row_dimensions[1].height = 18
    status_col = next((c.column for c in ws[1] if c.value == 'Status'), None)
    status_fills = {
        'Done':        PatternFill('solid', fgColor='C6EFCE'),
        'In Progress': PatternFill('solid', fgColor='DDEBF7'),
        'Not Started': PatternFill('solid', fgColor='F2F2F2'),
    }
    for row in ws.iter_rows(min_row=2):
        for cell in row:
            cell.alignment = Alignment(vertical='top', wrap_text=True)
            if status_col and cell.column == status_col and cell.value in status_fills:
                cell.fill = status_fills[cell.value]
    for i, w in enumerate(col_widths, 1):
        ws.column_dimensions[get_column_letter(i)].width = w


_XLSX_MIME = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'


@bp.route('/cri-e2e/export/cri')
def export_cri_table():
    """Export CRI Issues table as Excel."""
    cri, _, _ = _make_export_data()
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = 'CRI Issues'
    ws.append(['HSD ID', 'Type', 'Title', 'JIRA(s)', 'Status', 'Supported BMG', 'Comments'])
    for r in cri:
        ws.append([r['hsd_id'], r.get('section', ''), r.get('title', ''),
                   _chips_str(r.get('chips', [])), r.get('status_val', ''),
                   r.get('bmg', ''), r.get('comment', '')])
    _style_ws(ws, [12, 8, 46, 26, 12, 14, 32])
    buf = BytesIO(); wb.save(buf); buf.seek(0)
    return send_file(buf, download_name='CRI_Issues.xlsx', as_attachment=True, mimetype=_XLSX_MIME)


@bp.route('/cri-e2e/export/gtdcn')
def export_gtdcn_table():
    """Export GT DCN table (with editable Status + Notes) as Excel."""
    _, gtdcn, _ = _make_export_data()
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = 'GT DCN'
    ws.append(['HSD ID', 'Title', 'JIRA(s)', 'Status', 'Notes'])
    for r in gtdcn:
        ws.append([r['hsd_id'], r.get('title', ''), _chips_str(r.get('chips', [])),
                   r.get('status_val', ''), r.get('notes', '')])
    _style_ws(ws, [12, 50, 26, 12, 42])
    buf = BytesIO(); wb.save(buf); buf.seek(0)
    return send_file(buf, download_name='GT_DCN.xlsx', as_attachment=True, mimetype=_XLSX_MIME)


@bp.route('/cri-e2e/export/xpum')
def export_xpum_table():
    """Export XPUM/Sysman table as Excel."""
    _, _, xpum = _make_export_data()
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = 'XPUM Sysman'
    ws.append(['Feature', 'JIRA(s)', 'Status', 'Comments'])
    for r in xpum:
        ws.append([r.get('feature', ''), _chips_str(r.get('chips', [])),
                   r.get('status_val', ''), r.get('comment', '')])
    _style_ws(ws, [46, 26, 12, 42])
    buf = BytesIO(); wb.save(buf); buf.seek(0)
    return send_file(buf, download_name='XPUM_Sysman.xlsx', as_attachment=True, mimetype=_XLSX_MIME)


@bp.route('/cri-e2e/export/all')
def export_all_tables():
    """Export all three tables as one Excel workbook with separate sheets."""
    cri, gtdcn, xpum = _make_export_data()
    wb = openpyxl.Workbook()

    ws1 = wb.active
    ws1.title = 'CRI Issues'
    ws1.append(['HSD ID', 'Type', 'Title', 'JIRA(s)', 'Status', 'Supported BMG', 'Comments'])
    for r in cri:
        ws1.append([r['hsd_id'], r.get('section', ''), r.get('title', ''),
                    _chips_str(r.get('chips', [])), r.get('status_val', ''),
                    r.get('bmg', ''), r.get('comment', '')])
    _style_ws(ws1, [12, 8, 46, 26, 12, 14, 32])

    ws2 = wb.create_sheet('GT DCN')
    ws2.append(['HSD ID', 'Title', 'JIRA(s)', 'Status', 'Notes'])
    for r in gtdcn:
        ws2.append([r['hsd_id'], r.get('title', ''), _chips_str(r.get('chips', [])),
                    r.get('status_val', ''), r.get('notes', '')])
    _style_ws(ws2, [12, 50, 26, 12, 42])

    ws3 = wb.create_sheet('XPUM Sysman')
    ws3.append(['Feature', 'JIRA(s)', 'Status', 'Comments'])
    for r in xpum:
        ws3.append([r.get('feature', ''), _chips_str(r.get('chips', [])),
                    r.get('status_val', ''), r.get('comment', '')])
    _style_ws(ws3, [46, 26, 12, 42])

    buf = BytesIO(); wb.save(buf); buf.seek(0)
    return send_file(buf, download_name='CRI_E2E_Plan.xlsx', as_attachment=True, mimetype=_XLSX_MIME)


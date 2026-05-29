"""
CRI Daily TF — /cri-daily-tf

Downloads the latest Pending Driver Features Excel from SharePoint via the
Microsoft Graph API, parses the Xpum_Sysman and XeKMD_pending tabs, extracts
the latest WW entry from TF Meeting Minutes.

All columns are editable and persisted server-side (SQLite).
A 30-second polling endpoint lets every user see changes made by others.
"""
import io
import re
import json
import base64
import sqlite3
import logging
import requests
from contextlib import contextmanager
from datetime import datetime
from pathlib import Path

from flask import Blueprint, render_template, request, jsonify, send_file

from config import SOC_AUTOMATION_DIR, BASE_DIR

logger = logging.getLogger(__name__)
bp = Blueprint('cri_daily_tf', __name__)

# ── Constants ─────────────────────────────────────────────────────────────────
_SHARING_URL = (
    'https://intel.sharepoint.com/:x:/r/sites/druid/Shared%20Documents/'
    'Crescent%20Island/PSXT/Pending%20Driver%20Features%20enabling.xlsx'
    '?d=w2e76f07e41494533a34547eccf471b52&csf=1&web=1&e=vPlFTH'
)
_PROXIES = {
    'http':  'http://proxy-chain.intel.com:911',
    'https': 'http://proxy-chain.intel.com:912',
}
_CACHE_PATH = Path('/tmp/cri_daily_tf_cache.xlsx')
_CACHE_META  = Path('/tmp/cri_daily_tf_cache_meta.json')
_CACHE_TTL   = 60           # 1 minute excel cache (kept for burst protection)
_DB_PATH     = BASE_DIR / 'cri_daily_tf.db'
_VALID_FIELDS = {'feature', 'priority', 'status', 'eta', 'blocker'}


# ── SQLite helpers ────────────────────────────────────────────────────────────

@contextmanager
def _db_conn():
    conn = sqlite3.connect(str(_DB_PATH), check_same_thread=False)
    conn.row_factory = sqlite3.Row
    try:
        yield conn
        conn.commit()
    finally:
        conn.close()


def _init_db() -> None:
    with _db_conn() as conn:
        conn.execute('''
            CREATE TABLE IF NOT EXISTS tf_edits (
                row_id     TEXT PRIMARY KEY,
                feature    TEXT NOT NULL DEFAULT '',
                priority   TEXT NOT NULL DEFAULT '',
                status     TEXT NOT NULL DEFAULT '',
                eta        TEXT NOT NULL DEFAULT '',
                blocker    TEXT NOT NULL DEFAULT '',
                updated_at TEXT NOT NULL DEFAULT ''
            )
        ''')


def _load_edits() -> dict:
    """Return {row_id: {feature, priority, status, eta, blocker, updated_at}}."""
    _init_db()
    with _db_conn() as conn:
        rows = conn.execute(
            'SELECT row_id, feature, priority, status, eta, blocker, updated_at FROM tf_edits'
        ).fetchall()
    return {r['row_id']: dict(r) for r in rows}


def _save_edit(row_id: str, field: str, value: str) -> None:
    if field not in _VALID_FIELDS:
        raise ValueError(f'Invalid field: {field!r}')
    _init_db()
    now = datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')
    with _db_conn() as conn:
        existing = conn.execute(
            'SELECT feature, priority, status, eta, blocker FROM tf_edits WHERE row_id=?',
            (row_id,)
        ).fetchone()
        if existing:
            d = dict(existing)
            d[field] = value
            conn.execute(
                'UPDATE tf_edits SET feature=?, priority=?, status=?, eta=?, blocker=?, updated_at=?'
                ' WHERE row_id=?',
                (d['feature'], d['priority'], d['status'], d['eta'], d['blocker'], now, row_id)
            )
        else:
            d = {f: '' for f in _VALID_FIELDS}
            d[field] = value
            conn.execute(
                'INSERT INTO tf_edits (row_id, feature, priority, status, eta, blocker, updated_at)'
                ' VALUES (?,?,?,?,?,?,?)',
                (row_id, d['feature'], d['priority'], d['status'], d['eta'], d['blocker'], now)
            )


# ── Graph token ───────────────────────────────────────────────────────────────

def _get_graph_token() -> str | None:
    token_cache = SOC_AUTOMATION_DIR / '.graph_token_cache.json'
    if not token_cache.exists():
        return None
    td = json.load(open(token_cache))
    if datetime.now().timestamp() > td.get('expires_at', 0):
        try:
            import sys
            sys.path.insert(0, str(SOC_AUTOMATION_DIR))
            from refresh_graph_token import refresh_token
            if not refresh_token():
                return None
            td = json.load(open(token_cache))
        except Exception as e:
            logger.warning('Token refresh failed: %s', e)
            return None
    return td.get('access_token')


# ── Excel download with 5-min cache ──────────────────────────────────────────

def _fetch_excel() -> bytes | None:
    if _CACHE_PATH.exists() and _CACHE_META.exists():
        try:
            meta = json.load(open(_CACHE_META))
            if datetime.now().timestamp() < meta.get('expires_at', 0):
                return _CACHE_PATH.read_bytes()
        except Exception:
            pass

    token = _get_graph_token()
    if not token:
        return _CACHE_PATH.read_bytes() if _CACHE_PATH.exists() else None

    encoded = 'u!' + base64.urlsafe_b64encode(_SHARING_URL.encode()).rstrip(b'=').decode()
    headers = {'Authorization': f'Bearer {token}', 'Accept': 'application/json'}
    try:
        r = requests.get(
            f'https://graph.microsoft.com/v1.0/shares/{encoded}/driveItem',
            headers=headers, proxies=_PROXIES, timeout=20,
        )
        if r.status_code != 200:
            return _CACHE_PATH.read_bytes() if _CACHE_PATH.exists() else None
        dl_url = r.json().get('@microsoft.graph.downloadUrl')
        if not dl_url:
            return None
        r2 = requests.get(dl_url, proxies=_PROXIES, stream=True, timeout=30)
        if r2.status_code != 200:
            return _CACHE_PATH.read_bytes() if _CACHE_PATH.exists() else None
        data = b''.join(r2.iter_content(65536))
        _CACHE_PATH.write_bytes(data)
        with open(_CACHE_META, 'w') as f:
            json.dump({'expires_at': datetime.now().timestamp() + _CACHE_TTL}, f)
        return data
    except Exception as e:
        logger.warning('Excel fetch failed: %s', e)
        return _CACHE_PATH.read_bytes() if _CACHE_PATH.exists() else None


# ── Parsing helpers ───────────────────────────────────────────────────────────

def _extract_latest_ww(text) -> str:
    """Return the latest WW entry text from TF Meeting Minutes."""
    if not text:
        return ''
    t = str(text).strip()
    parts = re.split(r'(?=ww\s*\d+[\.\d]*\s*[:\-→>])', t, flags=re.IGNORECASE)
    entries = []
    for part in parts:
        m = re.match(
            r'(ww\s*(\d+)(?:\.(\d+))?(?:\.(\d+))?)\s*[:\-→>]+\s*(.*)',
            part, re.IGNORECASE | re.DOTALL,
        )
        if m:
            major = int(m.group(2))
            minor = int(m.group(3)) if m.group(3) else 0
            sub   = int(m.group(4)) if m.group(4) else 0
            entries.append((major, minor, sub, m.group(5).strip()))
    if not entries:
        return t
    entries.sort(key=lambda x: (x[0], x[1], x[2]))
    return entries[-1][3]


def _parse_xpum_sysman(ws) -> list[dict]:
    rows = []
    for i, row in enumerate(ws.iter_rows(values_only=True)):
        if i < 2:
            continue
        col_a = row[0] if row else None
        if not col_a:
            continue
        tf_min   = row[2]  if len(row) > 2  else None
        priority = row[12] if len(row) > 12 else None
        rows.append({
            'row_id':   f'xpum_{i}',
            'feature':  str(col_a).strip(),
            'priority': str(priority).strip() if priority else '',
            'status':   _extract_latest_ww(tf_min),
            'eta':      '',
            'blocker':  '',
        })
    return rows


def _parse_xekmd_pending(ws) -> list[dict]:
    rows = []
    for i, row in enumerate(ws.iter_rows(values_only=True)):
        if i < 1:
            continue
        col_a = row[0] if row else None
        if not col_a:
            continue
        col_a_s = str(col_a).strip()
        if col_a_s.lower().startswith(('other complete', 'target for po', 'page migration ',
                                        'feature - celestial', '[punit]', '[cls g10b + cri]: assert',
                                        '[cls g10b + cri] pcie peer')):
            continue
        tf_min   = row[1] if len(row) > 1 else None
        priority = row[8] if len(row) > 8 else None
        rows.append({
            'row_id':   f'xekmd_{i}',
            'feature':  col_a_s,
            'priority': str(priority).strip() if priority else '',
            'status':   _extract_latest_ww(tf_min),
            'eta':      '',
            'blocker':  '',
        })
    return rows


def _merge_edits(excel_rows: list[dict], edits: dict) -> list[dict]:
    """Overlay server-saved edits onto Excel data. Non-empty edits win."""
    merged = []
    for r in excel_rows:
        rid = r['row_id']
        e   = edits.get(rid, {})
        merged.append({
            'row_id':     rid,
            'feature':    e.get('feature') or r['feature'],
            'priority':   e.get('priority') or r['priority'],
            'status':     e.get('status')   or r['status'],
            'eta':        e.get('eta',  ''),
            'blocker':    e.get('blocker', ''),
            'updated_at': e.get('updated_at', ''),
        })
    return merged


def _load_data() -> tuple[list[dict], list[dict], str]:
    import openpyxl
    data = _fetch_excel()
    if not data:
        return [], [], 'Error: could not fetch Excel from SharePoint'
    wb    = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    xpum  = _parse_xpum_sysman(wb['Xpum_Sysman'])    if 'Xpum_Sysman'   in wb.sheetnames else []
    xekmd = _parse_xekmd_pending(wb['XeKMD_pending']) if 'XeKMD_pending' in wb.sheetnames else []
    edits = _load_edits()
    return _merge_edits(xpum, edits), _merge_edits(xekmd, edits), ''


# ── Routes ────────────────────────────────────────────────────────────────────

@bp.route('/cri-daily-tf')
def cri_daily_tf():
    xpum_rows, xekmd_rows, error = _load_data()
    return render_template(
        'cri_daily_tf.html',
        xpum_rows=xpum_rows,
        xekmd_rows=xekmd_rows,
        error=error,
        last_updated=datetime.now().strftime('%B %d, %Y %H:%M'),
    )


@bp.route('/cri-daily-tf/save', methods=['POST'])
def cri_daily_tf_save():
    body = request.get_json(force=True, silent=True) or {}
    row_id = str(body.get('row_id', '')).strip()
    field  = str(body.get('field',  '')).strip()
    value  = str(body.get('value',  '')).strip()
    if not row_id or not field:
        return jsonify({'error': 'row_id and field required'}), 400
    if field not in _VALID_FIELDS:
        return jsonify({'error': f'field must be one of {sorted(_VALID_FIELDS)}'}), 400
    try:
        _save_edit(row_id, field, value)
        return jsonify({'ok': True, 'updated_at': datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')})
    except Exception as exc:
        logger.exception('Save failed')
        return jsonify({'error': str(exc)}), 500


@bp.route('/cri-daily-tf/api/data')
def cri_daily_tf_api():
    """Polling endpoint — returns all current edits for the 30-second sync."""
    edits = _load_edits()
    return jsonify(edits)


@bp.route('/cri-daily-tf/reset-field', methods=['POST'])
def cri_daily_tf_reset_field():
    """Clear a single field's DB override so Excel value shows through."""
    body   = request.get_json(force=True, silent=True) or {}
    row_id = str(body.get('row_id', '')).strip()
    field  = str(body.get('field',  '')).strip()
    if not row_id or field not in _VALID_FIELDS:
        return jsonify({'error': 'row_id and valid field required'}), 400
    _init_db()
    with _db_conn() as conn:
        conn.execute(f"UPDATE tf_edits SET {field}='', updated_at='' WHERE row_id=?", (row_id,))
    return jsonify({'ok': True})


@bp.route('/cri-daily-tf/force-refresh', methods=['POST'])
def cri_daily_tf_force_refresh():
    """Delete the Excel cache so the next request fetches fresh from SharePoint."""
    try:
        if _CACHE_PATH.exists():
            _CACHE_PATH.unlink()
        if _CACHE_META.exists():
            _CACHE_META.unlink()
        # eagerly fetch so the next page load is fast
        _fetch_excel()
        return jsonify({'ok': True, 'fetched_at': datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')})
    except Exception as exc:
        logger.exception('Force-refresh failed')
        return jsonify({'ok': False, 'error': str(exc)}), 500


# ── Export to PowerPoint ──────────────────────────────────────────────────────

def _build_pptx(xpum_rows: list[dict], xekmd_rows: list[dict]) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    INTEL_BLUE  = RGBColor(0x00, 0x68, 0xB5)
    INTEL_DARK  = RGBColor(0x00, 0x3C, 0x71)
    WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_BLUE  = RGBColor(0xE8, 0xF4, 0xFD)
    PURPLE      = RGBColor(0x6A, 0x1B, 0x9A)
    TEAL        = RGBColor(0x00, 0x69, 0x5C)
    ROW_ALT     = RGBColor(0xF4, 0xF9, 0xFF)
    BORDER      = RGBColor(0xD0, 0xDC, 0xE8)
    TEXT_DARK   = RGBColor(0x1A, 0x25, 0x30)
    MUTED       = RGBColor(0x6B, 0x7B, 0x8D)
    RED         = RGBColor(0xC0, 0x39, 0x2B)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # completely blank

    ww_label = datetime.utcnow().strftime('WW%W.%u  ·  %d %b %Y')

    def _add_slide(title_text: str, accent: RGBColor, rows: list[dict]) -> None:
        slide = prs.slides.add_slide(blank)
        W, H = prs.slide_width, prs.slide_height

        # ── header band ──────────────────────────────────────────────────────
        hdr = slide.shapes.add_shape(1, 0, 0, W, Inches(0.72))  # MSO_SHAPE_TYPE.RECTANGLE=1
        hdr.fill.solid(); hdr.fill.fore_color.rgb = INTEL_DARK
        hdr.line.fill.background()

        def _txb(left, top, width, height, text, size, bold=False,
                 color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
            tb = slide.shapes.add_textbox(left, top, width, height)
            tf = tb.text_frame; tf.word_wrap = wrap
            p  = tf.paragraphs[0]; p.alignment = align
            run = p.add_run(); run.text = text
            run.font.size = Pt(size); run.font.bold = bold
            run.font.color.rgb = color
            return tb

        _txb(Inches(0.18), Inches(0.08), Inches(8), Inches(0.32),
             title_text, 14, bold=True)
        _txb(Inches(0.18), Inches(0.38), Inches(8), Inches(0.26),
             'Pending Driver Features  ·  Exec Review', 8, color=RGBColor(0xB3, 0xD7, 0xF5))
        _txb(Inches(9.5), Inches(0.08), Inches(3.6), Inches(0.28),
             ww_label, 8, color=RGBColor(0xB3, 0xD7, 0xF5), align=PP_ALIGN.RIGHT)
        _txb(Inches(9.5), Inches(0.36), Inches(3.6), Inches(0.24),
             'Intel® INTERNAL  ·  Data Centre Graphics', 7,
             color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.RIGHT)

        # ── accent strip below header ──────────────────────────────────────
        strip = slide.shapes.add_shape(1, 0, Inches(0.72), W, Pt(4))
        strip.fill.solid(); strip.fill.fore_color.rgb = accent
        strip.line.fill.background()

        # ── table ─────────────────────────────────────────────────────────
        COLS  = ['Feature', 'Priority', 'TF Meeting Minutes (latest WW update)', 'ETA', 'Blocker']
        WIDTHS = [Inches(3.2), Inches(0.6), Inches(5.8), Inches(0.9), Inches(2.63)]
        FIELDS = ['feature', 'priority', 'status', 'eta', 'blocker']

        top      = Inches(0.82)
        left     = Inches(0.18)
        n_rows   = len(rows)
        tbl_h    = H - top - Inches(0.35)
        row_h    = min(Pt(22), tbl_h / max(n_rows + 1, 1))

        tbl = slide.shapes.add_table(
            n_rows + 1, len(COLS), left, top,
            sum(WIDTHS), int(tbl_h)
        ).table

        # column widths
        for ci, w in enumerate(WIDTHS):
            tbl.columns[ci].width = int(w)

        def _cell(r, c, text, bg, fg=TEXT_DARK, bold=False, size=7, align=PP_ALIGN.LEFT):
            cell = tbl.cell(r, c)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            tf = cell.text_frame; tf.word_wrap = True
            p  = tf.paragraphs[0]; p.alignment = align
            run = p.add_run(); run.text = str(text or '')
            run.font.size = Pt(size); run.font.bold = bold
            run.font.color.rgb = fg
            # thin border
            for side in ('top', 'bottom', 'left', 'right'):
                border = getattr(cell, f'_{side}_border' , None)

        # header row
        hdr_cols = ['Feature', 'Pri', 'TF Meeting Minutes (latest WW)', 'ETA', 'Blocker']
        for ci, h in enumerate(hdr_cols):
            _cell(0, ci, h, INTEL_BLUE, WHITE, bold=True, size=8, align=PP_ALIGN.CENTER)
        tbl.rows[0].height = int(Pt(20))

        # data rows
        p1_map  = {'P1': RED, 'P2': RGBColor(0xE6, 0x7E, 0x22), 'P3': MUTED}
        for ri, row in enumerate(rows):
            bg = LIGHT_BLUE if ri % 2 == 0 else WHITE
            tbl.rows[ri + 1].height = int(row_h)
            for ci, field in enumerate(FIELDS):
                val  = row.get(field, '') or ''
                fg   = TEXT_DARK
                bold = False
                if field == 'priority':
                    fg   = p1_map.get(val.upper(), TEXT_DARK)
                    bold = val.upper() == 'P1'
                if field == 'blocker' and val:
                    fg = RED
                _cell(ri + 1, ci, val, bg, fg, bold, size=7)

        # ── footer ────────────────────────────────────────────────────────
        _txb(Inches(0.18), H - Inches(0.3), Inches(10), Inches(0.28),
             'Intel Corporation  ·  Data Centre Graphics & Compute Engineering  ·  INTEL INTERNAL',
             6.5, color=MUTED)
        _txb(Inches(10.5), H - Inches(0.3), Inches(2.6), Inches(0.28),
             f'{len(rows)} items', 6.5, color=MUTED, align=PP_ALIGN.RIGHT)

    _add_slide('🟣  XPUM / Sysman — Pending Features', PURPLE, xpum_rows)
    _add_slide('🟢  XeKMD — Pending Features',         TEAL,   xekmd_rows)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


@bp.route('/cri-daily-tf/export-slides')
def cri_daily_tf_export_slides():
    xpum_rows, xekmd_rows, _ = _load_data()
    try:
        pptx_bytes = _build_pptx(xpum_rows, xekmd_rows)
    except Exception as exc:
        logger.exception('PPTX export failed')
        return jsonify({'error': str(exc)}), 500
    fname = f"CRI_TF_{datetime.utcnow().strftime('%Y-WW%W')}.pptx"
    return send_file(
        io.BytesIO(pptx_bytes),
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
        as_attachment=True,
        download_name=fname,
    )
